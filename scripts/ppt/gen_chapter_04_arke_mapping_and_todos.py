"""
Chapter 04 deck generator: Map industry progress/challenges to Arke design + TODOs.

Output: docs/sharing/chapter-04-arke-mapping-and-todos.pptx
"""

from __future__ import annotations

from pathlib import Path
import importlib.util

from pptx import Presentation


def _load_master_module():
    master_path = Path(__file__).resolve().parents[1] / "gen_sharing_ppt.py"
    spec = importlib.util.spec_from_file_location("gen_sharing_ppt", master_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Failed to load module from {master_path}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)  # type: ignore[attr-defined]
    return mod


def build(out_path: Path) -> None:
    deck = _load_master_module()

    prs = Presentation()
    prs.slide_width = deck.SLIDE_W
    prs.slide_height = deck.SLIDE_H

    total = 10

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    deck.set_bg(s)
    deck.slide_header(s, "Chapter 04", "Arke 对齐：承接业界进展、解决关键难题，并显式列出 TODO")
    deck.panel(s, deck.Inches(0.9), deck.Inches(2.05), deck.Inches(11.5), deck.Inches(4.9), deck.BG_PANEL)
    deck.add_text(
        s,
        deck.Inches(1.2),
        deck.Inches(2.25),
        deck.Inches(11.0),
        deck.Inches(0.7),
        "输入：Chapter 02/03 的 7 方案洞察 + H1–H10 难题清单。\n输出：逐条映射到 Arke 四件套设计；未覆盖点标为 TODO（含验证方式）。",
        size=15,
        color=deck.FG,
    )
    deck.add_bullets(
        s,
        deck.Inches(1.2),
        deck.Inches(3.25),
        deck.Inches(11.0),
        deck.Inches(3.3),
        [
            "对齐目标：用「可验证的结构化系统」替代「自由代码搜索」的不可控成本",
            "落地优先：首场景 = AI 算力下的算子自动生成与调优",
            "工程标准：Gate/Benchmark 驱动，保证可复现、可回归、可运营",
        ],
        size=13,
        bullet_color=deck.ACCENT,
        line_spacing=1.18,
    )
    deck.footer(s, 1)

    # 2 Overall architecture
    deck.make_overall_arch(prs, page=2)

    # 3 Four-piece overview (reuse existing slide)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    deck.set_bg(s)
    deck.slide_header(s, "Chapter 04", "四件套能力一览：Language × IR × Compiler Toolchain × Agent")
    items = [
        ("Language", "AI-Native DSL", deck.ACCENT, [
            "算子级抽象 · where 符号维度 · strategy/annotation",
            "LLM 可读写 · token 经济 · 入口统一",
        ]),
        ("IR", "Multi-Layer IR", deck.ACCENT_ALT, [
            "Semantic / Strategy / Schedule / Instruction 四层",
            "按 LLM 参与度分层，策略可搜索、可回滚",
        ]),
        ("Compiler", "Registry + Pass + Verifier", deck.ACCENT_BLUE, [
            "OpRegistry SSOT · Pass Pipeline · BackendRegistry",
            "V0/V1/V2 三级验证 · checkpoint/rollback",
        ]),
        ("Agent", "Bounded Action + Tool Protocol", deck.ACCENT_RED, [
            "list_legal_actions → apply_decision",
            "budget / compact / trajectory / rationale",
        ]),
    ]
    card_w = deck.Inches(2.95)
    card_h = deck.Inches(3.9)
    gap = deck.Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start = (deck.SLIDE_W - total_w) / 2
    for i, (tag, subtitle, col, bullets) in enumerate(items):
        x = start + i * (card_w + gap)
        deck.panel(s, x, deck.Inches(2.0), card_w, card_h, deck.BG_PANEL)
        deck.accent_bar(s, x, deck.Inches(2.0), width=card_w, height=deck.Inches(0.08), color=col)
        deck.add_text(s, x + deck.Inches(0.25), deck.Inches(2.15), card_w - deck.Inches(0.5), deck.Inches(0.45),
                      tag, size=16, bold=True, color=deck.FG)
        deck.add_text(s, x + deck.Inches(0.25), deck.Inches(2.55), card_w - deck.Inches(0.5), deck.Inches(0.35),
                      subtitle, size=10, color=deck.FG_MUTED)
        deck.add_bullets(s, x + deck.Inches(0.25), deck.Inches(3.0), card_w - deck.Inches(0.5), deck.Inches(2.7),
                         bullets, size=11, bullet_color=col, line_spacing=1.15)
    deck.footer(s, 3)

    # 4 Map industry progress (P1–P8) → Arke
    s = prs.slides.add_slide(prs.slide_layouts[6])
    deck.set_bg(s)
    deck.slide_header(s, "Chapter 04", "承接业界好进展（P1–P8）→ Arke 设计落点")
    left = [
        ("P1 搜索化", "Agent loop + check/rollback，把优化建模为可迭代决策序列"),
        ("P2 验证门禁", "V0(<1ms) / V1(数值) / V2(性能) 三级验证"),
        ("P3 硬件信号", "get_hw_profile + compile_and_profile 作为工具协议"),
        ("P4 并行探索", "（TODO）并行 worker / beam-search 编排接口标准化"),
        ("P5 稳定抽象", "SemanticIR/StrategyIR 作为稳定中间层，下接 Triton/MLIR/LLVM"),
        ("P6 知识注入", "@rationale + 可检索的 KB（逐步替代纯 prompt 模板）"),
        ("P7 目标函数", "Benchmark BL/OT/ST/L + Gate 出口标准（可运营）"),
        ("P8 可观测性", "trajectory JSONL / report / artifacts 作为一等产物"),
    ]
    top = deck.Inches(2.05)
    row_h = deck.Inches(0.62)
    for i, (k, v) in enumerate(left):
        y = top + i * row_h
        deck.panel(s, deck.Inches(0.9), y, deck.Inches(11.5), row_h - deck.Inches(0.08), deck.BG_PANEL)
        deck.accent_bar(s, deck.Inches(0.9), y, width=deck.Inches(0.12), height=row_h - deck.Inches(0.08), color=deck.ACCENT)
        deck.add_text(s, deck.Inches(1.1), y + deck.Inches(0.1), deck.Inches(2.8), deck.Inches(0.4),
                      k, size=12, bold=True, color=deck.ACCENT, anchor=deck.MSO_ANCHOR.MIDDLE)
        deck.add_text(s, deck.Inches(3.3), y + deck.Inches(0.1), deck.Inches(8.9), deck.Inches(0.4),
                      v, size=11, color=deck.FG, anchor=deck.MSO_ANCHOR.MIDDLE)
    deck.footer(s, 4)

    # 5–6 Map H1–H10 challenges to Arke solutions + TODOs
    def map_slide(page: int, title: str, items):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        deck.set_bg(s)
        deck.slide_header(s, "Chapter 04", title)
        top = deck.Inches(1.95)
        row_h = deck.Inches(0.86)
        for i, (hid, sol, todo) in enumerate(items):
            y = top + i * row_h
            deck.panel(s, deck.Inches(0.9), y, deck.Inches(11.5), deck.Inches(0.78), deck.BG_PANEL)
            # badge
            badge = s.shapes.add_shape(deck.MSO_SHAPE.OVAL, deck.Inches(1.05), y + deck.Inches(0.14), deck.Inches(0.5), deck.Inches(0.5))
            deck._set_fill(badge, deck.ACCENT_RED)
            deck._no_line(badge)
            deck.add_text(s, deck.Inches(1.05), y + deck.Inches(0.14), deck.Inches(0.5), deck.Inches(0.5),
                          hid, size=10, bold=True, color=deck.BG_DARK,
                          align=deck.PP_ALIGN.CENTER, anchor=deck.MSO_ANCHOR.MIDDLE,
                          font=deck.FONT_MONO)
            deck.add_text(s, deck.Inches(1.7), y + deck.Inches(0.08), deck.Inches(6.9), deck.Inches(0.32),
                          "Arke 对应方案： " + sol, size=11, bold=True, color=deck.FG)
            deck.add_text(s, deck.Inches(1.7), y + deck.Inches(0.42), deck.Inches(10.6), deck.Inches(0.32),
                          "TODO： " + todo, size=10, color=deck.FG_MUTED)
        deck.footer(s, page)

    map_slide(5, "关键难题对齐（H1–H5）：结构化表达 × 跨硬件 × 深层控制", [
        ("H1", "Semantic/Strategy 分离 + Strategy IR（可回滚可审计）", "将自由代码优化转为 legal_actions 枚举的完备性评估"),
        ("H2", "@rationale 作为一等字段 + trajectory JSONL", "构建可检索 rationale KB（索引/相似度/跨硬件迁移评估）"),
        ("H3", "where + symbolic dimension system + conditional strategy", "动态 shape specialization 策略与 cache 设计的基准化"),
        ("H4", "target-aware strategy + BackendRegistry", "跨硬件 lift 的量化实验（≥10%）与失败模式归因"),
        ("H5", "Phase 3 MLIR / Phase 4 LLVM IR 路线", "L2/L3 决策空间（memory/warp/register/barrier）的 action schema 设计"),
    ])

    map_slide(6, "关键难题对齐（H6–H10）：成本治理 × 规模化落地 × 可控性", [
        ("H6", "budget（decisions/compiles）+ prompt cache/compact", "跨 session 的长期记忆与压缩策略（不丢 ground truth）"),
        ("H7", "Benchmark BL6 + Gate 驱动", "模型级 bottleneck→再优化→回归 的产品化闭环（orchestrator）"),
        ("H8", "BL/OT/ST/L 基准体系 + 多 baseline（P0–P5）", "端到端指标（吞吐/延迟/内存）统一报表与回归门禁"),
        ("H9", "SemanticIR 的组合性 + StrategyIR 的融合决策", "融合边界自动发现 + 语义合约（subgraph contract）标准化"),
        ("H10", "Bounded Action Space + V0 静态验证", "动作空间安全策略（资源上限/越界/并发）与 sandbox 执行"),
    ])

    # 7 Architecture pages (reuse existing)
    deck.make_lang_arch(prs, page=7)
    deck.make_ir_arch(prs, page=8)
    deck.make_compiler_arch(prs, page=9)
    deck.make_agent_arch(prs, page=10)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote: {out_path}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "docs" / "sharing" / "chapter-04-arke-mapping-and-todos.pptx"
    build(out)

