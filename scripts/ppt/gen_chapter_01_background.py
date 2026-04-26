"""
Chapter 01 deck generator: Why we need this insight (background + framing).

Output: docs/sharing/chapter-01-background.pptx
Theme: 16:9 dark navy + teal accents (match main deck style)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# ------------------------------ Theme (aligned with master)
BG_DARK = RGBColor(0x0B, 0x1B, 0x2E)
BG_PANEL = RGBColor(0x12, 0x2A, 0x43)
BG_PANEL_ALT = RGBColor(0x18, 0x35, 0x53)
FG = RGBColor(0xEA, 0xF2, 0xFA)
FG_MUTED = RGBColor(0x9F, 0xB5, 0xC9)

ACCENT = RGBColor(0x38, 0xD1, 0xB8)
ACCENT_ALT = RGBColor(0xF2, 0xC5, 0x5C)
ACCENT_BLUE = RGBColor(0x7A, 0xB8, 0xFF)
ACCENT_RED = RGBColor(0xEF, 0x6E, 0x6E)
STROKE = RGBColor(0x24, 0x44, 0x66)

FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ------------------------------ Helpers
def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _set_line(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_bg(slide, color: RGBColor = BG_DARK) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, color)
    _no_line(bg)
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def panel(slide, left, top, width, height, color=BG_PANEL, stroke=STROKE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.06
    _set_fill(shp, color)
    _set_line(shp, stroke, 0.75)
    return shp


def accent_bar(slide, left, top, width=Inches(0.12), height=Inches(0.42), color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill(shp, color)
    _no_line(shp)
    return shp


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = FG,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    *,
    size: int = 14,
    color: RGBColor = FG,
    bullet_color: RGBColor = ACCENT,
    line_spacing: float = 1.25,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        marker = "▸ " if level == 0 else "·  "
        r1 = p.add_run()
        r1.text = marker
        r1.font.name = FONT_BODY
        r1.font.size = Pt(size)
        r1.font.bold = level == 0
        r1.font.color.rgb = bullet_color if level == 0 else FG_MUTED
        r2 = p.add_run()
        r2.text = text
        r2.font.name = FONT_BODY
        r2.font.size = Pt(size - (0 if level == 0 else 1))
        r2.font.color.rgb = color if level == 0 else FG_MUTED
        p.level = level
    return tb


def slide_header(slide, kicker: str, title: str):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35), kicker,
             size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8), title,
             size=28, bold=True, color=FG)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55),
                                  Inches(0.6), Emu(26000))
    _set_fill(line, ACCENT)
    _no_line(line)


def footer(slide, page_num: int, total: int):
    add_text(slide, Inches(0.6), Inches(7.15), Inches(6), Inches(0.28),
             "AI-Native Compile Stack · Insight Deck · Chapter 01",
             size=9, color=FG_MUTED)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.28),
             f"{page_num:02d} / {total:02d}",
             size=9, color=FG_MUTED, align=PP_ALIGN.RIGHT)


# ------------------------------ Slides
def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    for y in (Inches(0.0), Inches(7.42)):
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, Inches(0.08))
        _set_fill(b, ACCENT)
        _no_line(b)
    add_text(s, Inches(0.9), Inches(1.2), Inches(11), Inches(0.5),
             "CHAPTER 01 · BACKGROUND", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.9), Inches(1.8), Inches(12), Inches(1.2),
             "为什么需要这份洞察？", size=46, bold=True, color=FG)
    add_text(s, Inches(0.9), Inches(3.0), Inches(12), Inches(0.8),
             "AI 原生算子生成与调优正在从“写代码”走向“做决策 + 可验证闭环”",
             size=20, color=FG_MUTED)


def agenda(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "01 · Background", "本章大纲（我们要解释的 4 件事）")
    items = [
        "为什么“现在”必须做系统洞察（硬件×工作负载×工程）",
        "传统路径的失效点：人写/规则/纯 autotune/LLM 直写",
        "我们如何评估方案：正确性/性能/成本/可迁移/可运营",
        "洞察输出：用案例证据抽象出趋势→进展→难题→构建方案",
    ]
    panel(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.8), BG_PANEL)
    accent_bar(s, Inches(0.9), Inches(2.0), width=Inches(0.12), height=Inches(4.8), color=ACCENT)
    add_bullets(s, Inches(1.2), Inches(2.3), Inches(11.0), Inches(4.3), items,
                size=16, bullet_color=ACCENT, line_spacing=1.22)
    footer(s, page, total)


def why_now(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "01 · Background", "Why Now：三重驱动导致“kernel 规模化危机”")
    cols = [
        ("硬件异构", ACCENT, [
            "SIMT → SIMD → 多类 NPU / DSA",
            "代际更快，优化碎片化",
            "同一算子：不同硬件要不同策略",
        ]),
        ("工作负载复杂化", ACCENT_ALT, [
            "Attention/MLA/GQA/MoE 成为主战场",
            "融合组合爆炸（L2/L3）",
            "动态 shape 常态化",
        ]),
        ("工程成本结构变化", ACCENT_RED, [
            "专家稀缺，需求倍增",
            "验证/回归成本上升",
            "“写一个 kernel”→“维护一个生态”",
        ]),
    ]
    w = Inches(4.05)
    gap = Inches(0.12)
    start = Inches(0.6)
    for i, (t, c, blt) in enumerate(cols):
        x = start + i * (w + gap)
        panel(s, x, Inches(2.0), w, Inches(3.7), BG_PANEL)
        accent_bar(s, x, Inches(2.0), width=w, height=Inches(0.08), color=c)
        add_text(s, x + Inches(0.3), Inches(2.25), w - Inches(0.6), Inches(0.45),
                 t, size=18, bold=True, color=c)
        add_bullets(s, x + Inches(0.3), Inches(2.85), w - Inches(0.6), Inches(2.7),
                    blt, size=13, bullet_color=c, line_spacing=1.18)
    panel(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(1.0), BG_PANEL_ALT)
    add_text(s, Inches(1.1), Inches(6.05), Inches(11.2), Inches(0.8),
             "结论：kernel 优化从“单点手艺活”变成“规模化系统工程”，必须用可验证闭环与可迁移知识资产承载。",
             size=14, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total)


def baselines_fail(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "01 · Background", "传统路径的失效点：4 类基线各自卡住什么？")
    cards = [
        ("人工手写/调优", ACCENT, ["不可扩展", "专家稀缺", "回归维护难"]),
        ("纯 autotune / heuristics", ACCENT_ALT, ["搜索空间爆炸", "启发式难泛化", "缺“知识注入”接口"]),
        ("LLM 直写 Triton/CUDA", ACCENT_RED, ["正确性不稳", "token 成本高", "调试链路长"]),
        ("库/编译器黑盒", ACCENT_BLUE, ["长尾覆盖差", "不可解释", "跨硬件/定制受限"]),
    ]
    w = Inches(5.85)
    h = Inches(1.65)
    for i, (t, c, b) in enumerate(cards):
        x = Inches(0.6) + (i % 2) * Inches(6.2)
        y = Inches(2.0) + (i // 2) * Inches(1.85)
        panel(s, x, y, w, h, BG_PANEL)
        accent_bar(s, x, y, width=Inches(0.12), height=h, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.12), w - Inches(0.4), Inches(0.35),
                 t, size=15, bold=True, color=c)
        add_bullets(s, x + Inches(0.3), y + Inches(0.55), w - Inches(0.4), Inches(1.0),
                    b, size=12, bullet_color=c, line_spacing=1.15)
    panel(s, Inches(0.6), Inches(5.8), Inches(12.15), Inches(1.15), BG_PANEL_ALT)
    add_bullets(s, Inches(0.9), Inches(6.05), Inches(11.7), Inches(0.85), [
        "洞察价值：把“如何规模化产出/验证/沉淀 kernel”抽象成通用机制，而不是某个技巧。",
        "接下来：用 7 个业界系统案例，提炼正在收敛的趋势与尚未解决的难题。",
    ], size=13, bullet_color=ACCENT_ALT, line_spacing=1.15)
    footer(s, page, total)


def evaluation_axes(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "01 · Background", "评估维度：一个“可落地”的系统必须同时满足什么？")
    axes = [
        ("Correctness", "编译通过 + 数值一致 + 回归可测", ACCENT),
        ("Performance", "接近 roofline / vendor library，且方差可控", ACCENT_ALT),
        ("Cost", "token / compile / profile 预算可控", ACCENT_RED),
        ("Generality", "动态 shape + 多算子组合 + 多场景", ACCENT_BLUE),
        ("Portability", "跨硬件（NVIDIA/Ascend/AMD…）策略可迁移", ACCENT),
        ("Operability", "可观测、可回滚、可复现、可审计", ACCENT_ALT),
    ]
    w = Inches(3.85)
    h = Inches(1.55)
    for i, (k, d, c) in enumerate(axes):
        x = Inches(0.6) + (i % 3) * Inches(4.1)
        y = Inches(2.0) + (i // 3) * Inches(1.75)
        panel(s, x, y, w, h, BG_PANEL)
        accent_bar(s, x, y, width=Inches(0.12), height=h, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.15), w - Inches(0.4), Inches(0.35),
                 k, size=14, bold=True, color=c, font=FONT_MONO)
        add_text(s, x + Inches(0.3), y + Inches(0.6), w - Inches(0.4), Inches(0.8),
                 d, size=12, color=FG_MUTED)
    panel(s, Inches(0.6), Inches(5.9), Inches(12.15), Inches(1.0), BG_PANEL_ALT)
    add_text(s, Inches(0.9), Inches(6.05), Inches(11.7), Inches(0.8),
             "这些维度将作为后续“案例趋势 → 难题 → Arke 对齐”的统一坐标系。",
             size=14, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total)


def output_shape(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "01 · Background", "洞察的输出形态：从案例证据到可构建的技术路线")
    panel(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.9), BG_PANEL)
    accent_bar(s, Inches(0.9), Inches(2.0), width=Inches(0.12), height=Inches(4.9), color=ACCENT)
    add_bullets(s, Inches(1.2), Inches(2.3), Inches(11.0), Inches(4.4), [
        "先看案例：KernelEvolve / KernelAgent / AutoKernel / K-Search / AVO / CuTeGen / KernelGen-LM",
        "再抽象趋势：search 化、工具化、硬件信号化、稳定抽象层、知识注入与沉淀",
        "再总结进展：哪些机制已被证明有效（可复现的工程模式）",
        "再总结难题：哪些问题仍缺“体系化解决”（策略资产化、跨硬件迁移、动态 shape 等）",
        "最后引出构建：用语言/IR/工具链/agent 四件套把优势固化、把难题逐条击穿",
    ], size=16, bullet_color=ACCENT_ALT, line_spacing=1.18)
    footer(s, page, total)


def build(out_path: Path) -> None:
    total = 7
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover(prs)                    # 1
    agenda(prs, 2, total)         # 2
    why_now(prs, 3, total)        # 3
    baselines_fail(prs, 4, total) # 4
    evaluation_axes(prs, 5, total)# 5
    output_shape(prs, 6, total)   # 6
    # Chapter bridge slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "01 · Background", "Next：进入业界案例（7 个系统）")
    panel(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(4.5), BG_PANEL)
    add_text(s, Inches(1.2), Inches(2.5), Inches(11.0), Inches(0.6),
             "下一章将逐一拆解 7 个 LLM-driven kernel 系统：",
             size=18, bold=True, color=ACCENT)
    add_bullets(s, Inches(1.2), Inches(3.2), Inches(11.0), Inches(3.2), [
        "KernelEvolve · KernelAgent · AutoKernel · K-Search · AVO · CuTeGen · KernelGen-LM",
        "每方案：技术架构图 + 关键技术特征 + 与 Arke 的对照点",
        "之后：提炼共性进展 & 关键难题，为 Arke 设计对齐铺垫",
    ], size=16, bullet_color=ACCENT_ALT, line_spacing=1.15)
    footer(s, 7, total)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote: {out_path}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "docs" / "sharing" / "chapter-01-background.pptx"
    build(out)

