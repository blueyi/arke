"""
Chapter 02 deck generator: Industry deep-dive on 7 LLM-driven kernel generation/optimization systems.

Output: docs/sharing/chapter-02-industry-methods.pptx
Theme: 16:9 dark navy + teal accents (match main deck style)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# ------------------------------ Theme
BG_DARK = RGBColor(0x0B, 0x1B, 0x2E)
BG_PANEL = RGBColor(0x12, 0x2A, 0x43)
BG_PANEL_ALT = RGBColor(0x18, 0x35, 0x53)
FG = RGBColor(0xEA, 0xF2, 0xFA)
FG_MUTED = RGBColor(0x9F, 0xB5, 0xC9)

ACCENT = RGBColor(0x38, 0xD1, 0xB8)
ACCENT_ALT = RGBColor(0xF2, 0xC5, 0x5C)
ACCENT_BLUE = RGBColor(0x7A, 0xB8, 0xFF)
ACCENT_RED = RGBColor(0xEF, 0x6E, 0x6E)
STROKE = RGBColor(0x24, 0x44, 0x66)

FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ------------------------------ Helpers


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _set_line(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_bg(slide, color: RGBColor = BG_DARK) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, color)
    _no_line(bg)
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = FG,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def panel(slide, left, top, width, height, color=BG_PANEL, stroke=STROKE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.06
    _set_fill(shp, color)
    _set_line(shp, stroke, 0.75)
    return shp


def accent_bar(slide, left, top, width=Inches(0.12), height=Inches(0.42), color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill(shp, color)
    _no_line(shp)
    return shp


def add_bullets(slide, left, top, width, height, bullets, *, size=12, bullet_color=ACCENT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = level
        p.space_after = Pt(2)
        marker = "▸ " if level == 0 else "·  "
        r1 = p.add_run()
        r1.text = marker
        r1.font.name = FONT_BODY
        r1.font.size = Pt(size)
        r1.font.bold = level == 0
        r1.font.color.rgb = bullet_color if level == 0 else FG_MUTED
        r2 = p.add_run()
        r2.text = text
        r2.font.name = FONT_BODY
        r2.font.size = Pt(size - (0 if level == 0 else 1))
        r2.font.color.rgb = FG if level == 0 else FG_MUTED
    return tb


def slide_header(slide, kicker: str, title: str):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35), kicker,
             size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8), title,
             size=28, bold=True, color=FG)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55),
                                  Inches(0.6), Emu(26000))
    _set_fill(line, ACCENT)
    _no_line(line)


def footer(slide, page_num: int, total: int):
    add_text(slide, Inches(0.6), Inches(7.15), Inches(8), Inches(0.28),
             "Chapter 02 · Industry Methods · LLM-driven kernel generation/optimization",
             size=9, color=FG_MUTED)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.28),
             f"{page_num:02d} / {total:02d}",
             size=9, color=FG_MUTED, align=PP_ALIGN.RIGHT)


def method_slide(prs, page, total, *, name, meta, arch_steps, features, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, f"02 · Industry Methods · {page-2:02d} / 07", name)

    # Meta strip
    panel(slide, Inches(0.6), Inches(1.9), Inches(12.1), Inches(0.55), BG_PANEL_ALT)
    add_text(slide, Inches(0.8), Inches(1.98), Inches(11.8), Inches(0.4),
             meta, size=11, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE)

    # Left: architecture
    panel(slide, Inches(0.6), Inches(2.65), Inches(6.5), Inches(4.3), BG_PANEL)
    accent_bar(slide, Inches(0.6), Inches(2.65), width=Inches(0.12), height=Inches(4.3), color=accent_color)
    add_text(slide, Inches(0.9), Inches(2.75), Inches(6.1), Inches(0.35),
             "技术架构", size=13, bold=True, color=accent_color)
    y = Inches(3.15)
    for h, d in arch_steps:
        panel(slide, Inches(0.9), y, Inches(6.0), Inches(0.52), BG_PANEL_ALT)
        add_text(slide, Inches(1.05), y + Inches(0.05), Inches(2.6), Inches(0.42),
                 h, size=10, bold=True, color=accent_color, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(3.75), y + Inches(0.05), Inches(3.1), Inches(0.42),
                 d, size=9, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.6)

    # Right: features
    panel(slide, Inches(7.35), Inches(2.65), Inches(5.4), Inches(4.3), BG_PANEL)
    accent_bar(slide, Inches(7.35), Inches(2.65), width=Inches(0.12), height=Inches(4.3), color=ACCENT_ALT)
    add_text(slide, Inches(7.65), Inches(2.75), Inches(5.0), Inches(0.35),
             "关键技术特征", size=13, bold=True, color=ACCENT_ALT)
    add_bullets(slide, Inches(7.65), Inches(3.15), Inches(5.0), Inches(3.7),
                features, size=10, bullet_color=ACCENT_ALT)

    footer(slide, page, total)
    return slide


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Deck plan:
    # 01 cover
    # 02 agenda
    # 03 section intro
    # 04-10 7 methods
    # 11 matrix
    # 12 pros/cons
    # 13 trend + arke alignment
    total = 13

    # 1 cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_text(s, Inches(0.9), Inches(1.2), Inches(11), Inches(0.5),
             "CHAPTER 02", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.9), Inches(1.7), Inches(12), Inches(1.2),
             "LLM-Driven Kernel Generation / Optimization", size=40, bold=True, color=FG)
    add_text(s, Inches(0.9), Inches(2.9), Inches(12), Inches(0.8),
             "7 个代表方案的技术架构与关键特征", size=24, bold=True, color=ACCENT_ALT)
    panel(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.4), BG_PANEL)
    add_text(s, Inches(1.1), Inches(5.35), Inches(11), Inches(0.8),
             "KernelEvolve · KernelAgent · AutoKernel · K-Search · AVO · CuTeGen · KernelGen-LM",
             size=14, color=FG)
    add_text(s, Inches(1.1), Inches(6.0), Inches(11), Inches(0.4),
             "基于公开论文 / 官方博客 / 开源仓库抽象，不引入未证实细节。",
             size=11, color=FG_MUTED)
    footer(s, 1, total)

    # 2 agenda
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "AGENDA", "本章内容")
    items = [
        ("01", "为什么业界都在做 LLM-driven kernels", "从 one-shot 到 search+verify+profile 的收敛"),
        ("02", "7 个代表方案逐一拆解", "每个方案：架构流水 + 关键技术特征"),
        ("03", "汇总矩阵 + 优劣势速览", "按维度归纳差异点与共性收敛"),
        ("04", "演进趋势与 Arke 呼应", "从自由代码搜索 → 结构化决策 + 多层 IR"),
    ]
    top = Inches(2.0)
    for i, (no, title, sub) in enumerate(items):
        y = top + i * Inches(0.9)
        panel(s, Inches(0.9), y, Inches(11.5), Inches(0.75), BG_PANEL)
        add_text(s, Inches(1.05), y + Inches(0.15), Inches(0.9), Inches(0.45),
                 no, size=18, bold=True, color=ACCENT, font=FONT_MONO)
        add_text(s, Inches(2.0), y + Inches(0.12), Inches(9.9), Inches(0.3),
                 title, size=15, bold=True, color=FG)
        add_text(s, Inches(2.0), y + Inches(0.42), Inches(9.9), Inches(0.3),
                 sub, size=11, color=FG_MUTED)
    footer(s, 2, total)

    # 3 section intro
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "02 · Industry Methods", "7 个代表方案一览")
    methods = [
        ("KernelEvolve", "Meta · ISCA'26 · 生产级异构 + RAG + 搜索", ACCENT_ALT),
        ("KernelFalcon + KernelAgent", "Meta PyTorch · Deep Agents + HW signal", ACCENT_BLUE),
        ("KernelGen-LM (AscendKernelGen)", "PCL · 领域模型 + SFT/RLEF + NPU DSL", ACCENT_RED),
        ("AutoKernel", "RightNow AI · autoresearch loop + dual backend", ACCENT),
        ("K-Search", "UC Berkeley · co-evolving world model", ACCENT_ALT),
        ("AVO", "NVIDIA · agentic variation operators + evolution", ACCENT_RED),
        ("CuTeGen", "UofT · CuTe abstraction + progressive refine", ACCENT_BLUE),
    ]
    top = Inches(2.05)
    for i, (n, d, c) in enumerate(methods):
        y = top + i * Inches(0.65)
        panel(s, Inches(0.9), y, Inches(11.5), Inches(0.56), BG_PANEL)
        accent_bar(s, Inches(0.9), y, width=Inches(0.12), height=Inches(0.56), color=c)
        add_text(s, Inches(1.15), y + Inches(0.1), Inches(2.2), Inches(0.36),
                 n, size=13, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.45), y + Inches(0.1), Inches(9.0), Inches(0.36),
                 d, size=11, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 3, total)

    # 4-10 methods (content is consistent with the main deck additions)
    method_slide(
        prs, 4, total,
        name="KernelEvolve · Meta 生产级异构 Agent Kernel 编写（ISCA 2026）",
        meta="Meta Platforms · arXiv 2512.23236 · ISCA 2026 · 生产部署 · KernelBench 250/250 PASS · ATen 160 ops × 3 HW · 最高 17× vs PyTorch",
        arch_steps=[
            ("Hardware Knowledge Base", "RAG 注入硬件手册/ISA/优化模板（含专有架构）"),
            ("LLM Synthesizer", "多 DSL 候选：Triton/CuTe/FlyDSL/CUDA/HIP/MTIA C++"),
            ("Tree Search + State", "把优化视为搜索问题，维护自改进 state"),
            ("Job Harness", "规模化并行 compile + profile + benchmark"),
            ("RAG Prompt Synthesis", "运行时上下文 → 动态 prompt 合成"),
            ("Deployment Loop", "生产化入链 · DLRM/Ads 推理吞吐提升"),
        ],
        features=[
            "Search-based（非 one-shot）",
            "覆盖多抽象层与多语言（DSL→CUDA/HIP→专有 C++）",
            "RAG 硬件知识库：将“缺训练语料”的专有硬件纳入优化",
            "异构规模化：NVIDIA / AMD / MTIA / CPU",
            "验证依赖 job-harness（compile/正确性/性能）",
        ],
        accent_color=ACCENT_ALT,
    )

    method_slide(
        prs, 5, total,
        name="KernelFalcon + KernelAgent · PyTorch Deep Agents",
        meta="Meta PyTorch Team · github.com/meta-pytorch/KernelAgent · KernelBench L1/L2/L3 250/250 · 1.56× vs torch.compile · 89% H100 roofline",
        arch_steps=[
            ("Orchestrator", "确定性编排（Python control plane）"),
            ("FuserAgent", "Code-to-code 融合，保持 PyTorch 语义"),
            ("ExtractorAgent", "子图边界 + shape contract → JSON 合同"),
            ("Dispatcher+Workers", "并行 Triton kernel 合成 + early-win"),
            ("Profiler/Judge/Analyze", "NCU + roofline + LLM bottleneck diagnosis"),
            ("Composer", "用合成 kernel 重建 forward 端到端替换"),
        ],
        features=[
            "Deep-agent 分层：任务分解→并行探索→验证门禁",
            "Deterministic orchestration：把控制逻辑从 LLM 拿出来",
            "Grounded tool use：每步经编译器与硬件信号",
            "Hardware-guided loop：NCU/roofline 指导优化",
            "输出为 Triton（仍是自由代码空间）",
        ],
        accent_color=ACCENT_BLUE,
    )

    method_slide(
        prs, 6, total,
        name="KernelGen-LM / AscendKernelGen · NPU 领域模型（PCL, 2026）",
        meta="PCL / 中山大学 / 华为 · arXiv 2601.07160 · Ascend NPU · AscendC DSL · L2 编译成功率 0%→95.5% (Pass@10)",
        arch_steps=[
            ("Ascend-CoT Dataset", "真实 kernel + 文档推理 + 代码推理链"),
            ("KernelGen-LM", "Qwen3 backbone · NPU-aware SFT"),
            ("RLEF", "执行反馈强化：compile/正确性/性能信号"),
            ("NPUKernelBench", "编译/正确性/性能三维评测，难度分层"),
            ("Generate-Eval Loop", "候选→硬件执行→反馈→再生成"),
            ("Open Artifacts", "开源数据/权重/评测工件"),
        ],
        features=[
            "走“训一个领域模型”路线：知识进入权重",
            "强耦合单硬件（AscendC/昇腾）换取深度",
            "将 DSL 数据稀缺用 CoT + RLEF 填补",
            "评测一等公民：编译/正确性/性能全量指标",
            "跨硬件迁移能力有限（需要再训/再适配）",
        ],
        accent_color=ACCENT_RED,
    )

    method_slide(
        prs, 7, total,
        name="AutoKernel · Autoresearch for GPU kernels（RightNow AI, 2026）",
        meta="arXiv 2603.21331 · github.com/RightNow-AI/autokernel · H100：RMSNorm 5.29× / softmax 2.82× / cross-entropy 2.21× vs eager · dual backend",
        arch_steps=[
            ("Profile", "torch.profiler 定位瓶颈，Amdahl 排序"),
            ("Extract", "瓶颈算子抽取为独立 kernel 文件"),
            ("Optimize Loop", "agent edit → fixed bench → keep/revert"),
            ("Dual Backend", "Triton（快迭代）+ CUDA C++（深硬件）"),
            ("Correctness Harness", "5-stage correctness gate"),
            ("Playbook", "6-tier playbook（长 instructions）"),
        ],
        features=[
            "autoresearch 闭环：可长时间无人值守探索",
            "Amdahl 优先级：先优化“影响最大”的 kernel",
            "双后端：迭代速度与极限性能兼顾",
            "知识注入以 playbook 为主（提示工程资产）",
            "优化粒度是“编辑文件”，缺少 IR 级可逆决策",
        ],
        accent_color=ACCENT,
    )

    method_slide(
        prs, 8, total,
        name="K-Search · Co-Evolving Intrinsic World Model（UC Berkeley, 2026）",
        meta="arXiv 2602.19128 · github.com/caoshiyi/K-Search · FlashInfer GQA/MLA/MoE：平均 2.10× / MoE 最高 14.3× · GPUMode TriMul H100 SOTA 1030μs",
        arch_steps=[
            ("Problem→Search Tree", "把 kernel 生成视作规划问题"),
            ("World Model", "LLM 估计状态转移与价值，指导搜索"),
            ("Plan↔Codegen", "显式解耦高层策略与低层实现"),
            ("Co-Evolving", "世界模型随搜索共同演化（优先级/假设）"),
            ("Stagnation-Aware", "无提升自动切换 action"),
            ("Exec Feedback", "执行反馈更新 world model"),
        ],
        features=[
            "把 LLM 从“生成器”升级为“世界模型/规划器”",
            "策略/实现解耦：对临时编译失败更鲁棒",
            "复杂 kernel 优势显著（MoE/MLA/GQA）",
            "世界模型大多隐式在 LLM 内部（难跨 session 持久化）",
            "输出仍是自由代码空间（Triton/CUDA）",
        ],
        accent_color=ACCENT_ALT,
    )

    method_slide(
        prs, 9, total,
        name="AVO · Agentic Variation Operators（NVIDIA, 2026）",
        meta="arXiv 2603.24517 · Blackwell B200 · MHA 7 天连续演化 · 超 cuDNN 3.5% / 超 FlashAttention-4 10.5% · GQA 迁移 30min",
        arch_steps=[
            ("EA Skeleton", "population/lineage 演化搜索骨架"),
            ("Vary()=Agent()", "把变异算子替换为自主 agent"),
            ("Self-directed Loop", "propose→repair→critique→verify"),
            ("Domain KB", "硬件规约 + 现有实现（FA4/CUTLASS）"),
            ("Exec Feedback", "性能/正确性反馈驱动下一轮变异"),
            ("Long-horizon", "长周期演化发现微架构级优化"),
        ],
        features=[
            "把“变异算子”本身升级为 agent（从 pipeline → agent）",
            "在 attention 上逼近/超过专家实现",
            "长周期/高预算：探索深但 token/算力成本高",
            "知识多依赖外部 KB（手工维护成本）",
            "输出多为 CUDA/C++ 级实现（门槛高）",
        ],
        accent_color=ACCENT_RED,
    )

    method_slide(
        prs, 10, total,
        name="CuTeGen · LLM × CuTe 抽象层（U. Toronto, 2026）",
        meta="arXiv 2604.01489 · CUTLASS v4.3 (CuTe) · 12 matmul + 14 activation kernels · generate–test–refine · 渐进精炼",
        arch_steps=[
            ("Initial Prompt", "PyTorch reference + CuTe 示例作为锚点"),
            ("Single-kernel Refinement", "精炼一个进化 kernel（非大规模搜索）"),
            ("Exec Validation", "每次改动 compile + ref 对拍"),
            ("Structured Debugging", "失败走结构化模板而非纯报错粘贴"),
            ("Staged Optimization", "先正确性再性能，逐步加压"),
            ("Delayed Profiling", "profile 反馈后置注入减少早期噪声"),
        ],
        features=[
            "关键选择：用 CuTe 作为稳定抽象层（比 CUDA 易迭代）",
            "单 kernel 渐进精炼：token 成本更可控",
            "对性能关键结构（tiling/data-movement）更显式",
            "仍是 C++ 模板体系（非 LLM-native 表示）",
            "跨硬件与知识沉淀机制较弱",
        ],
        accent_color=ACCENT_BLUE,
    )

    # 11 matrix
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "02 · Industry Methods", "7 方案 × 9 维度特征矩阵（概览）")
    dims = ["范式", "结构", "LLM角色", "知识注入", "验证", "目标语言", "硬件", "沉淀介质", "成果"]
    rows = [
        ("KernelEvolve", ["Search", "Tree+State", "生成器", "RAG", "Harness", "Multi-DSL", "Hetero", "KB", "Prod"]),
        ("KernelAgent", ["Agents", "Deep", "多Agent", "NCU", "Gated", "Triton", "NVIDIA", "Templates", "KB100%"]),
        ("KernelGen-LM", ["Train", "Loop", "领域LM", "CoT", "Bench", "AscendC", "Ascend", "Weights", "95.5%"]),
        ("AutoKernel", ["Autoresearch", "Edit/Keep", "Agent", "Playbook", "Harness", "Triton+CUDA", "NVIDIA", "Artifacts", "5.29×"]),
        ("K-Search", ["Planning", "WM Tree", "WorldModel", "Prior", "Exec", "Triton/CUDA", "NVIDIA", "WM", "14.3×"]),
        ("AVO", ["Evolution", "Lineage", "VarOp", "KB", "Verify", "CUDA/C++", "B200", "KB", "FA4+"]),
        ("CuTeGen", ["Refine", "Single", "生成+调试", "CuTe ex", "Ref", "CuTe", "NVIDIA", "—", "GEMM/Act"]),
    ]
    tbl_x, tbl_y = Inches(0.6), Inches(2.0)
    name_w, col_w, row_h = Inches(2.1), Inches(1.13), Inches(0.55)
    panel(s, tbl_x, tbl_y, name_w + col_w * len(dims), row_h, BG_PANEL_ALT)
    add_text(s, tbl_x + Inches(0.1), tbl_y + Inches(0.08), name_w, Inches(0.35),
             "方案", size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    for j, d in enumerate(dims):
        add_text(s, tbl_x + name_w + j * col_w, tbl_y + Inches(0.08),
                 col_w, Inches(0.35), d, size=9, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, (n, vals) in enumerate(rows):
        y = tbl_y + (i + 1) * row_h
        panel(s, tbl_x, y, name_w + col_w * len(dims), row_h - Inches(0.05),
              BG_PANEL if i % 2 == 0 else BG_PANEL_ALT)
        add_text(s, tbl_x + Inches(0.1), y + Inches(0.06), name_w, Inches(0.35),
                 n, size=10, bold=True, color=FG, anchor=MSO_ANCHOR.MIDDLE)
        for j, v in enumerate(vals):
            add_text(s, tbl_x + name_w + j * col_w, y + Inches(0.06),
                     col_w, Inches(0.35), v, size=8, color=FG_MUTED,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 11, total)

    # 12 pros/cons
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "02 · Industry Methods", "共性进展 vs 关键难题（为 Arke 对齐做铺垫）")
    panel(s, Inches(0.6), Inches(1.9), Inches(6.2), Inches(5.1), BG_PANEL)
    accent_bar(s, Inches(0.6), Inches(1.9), width=Inches(0.12), height=Inches(5.1), color=ACCENT)
    add_text(s, Inches(0.9), Inches(2.0), Inches(5.8), Inches(0.4),
             "好的进展（已被验证有效）", size=14, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.9), Inches(2.45), Inches(5.8), Inches(4.4), [
        "从 one-shot 走向 search / evolve / autoresearch",
        "验证门禁：compile→correctness→perf 的 gating 收敛",
        "硬件信号进入 loop：NCU/roofline/metrics 驱动",
        "并行探索 + early-win：多 worker 稳定提升成功率",
        "选择稳定抽象层（Triton/CuTe）降低迭代成本",
    ], size=11, bullet_color=ACCENT)

    panel(s, Inches(7.0), Inches(1.9), Inches(5.7), Inches(5.1), BG_PANEL_ALT)
    accent_bar(s, Inches(7.0), Inches(1.9), width=Inches(0.12), height=Inches(5.1), color=ACCENT_RED)
    add_text(s, Inches(7.3), Inches(2.0), Inches(5.3), Inches(0.4),
             "关键难题（仍未体系化解决）", size=14, bold=True, color=ACCENT_RED)
    add_bullets(s, Inches(7.3), Inches(2.45), Inches(5.3), Inches(4.4), [
        "知识沉淀介质：RAG/权重/模板/日志，难跨硬件迁移",
        "策略与实现仍耦合在自由代码里（难审计、难回滚）",
        "动态 shape 与泛化：正确性与性能同时保证困难",
        "从单 kernel 到模型级自治优化：bottleneck→再优化→回归",
        "Triton 抽象天花板：更深层（寄存器/屏障/指令）决策难表达",
        "token/预算治理：长会话稳定性与成本",
    ], size=11, bullet_color=ACCENT_RED)
    footer(s, 12, total)

    # 13 trend
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    slide_header(s, "02 · Industry Methods", "演进方向：从自由代码搜索 → 结构化决策 + 多层 IR（与 Arke 呼应）")
    add_bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.8), [
        "趋势 1：LLM 从“写代码”→“做决策”（bounded decision search）",
        "趋势 2：语义/策略解耦，从单层代码到多层可验证 IR",
        "趋势 3：验证器前移：V0 快剪枝 + V1 数值 + V2 性能闭环",
        "趋势 4：知识显式化：从权重/模板/日志 → 可迁移的 rationale/trajectory 资产",
        "趋势 5：多后端与多硬件：同一语义，多策略适配不同目标",
    ], size=16, bullet_color=ACCENT_ALT)
    panel(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.65), BG_PANEL_ALT)
    add_text(s, Inches(1.1), Inches(6.4), Inches(11.1), Inches(0.45),
             "下一章（Arke 设计）将逐条对齐这些趋势与难题，并把未覆盖点列为 TODO。",
             size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 13, total)

    prs.save(out_path)


