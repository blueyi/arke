"""
Insight deck generator (single PPTX) — implements outline v0.3.

Output: docs/sharing/ai-native-insight-deck.pptx
Theme:  16:9 dark-navy + teal accents (consistent with existing chapter decks).

Outline mapping (slide → outline section, v0.5):
  01  Cover                                                       (A1)
  02  TOC
  03  Glossary (术语速读)
  04  PART A divider
  05  A1   封面与目标 / 3 questions
  06  A2   模型能力侧 — 6 维能力跃迁
  07  A2.5 范式判断 — 6 条预判 (NEW · 全篇技术前提)
  08  A3   问题域侧 — 硬件 / 工作负载 / 工程现实
  09  A3   真实数据时间线（模型 / 硬件 / 算子复杂度）
  10  A4   传统路径失效（4 类基线）
  11  A5   洞察输出形态（六段式）
  12  PART B divider
  13  B1   全景图
  14  B2   7 案例索引表（含真实数据列）
  15  B2.1 KernelEvolve  ← 真实数据
  16  B2.2 KernelAgent / KernelFalcon
  17  B2.3 AutoKernel
  18  B2.4 K-Search
  19  B2.5 AVO
  20  B2.6 CuTeGen
  21  B2.7 KernelGen-LM / AscendKernelGen
  22  B3   趋势归纳 T1–T6（每条标注代表案例 + ↔ A2.5 预判）
  23  B4   好进展 G1–G6（每条标注代表案例）
  24  B5   关键难题 H1–H9（每条标注业界证据）
  25  TARGET-STATE divider
  26  S1   北极星定义
  27  S2   7 大技术特征 F1–F7
  28  S3   对通用编程的启示
  29  S4   过渡到 Part C
  30  PART C divider
  31  C1   设计逻辑（5 条公理 ← A2.5 反向追溯）
  32  C2   整体架构总图
  33  C3   ① Arke Language
  34  C4   ② Arke IR
  35  C5   ③ Compiler Toolchain
  36  C6   ④ Agent Engineering
  37  C7   ⑤ Benchmark 体系
  38  C8   与 F1–F7 最终对齐表
  39  C9   讨论题 / Q&A
  40  References

Audience: 高层管理者 + 技术专家。
- Pre-Part-C 严禁 Arke 专有概念（@rationale / V0/V1/V2 / SemanticIR / .ak / KernelCache）
- 首次出现的专业术语在 footer 标注 "→ Glossary P03"
- A3 三条曲线 / B2 案例数据 全部使用公开真实数据，标注来源
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# ============================================================
# Theme (consistent with existing chapter decks)
# ============================================================
BG_DARK = RGBColor(0x0B, 0x1B, 0x2E)
BG_PANEL = RGBColor(0x12, 0x2A, 0x43)
BG_PANEL_ALT = RGBColor(0x18, 0x35, 0x53)
BG_PANEL_DEEP = RGBColor(0x0E, 0x21, 0x36)
FG = RGBColor(0xEA, 0xF2, 0xFA)
FG_MUTED = RGBColor(0x9F, 0xB5, 0xC9)
FG_DIM = RGBColor(0x6E, 0x86, 0x9F)

ACCENT = RGBColor(0x38, 0xD1, 0xB8)          # teal
ACCENT_ALT = RGBColor(0xF2, 0xC5, 0x5C)      # gold
ACCENT_BLUE = RGBColor(0x7A, 0xB8, 0xFF)
ACCENT_RED = RGBColor(0xEF, 0x6E, 0x6E)
ACCENT_PURPLE = RGBColor(0xB8, 0x8A, 0xF0)
STROKE = RGBColor(0x24, 0x44, 0x66)

FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DECK_NAME = "AI-Native 算子工程洞察"


# ============================================================
# Drawing primitives
# ============================================================
def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _set_line(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_bg(slide, color: RGBColor = BG_DARK) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, color)
    _no_line(bg)
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def panel(slide, left, top, width, height, color=BG_PANEL, stroke=STROKE,
          rounded: bool = True, adj: float = 0.06):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if rounded:
        shp.adjustments[0] = adj
    _set_fill(shp, color)
    if stroke is None:
        _no_line(shp)
    else:
        _set_line(shp, stroke, 0.75)
    return shp


def accent_bar(slide, left, top, width=Inches(0.12), height=Inches(0.42),
               color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill(shp, color)
    _no_line(shp)
    return shp


def add_text(slide, left, top, width, height, text: str, *,
             size: int = 14, bold: bool = False, color: RGBColor = FG,
             font: str = FONT_BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic: bool = False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, *,
                size: int = 14, color: RGBColor = FG,
                bullet_color: RGBColor = ACCENT,
                line_spacing: float = 1.25,
                marker_l0: str = "▸ ", marker_l1: str = "·  "):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        marker = marker_l0 if level == 0 else marker_l1
        r1 = p.add_run()
        r1.text = marker
        r1.font.name = FONT_BODY
        r1.font.size = Pt(size)
        r1.font.bold = level == 0
        r1.font.color.rgb = bullet_color if level == 0 else FG_DIM
        r2 = p.add_run()
        r2.text = text
        r2.font.name = FONT_BODY
        r2.font.size = Pt(size - (0 if level == 0 else 1))
        r2.font.color.rgb = color if level == 0 else FG_MUTED
        p.level = level
    return tb


def add_runs(slide, left, top, width, height, segments, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing: float = 1.25):
    """segments: list of dicts with text/size/bold/color/font."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for seg in segments:
        run = p.add_run()
        run.text = seg["text"]
        run.font.name = seg.get("font", FONT_BODY)
        run.font.size = Pt(seg.get("size", 14))
        run.font.bold = seg.get("bold", False)
        run.font.italic = seg.get("italic", False)
        run.font.color.rgb = seg.get("color", FG)
    return tb


def slide_header(slide, kicker: str, title: str):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(11), Inches(0.35), kicker,
             size=12, bold=True, color=ACCENT, font=FONT_MONO)
    add_text(slide, Inches(0.6), Inches(0.78), Inches(12.2), Inches(0.85),
             title, size=26, bold=True, color=FG)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.62),
        Inches(0.6), Emu(28000)
    )
    _set_fill(line, ACCENT)
    _no_line(line)


def footer(slide, page_num: int, total: int, section: str = ""):
    left_text = f"{DECK_NAME} · {section}" if section else DECK_NAME
    add_text(slide, Inches(0.6), Inches(7.15), Inches(8), Inches(0.28),
             left_text, size=9, color=FG_DIM)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.28),
             f"{page_num:02d} / {total:02d}",
             size=9, color=FG_DIM, align=PP_ALIGN.RIGHT, font=FONT_MONO)


def glossary_hint(slide, terms: str = "→ 术语见速读 P03"):
    """Bottom-right small annotation for slides that introduce new terms."""
    add_text(slide, Inches(8.5), Inches(7.15), Inches(2.95), Inches(0.28),
             terms, size=9, color=FG_DIM, italic=True,
             align=PP_ALIGN.RIGHT, font=FONT_MONO)


def source_note(slide, text: str):
    """Bottom-left small annotation for data sources / citations."""
    add_text(slide, Inches(0.6), Inches(6.78), Inches(8.5), Inches(0.28),
             text, size=9, color=FG_DIM, italic=True)


def chapter_chip(slide, label: str, color: RGBColor = ACCENT):
    """Top-right kicker chip (e.g. PART A / B / TARGET / PART C)."""
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11.2), Inches(0.45), Inches(1.55), Inches(0.32))
    chip.adjustments[0] = 0.4
    _set_fill(chip, color)
    _no_line(chip)
    add_text(slide, Inches(11.2), Inches(0.45), Inches(1.55), Inches(0.32),
             label, size=10, bold=True, color=BG_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_MONO)


# ============================================================
# Slide 01 · Cover
# ============================================================
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    # top + bottom accent bands
    for y in (0, Inches(7.42)):
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, Inches(0.08))
        _set_fill(b, ACCENT)
        _no_line(b)

    # decorative ring
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(10.2), Inches(0.9),
                              Inches(2.4), Inches(2.4))
    deco.fill.background()
    _set_line(deco, ACCENT, 1.5)
    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(10.7), Inches(1.4),
                               Inches(1.4), Inches(1.4))
    deco2.fill.background()
    _set_line(deco2, ACCENT_ALT, 1.0)

    add_text(s, Inches(0.9), Inches(1.1), Inches(11), Inches(0.5),
             "TECH INSIGHT · 2026", size=14, bold=True, color=ACCENT,
             font=FONT_MONO)
    add_text(s, Inches(0.9), Inches(1.65), Inches(12), Inches(1.4),
             "AI-Native 时代", size=42, bold=True, color=FG)
    add_text(s, Inches(0.9), Inches(2.55), Inches(12), Inches(1.4),
             "算子工程的范式正在被重写", size=42, bold=True, color=FG)

    # subtitle bar
    panel(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(1.0),
          BG_PANEL, stroke=None)
    accent_bar(s, Inches(0.9), Inches(4.2),
               width=Inches(0.14), height=Inches(1.0), color=ACCENT_ALT)
    add_text(s, Inches(1.2), Inches(4.32), Inches(11.0), Inches(0.8),
             "LLM-driven kernel generation/optimization · 业界趋势 · 目标态 · Arke 技术构建",
             size=18, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE)

    # 3-question strip
    qs = [
        ("01", "模型能力", "强到能做什么？"),
        ("02", "业界现状", "已经走到哪、卡在哪？"),
        ("03", "未来构建", "目标态长什么样、怎么搭？"),
    ]
    for i, (n, t, q) in enumerate(qs):
        x = Inches(0.9) + i * Inches(3.95)
        panel(s, x, Inches(5.5), Inches(3.75), Inches(1.3), BG_PANEL_ALT,
              stroke=None)
        add_text(s, x + Inches(0.3), Inches(5.6), Inches(0.7), Inches(0.4),
                 n, size=22, bold=True, color=ACCENT, font=FONT_MONO)
        add_text(s, x + Inches(1.0), Inches(5.62), Inches(2.6), Inches(0.4),
                 t, size=14, bold=True, color=FG)
        add_text(s, x + Inches(1.0), Inches(6.0), Inches(2.6), Inches(0.7),
                 q, size=12, color=FG_MUTED)

    add_text(s, Inches(0.9), Inches(7.0), Inches(11), Inches(0.3),
             "演讲者 · 内部技术分享     |     " + DECK_NAME,
             size=10, color=FG_DIM, font=FONT_MONO)


# ============================================================
# Slide 02 · Table of Contents
# ============================================================
def slide_toc(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_text(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35),
             "TABLE OF CONTENTS", size=12, bold=True, color=ACCENT,
             font=FONT_MONO)
    add_text(s, Inches(0.6), Inches(0.78), Inches(12), Inches(0.85),
             "总目录", size=26, bold=True, color=FG)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.62), Inches(0.6),
                              Emu(28000))
    _set_fill(line, ACCENT)
    _no_line(line)

    entries = [
        ("Glossary · 术语速读", "LLM 工程 · GPU/NPU 硬件 · 业界评测系统",
         "03", ACCENT_RED),
        ("Part A · 背景", "模型能力 · 范式判断 (A2.5) · 业界现状 · 失效",
         "05–11", ACCENT),
        ("Part B · 业界趋势", "7 个 LLM-driven kernel 案例 → T1–T6 / G1–G6 / H1–H9",
         "13–24", ACCENT_BLUE),
        ("🌟 目标态总览", "北极星 · F1–F7 七大特征 · 对通用编程的启示",
         "26–29", ACCENT_ALT),
        ("Part C · Arke 技术构建", "设计逻辑 · 总架构 · 四件套 + Benchmark · F 对齐",
         "31–39", ACCENT_PURPLE),
    ]
    top = Inches(2.0)
    row_h = Inches(0.88)
    for i, (name, sub, rng, col) in enumerate(entries):
        y = top + i * row_h
        p = panel(s, Inches(0.9), y, Inches(11.5), Inches(0.78), BG_PANEL,
                  stroke=None)
        accent_bar(s, Inches(0.9), y, width=Inches(0.14),
                   height=Inches(0.78), color=col)
        add_text(s, Inches(1.2), y + Inches(0.1), Inches(8.5), Inches(0.36),
                 name, size=15, bold=True, color=FG)
        add_text(s, Inches(1.2), y + Inches(0.43),
                 Inches(8.5), Inches(0.32),
                 sub, size=11.5, color=FG_MUTED)
        add_text(s, Inches(10.0), y + Inches(0.2), Inches(2.3), Inches(0.45),
                 rng, size=17, bold=True, color=col,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_MONO)

    footer(s, page, total, "TOC")


# ============================================================
# Slide 03 · Glossary (术语速读)
# ============================================================
def slide_glossary(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_text(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35),
             "GLOSSARY · 术语速读", size=12, bold=True, color=ACCENT,
             font=FONT_MONO)
    add_text(s, Inches(0.6), Inches(0.78), Inches(12), Inches(0.85),
             "面向高层 + 专家的术语速读表", size=26, bold=True, color=FG)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.62), Inches(0.6),
                              Emu(28000))
    _set_fill(line, ACCENT)
    _no_line(line)
    add_text(s, Inches(0.6), Inches(1.72), Inches(12.2), Inches(0.32),
             "本 deck 后续页若引入下列概念，统一指代下表定义；详细说明请扫题对号入座。",
             size=11, color=FG_MUTED, italic=True)

    # Three columns × multiple rows
    cols = [
        ("LLM 工程层", ACCENT, [
            ("Agent / Tool-use",
             "LLM 通过结构化工具调用与外部工具（编译器/profiler）交互"),
            ("RAG",
             "Retrieval-Augmented Generation；将外部知识检索后注入上下文"),
            ("Tool harness",
             "把编译/对拍/性能测试封装为可被 Agent 调用的工具集"),
            ("RLEF",
             "RL with Execution Feedback；用执行结果做强化学习反馈"),
            ("SFT",
             "Supervised Fine-Tuning；用标注样本做监督微调"),
            ("World model",
             "Agent 内部对环境状态的可预测表征，用于规划"),
            ("Decision rationale",
             "对每条优化决策的自然语言解释，用于审计与迁移"),
            ("Bounded action space",
             "决策从编译器枚举的合法动作集中选择，禁止自由代码"),
        ]),
        ("GPU / NPU 硬件层", ACCENT_BLUE, [
            ("SIMT",
             "Single-Instruction Multiple-Thread（NVIDIA / AMD GPU）"),
            ("SIMD",
             "Single-Instruction Multiple-Data（Ascend NPU 等）"),
            ("Tensor Core",
             "NVIDIA GPU 上的矩阵乘加加速单元"),
            ("Roofline",
             "性能上限模型：计算强度 vs 带宽/算力交点"),
            ("NCU",
             "NVIDIA Nsight Compute；kernel 级 profiler"),
            ("Occupancy",
             "SM 同时驻留的 warp 数与硬件上限之比"),
            ("Triton / CuTe / MLIR",
             "GPU 编程抽象层：Triton (块级) / CuTe (布局代数) / MLIR (多方言)"),
            ("FlashAttention (FA-2/3/4)",
             "Tri Dao 等提出的 IO-aware attention 系列实现"),
        ]),
        ("业界评测 / 系统", ACCENT_ALT, [
            ("KernelBench",
             "Stanford 推的 GPU kernel 生成基准；分 L1/L2/L3 三级"),
            ("ATen",
             "PyTorch 后端算子集，主流神经网络的「指令集」 (≈2000+)"),
            ("FlagGems / Liger-Kernel",
             "开源 Triton 算子库；常作为 LLM 生成 kernel 的对照基线"),
            ("torch.compile / Inductor",
             "PyTorch 2.x 图编译器，AOT 路径的事实标准"),
            ("cuBLAS / cuDNN",
             "NVIDIA 官方矩阵 / 卷积库；性能上限基线"),
            ("MLA / GQA / MoE",
             "Multi-Head Latent Attention / Grouped-Query Attention / Mixture-of-Experts"),
            ("Pass@k",
             "Top-k 候选中至少 1 个通过的概率；代码生成评测常用"),
            ("Amdahl",
             "并行加速上限 = 1 / ((1-p) + p/N)；用作瓶颈优先级"),
        ]),
    ]
    cw = Inches(4.0)
    cy = Inches(2.15)
    cx = Inches(0.6)
    for i, (title, c, items) in enumerate(cols):
        x = cx + i * (cw + Inches(0.1))
        panel(s, x, cy, cw, Inches(4.85), BG_PANEL, stroke=None)
        accent_bar(s, x, cy, width=cw, height=Inches(0.08), color=c)
        add_text(s, x + Inches(0.25), cy + Inches(0.2),
                 cw - Inches(0.4), Inches(0.4),
                 title, size=14, bold=True, color=c, font=FONT_MONO)
        ty = cy + Inches(0.65)
        for term, defn in items:
            add_text(s, x + Inches(0.25), ty, cw - Inches(0.4), Inches(0.2),
                     term, size=10.5, bold=True, color=FG)
            add_text(s, x + Inches(0.25), ty + Inches(0.2),
                     cw - Inches(0.4), Inches(0.32),
                     defn, size=9, color=FG_MUTED)
            ty = ty + Inches(0.52)

    add_text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.25),
             "提示：Arke 项目专有概念（@rationale / 多层 IR / 多级验证 / KernelCache 等）将在 Part C 引入；本表只列业界通用术语。",
             size=9, color=FG_DIM, italic=True)

    footer(s, page, total, "Glossary")


# ============================================================
# Generic chapter divider
# ============================================================
def slide_divider(prs, page, total, kicker, title, subtitle, accent=ACCENT):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    for y in (0, Inches(7.42)):
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, Inches(0.08))
        _set_fill(b, accent)
        _no_line(b)
    add_text(s, Inches(0.9), Inches(1.6), Inches(12), Inches(0.55),
             kicker, size=18, bold=True, color=accent, font=FONT_MONO)
    add_text(s, Inches(0.9), Inches(2.2), Inches(12), Inches(1.6),
             title, size=44, bold=True, color=FG)
    panel(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.4), BG_PANEL,
          stroke=None)
    accent_bar(s, Inches(0.9), Inches(4.5), width=Inches(0.14),
               height=Inches(1.4), color=ACCENT_ALT)
    add_text(s, Inches(1.2), Inches(4.65), Inches(11.0), Inches(1.1),
             subtitle, size=15, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total, kicker.split(" · ")[0] if " · " in kicker else kicker)


# ============================================================
# PART A
# ============================================================
def slide_a1_questions(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART A", ACCENT)
    slide_header(s, "A1 · 背景与目标",
                 "本次分享要回答的 3 个问题")
    panel(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.0),
          BG_PANEL_ALT, stroke=None)
    add_text(s, Inches(1.2), Inches(2.1), Inches(11.0), Inches(0.8),
             "AI-Native 时代，算子工程的范式正在被重写",
             size=20, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

    qs = [
        ("01", "模型能力跃迁", "强到能做什么？", ACCENT,
         ["代码合成 / 结构化推理 / 长上下文",
          "可验证闭环 + Agent 化",
          "后训练让经验可累计"]),
        ("02", "业界相关方向", "走到哪、卡在哪？", ACCENT_BLUE,
         ["7 个代表性 LLM-driven kernel 案例",
          "趋势 T1–T6 已开始收敛",
          "9 个关键技术难题待解"]),
        ("03", "未来如何构建", "目标态长什么样？", ACCENT_ALT,
         ["7 大技术特征 F1–F7",
          "对通用编程的 6 条启示",
          "Arke 给出一种可落地的回答"]),
    ]
    for i, (n, t, q, col, items) in enumerate(qs):
        x = Inches(0.6) + i * Inches(4.15)
        panel(s, x, Inches(3.2), Inches(4.0), Inches(3.7), BG_PANEL,
              stroke=None)
        accent_bar(s, x, Inches(3.2), width=Inches(4.0),
                   height=Inches(0.08), color=col)
        add_text(s, x + Inches(0.3), Inches(3.4), Inches(0.7), Inches(0.5),
                 n, size=24, bold=True, color=col, font=FONT_MONO)
        add_text(s, x + Inches(1.1), Inches(3.42), Inches(2.7), Inches(0.4),
                 t, size=15, bold=True, color=FG)
        add_text(s, x + Inches(1.1), Inches(3.78), Inches(2.7), Inches(0.4),
                 q, size=12, color=FG_MUTED, italic=True)
        add_bullets(s, x + Inches(0.3), Inches(4.5), Inches(3.5),
                    Inches(2.3), items, size=12, bullet_color=col,
                    line_spacing=1.22)
    footer(s, page, total, "Part A")


def slide_a2_model_capability(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART A", ACCENT)
    slide_header(s, "A2 · 模型能力侧",
                 "LLM 能力跃迁带来的「新可能」——为什么以前做不了，今天可以做")

    caps = [
        ("代码生成", "片段补全 → 跨文件、跨抽象层程序合成", ACCENT),
        ("结构化推理", "Tool-use / Function-calling / 结构化输出", ACCENT_BLUE),
        ("长上下文 + RAG", "硬件手册 / ISA / Profile / 历史轨迹入上下文",
         ACCENT_ALT),
        ("可验证闭环", "与编译器/profiler 多轮反馈，而非一次性写出",
         ACCENT_PURPLE),
        ("代理化", "规划 / 执行 / 反思 / 回滚 的 Agent 编排", ACCENT_RED),
        ("后训练范式", "SFT / RLHF / RLEF / 轨迹回放 让经验可累计", ACCENT),
    ]
    cell_w = Inches(4.0)
    cell_h = Inches(1.55)
    for i, (t, d, c) in enumerate(caps):
        x = Inches(0.6) + (i % 3) * (cell_w + Inches(0.1))
        y = Inches(2.0) + (i // 3) * (cell_h + Inches(0.15))
        panel(s, x, y, cell_w, cell_h, BG_PANEL, stroke=None)
        accent_bar(s, x, y, width=Inches(0.12), height=cell_h, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.18),
                 cell_w - Inches(0.4), Inches(0.4),
                 t, size=15, bold=True, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.65),
                 cell_w - Inches(0.4), Inches(0.85),
                 d, size=12, color=FG_MUTED)

    panel(s, Inches(0.6), Inches(5.65), Inches(12.15), Inches(1.15),
          BG_PANEL_ALT, stroke=None)
    accent_bar(s, Inches(0.6), Inches(5.65), width=Inches(0.14),
               height=Inches(1.15), color=ACCENT)
    add_text(s, Inches(0.95), Inches(5.77), Inches(11.7), Inches(0.45),
             "小结", size=12, bold=True, color=ACCENT, font=FONT_MONO)
    add_text(s, Inches(0.95), Inches(6.10), Inches(11.7), Inches(0.7),
             "模型从「会写代码」升级为「会做带反馈的工程决策」——这是后续所有趋势的能力底座。",
             size=15, bold=True, color=FG)
    footer(s, page, total, "Part A")


# ============================================================
# A2.5 · Paradigm-shift thesis (key foundation for the whole deck)
# ============================================================
def slide_a2p5_thesis(prs, page, total):
    """A2.5 — overview: the 6-prediction thesis that anchors the deck."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART A", ACCENT)
    slide_header(
        s, "A2.5 · 范式判断（本洞察的技术前提）",
        "从「代码编写」到「意图治理」的范式迁移——后续所有判断的根")

    # Banner statement
    panel(s, Inches(0.6), Inches(1.95), Inches(12.15), Inches(0.85),
          BG_PANEL_ALT, stroke=None)
    accent_bar(s, Inches(0.6), Inches(1.95), width=Inches(0.16),
               height=Inches(0.85), color=ACCENT_ALT)
    add_runs(
        s, Inches(0.95), Inches(2.0), Inches(11.5), Inches(0.75),
        [
            {"text": "未来开发生态的胜负手不再是 ", "size": 16,
             "color": FG_MUTED},
            {"text": "「如何写代码」", "size": 16, "bold": True,
             "color": FG_DIM, "italic": True},
            {"text": "，而是 ", "size": 16, "color": FG_MUTED},
            {"text": "「如何定义意图 + 如何闭环验证」", "size": 16,
             "bold": True, "color": ACCENT},
            {"text": "——本 deck 后续每一条主张，都是这个判断的延伸。",
             "size": 16, "color": FG_MUTED},
        ], anchor=MSO_ANCHOR.MIDDLE)

    # 2x3 grid of 6 predictions
    preds = [
        ("❶", "能力平庸化",
         "顶级模型代差缩小，竞争从「博学」转向\n「特定链路稳定性 / 成功率」",
         "Commodity · FeatureBench",
         "→ T2 / T6 / F2",
         ACCENT),
        ("❷", "Harness 治理",
         "「上下文焦虑」→「智能脚手架」\n图谱索引 + 静态/动态依赖 + 语义路由",
         "Harness · 语义路由器",
         "→ T2 / T3 / F4",
         ACCENT_BLUE),
        ("❸", "颗粒度跃迁",
         "片段式编程被吞噬，开发坍缩为\n「高层架构」+「底层性能」两端",
         "片段消失 · 两端坍缩",
         "→ T1 / T4 / F1",
         ACCENT_ALT),
        ("❹", "Agent 自主流转",
         "从「辅助生成」到「自愈工作流」\n自主拉仓库 / 跑测试 / 修 bug",
         "Self-healing · 端到端自治",
         "→ T1 / G6 / F2",
         ACCENT_PURPLE),
        ("❺", "意图 + 形式化验证",
         "代码边际成本 → 0；「确保 AI 做对」成为\n最贵环节，验证驱动开发",
         "形式化验证 · 抗幻觉",
         "→ G1 / F3 / F5",
         ACCENT_RED),
        ("❻", "AI 原生基座",
         "软件主动适配 AI：原生编译器 +\nAI-friendly IR + 强契约模块化",
         "AI-native compiler · 解耦契约",
         "→ T4 / F3 / F6",
         ACCENT),
    ]
    cw = Inches(4.0)
    ch = Inches(1.95)
    cx = Inches(0.6)
    cy = Inches(2.95)
    for i, (n, t, d, kw, mapping, c) in enumerate(preds):
        x = cx + (i % 3) * (cw + Inches(0.1))
        y = cy + (i // 3) * (ch + Inches(0.13))
        panel(s, x, y, cw, ch, BG_PANEL, stroke=None)
        accent_bar(s, x, y, width=Inches(0.14), height=ch, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.1),
                 Inches(0.5), Inches(0.4),
                 n, size=18, bold=True, color=c)
        add_text(s, x + Inches(0.85), y + Inches(0.13),
                 cw - Inches(1.0), Inches(0.4),
                 t, size=14, bold=True, color=FG)
        add_text(s, x + Inches(0.3), y + Inches(0.55),
                 cw - Inches(0.4), Inches(0.85),
                 d, size=10.5, color=FG_MUTED)
        add_text(s, x + Inches(0.3), y + ch - Inches(0.6),
                 cw - Inches(0.4), Inches(0.28),
                 kw, size=9, color=ACCENT_ALT, italic=True,
                 font=FONT_MONO)
        add_text(s, x + Inches(0.3), y + ch - Inches(0.32),
                 cw - Inches(0.4), Inches(0.28),
                 mapping, size=9, bold=True, color=c, font=FONT_MONO)

    # bottom takeaway
    add_text(s, Inches(0.6), Inches(7.04), Inches(12.15), Inches(0.3),
             "当代码生成的边际成本趋近于零，「定义意图 + 闭环验证」成为生产力重心；本 deck 的所有趋势 / 目标态 / Arke 设计都建立在此之上。",
             size=10.5, color=FG_DIM, italic=True,
             align=PP_ALIGN.CENTER)
    footer(s, page, total, "Part A · 范式判断")


def slide_a3_problem_domain(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART A", ACCENT)
    slide_header(s, "A3 · 问题域侧",
                 "算子工程的业界现状——为什么这个领域亟需被 LLM 重塑")

    cols = [
        ("硬件维度", ACCENT, [
            "SIMT (NVIDIA) → SIMD (Ascend)",
            "→ 多类 NPU / DSA",
            "代际叠加 + 异构叠加",
            "同一算子，不同硬件不同策略",
        ]),
        ("工作负载维度", ACCENT_ALT, [
            "Attention / MLA / GQA / MoE",
            "融合组合爆炸 (L2 / L3)",
            "动态 shape 常态化",
            "长上下文 → KV / Paged",
        ]),
        ("工程现实维度", ACCENT_RED, [
            "kernel SKU 爆炸（算子×形状×精度×后端）",
            "顶级 kernel 工程师稀缺且不可扩展",
            "验证 / 回归成本随版本上升",
            "厂商库覆盖不全 + 演进不同步",
        ]),
    ]
    w = Inches(4.05)
    gap = Inches(0.12)
    start = Inches(0.6)
    for i, (t, c, blt) in enumerate(cols):
        x = start + i * (w + gap)
        panel(s, x, Inches(2.0), w, Inches(3.85), BG_PANEL, stroke=None)
        accent_bar(s, x, Inches(2.0), width=w, height=Inches(0.08), color=c)
        add_text(s, x + Inches(0.3), Inches(2.22),
                 w - Inches(0.6), Inches(0.45),
                 t, size=18, bold=True, color=c)
        add_bullets(s, x + Inches(0.3), Inches(2.85),
                    w - Inches(0.6), Inches(2.85),
                    blt, size=13, bullet_color=c, line_spacing=1.22)
    panel(s, Inches(0.6), Inches(6.05), Inches(12.15), Inches(0.95),
          BG_PANEL_ALT, stroke=None)
    add_text(s, Inches(0.85), Inches(6.18), Inches(11.7), Inches(0.7),
             "三股力量同步发散：模型规模 ↑、算子复杂度 ↑、硬件代际 ↑——传统路径无法线性扩展。",
             size=14, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total, "Part A")


def _timeline_panel(slide, x, y, w, h, *, title, color,
                    points, y_label, source_text):
    """One mini timeline chart with real milestones.

    points: list of (x_label, y_label_value, numeric_y_for_position)
    where numeric_y_for_position is in the [0, 1] range used for height.
    """
    panel(slide, x, y, w, h, BG_PANEL_DEEP, stroke=STROKE)
    accent_bar(slide, x, y, width=w, height=Inches(0.06), color=color)
    add_text(slide, x + Inches(0.2), y + Inches(0.12),
             w - Inches(0.4), Inches(0.32),
             title, size=12, bold=True, color=color, font=FONT_MONO)

    # plot area
    pad_l = Inches(0.55)
    pad_r = Inches(0.25)
    pad_t = Inches(0.55)
    pad_b = Inches(0.95)
    px = x + pad_l
    py = y + pad_t
    pw = w - pad_l - pad_r
    ph = h - pad_t - pad_b

    # baseline
    bl = slide.shapes.add_connector(1, px, py + ph, px + pw, py + ph)
    _set_line(bl, FG_MUTED, 1.0)
    # y-axis
    ya = slide.shapes.add_connector(1, px, py, px, py + ph)
    _set_line(ya, FG_MUTED, 0.8)
    # 3 light grid lines
    for k in (0.25, 0.5, 0.75):
        gy = py + ph - Emu(int(ph * k))
        gl = slide.shapes.add_connector(1, px, gy, px + pw, gy)
        _set_line(gl, STROKE, 0.4)
    # y axis label
    add_text(slide, x + Inches(0.05), py - Inches(0.05),
             Inches(2.5), Inches(0.25),
             y_label, size=8, color=FG_MUTED, font=FONT_MONO)

    n = len(points)
    if n < 2:
        return
    # plot polyline + markers + labels
    coords = []
    for i, (xl, yl, yv) in enumerate(points):
        t = i / (n - 1)
        cx = px + Emu(int(pw * t))
        cy = py + ph - Emu(int(ph * max(0.02, min(0.98, yv))))
        coords.append((cx, cy))
    for i in range(n - 1):
        ln = slide.shapes.add_connector(1, coords[i][0], coords[i][1],
                                        coords[i + 1][0], coords[i + 1][1])
        _set_line(ln, color, 2.0)
    for i, (xl, yl, _) in enumerate(points):
        cx_, cy_ = coords[i]
        m = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   cx_ - Inches(0.06),
                                   cy_ - Inches(0.06),
                                   Inches(0.12), Inches(0.12))
        _set_fill(m, color)
        _no_line(m)
        # value label above
        add_text(slide, cx_ - Inches(0.7), cy_ - Inches(0.42),
                 Inches(1.4), Inches(0.25),
                 yl, size=8.5, bold=True, color=FG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)
        # x label below baseline
        add_text(slide, cx_ - Inches(0.5), py + ph + Inches(0.05),
                 Inches(1.0), Inches(0.4),
                 xl, size=9, color=FG_MUTED, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)

    # source citation strip
    add_text(slide, x + Inches(0.2), y + h - Inches(0.42),
             w - Inches(0.4), Inches(0.36),
             source_text, size=7.5, color=FG_DIM, italic=True)


def slide_a3_chart(prs, page, total):
    """Single 2D line chart — three real-data curves on a shared
    time axis (2017→2026), y = log10(× baseline) so divergence is visible."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART A", ACCENT)
    slide_header(s, "A3 · 三条发散曲线（公开数据 · 二维曲线图）",
                 "同一时间轴上的相对增速 — y 轴 = log10(× baseline) (2017→2026)")

    # ============================================================
    # Left panel · 2D chart
    # ============================================================
    cx = Inches(0.6)
    cy = Inches(2.0)
    cw = Inches(8.6)
    ch = Inches(4.7)
    panel(s, cx, cy, cw, ch, BG_PANEL_DEEP, stroke=STROKE)

    # plot area paddings
    pad_l = Inches(0.85)
    pad_r = Inches(0.3)
    pad_t = Inches(0.55)
    pad_b = Inches(0.95)
    px = cx + pad_l
    py = cy + pad_t
    pw = cw - pad_l - pad_r
    ph = ch - pad_t - pad_b

    # axes
    xa = s.shapes.add_connector(1, px, py + ph, px + pw, py + ph)
    _set_line(xa, FG_MUTED, 1.4)
    ya = s.shapes.add_connector(1, px, py, px, py + ph)
    _set_line(ya, FG_MUTED, 1.4)

    YEAR_MIN, YEAR_MAX = 2017, 2026
    span = YEAR_MAX - YEAR_MIN
    Y_MAX = 5.0  # log10 ticks: 0..5  (10^5 = 100000×)

    # X gridlines + year tick labels
    for yr in range(YEAR_MIN, YEAR_MAX + 1):
        t = (yr - YEAR_MIN) / span
        tx = px + Emu(int(pw * t))
        gl = s.shapes.add_connector(1, tx, py, tx, py + ph)
        _set_line(gl, STROKE, 0.4)
        add_text(s, tx - Inches(0.25), py + ph + Inches(0.06),
                 Inches(0.5), Inches(0.28),
                 str(yr), size=9, color=FG_MUTED, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)

    # Y gridlines + log labels
    for k in range(0, int(Y_MAX) + 1):
        v = k / Y_MAX
        gy = py + ph - Emu(int(ph * v))
        gl = s.shapes.add_connector(1, px, gy, px + pw, gy)
        _set_line(gl, STROKE, 0.4)
        label = "1×" if k == 0 else f"10^{k}×"
        add_text(s, px - Inches(0.78), gy - Inches(0.13),
                 Inches(0.7), Inches(0.28),
                 label, size=8.5, color=FG_MUTED, font=FONT_MONO,
                 align=PP_ALIGN.RIGHT)

    # axis titles
    add_text(s, px - Inches(0.85), py - Inches(0.38),
             Inches(2.6), Inches(0.3),
             "log10(× baseline)", size=9, bold=True,
             color=FG_MUTED, font=FONT_MONO)
    add_text(s, px + pw - Inches(1.0), py + ph + Inches(0.36),
             Inches(1.0), Inches(0.3),
             "year", size=9, bold=True, color=FG_MUTED,
             font=FONT_MONO, align=PP_ALIGN.RIGHT)

    # ============================================================
    # 3 series (real public data, normalized to each curve's baseline)
    # ============================================================
    # values are log10(value / baseline_value); selected milestones only
    series = [
        # (color, name, baseline_label, points = [(year, log10_ratio, label)])
        (ACCENT, "① 模型规模 (params)", "GPT-1 117M = 1×", [
            (2018, 0.00, "GPT-1 117M"),
            (2019, 1.11, "GPT-2 1.5B"),
            (2020, 3.17, "GPT-3 175B"),
            (2022, 3.66, "PaLM 540B"),
            (2023, 4.16, "GPT-4 ≈1.7T*"),
            (2024, 3.76, "DeepSeek-V3 671B"),
            (2025, 4.23, "LLaMA-4 ≈2T*"),
        ]),
        (ACCENT_RED, "② NVIDIA HW (BF16 算力)", "V100 125 TF = 1×", [
            (2017, 0.00, "V100 125 TF"),
            (2020, 0.40, "A100 312 TF"),
            (2022, 0.90, "H100 989 TF"),
            (2024, 1.16, "B100 1.8 PF"),
            (2025, 1.26, "B200 2.25 PF"),
            (2026, 1.30, "GB200 2.5 PF"),
        ]),
        (ACCENT_ALT, "③ 算子复杂度 (定性)",
         "ATen 350 / vanilla attn = 1×", [
            (2018, 0.00, "ATen ≈350 · vanilla"),
            (2020, 0.36, "ATen 800+ · MHA"),
            (2022, 0.50, "FA-1"),
            (2023, 0.70, "FA-2 + GQA"),
            (2024, 0.90, "FA-3 + MLA + MoE"),
            (2026, 1.10, "FA-4 (Blackwell)"),
         ]),
    ]

    # plot polylines + markers; collect end-points for series labels
    end_points = []
    for col, name, baseline, pts in series:
        coords = []
        for (yr, v, lbl) in pts:
            t = (yr - YEAR_MIN) / span
            mx = px + Emu(int(pw * t))
            my = py + ph - Emu(int(ph * (max(0.0, min(Y_MAX, v)) / Y_MAX)))
            coords.append((mx, my))
        # polyline
        for i in range(len(coords) - 1):
            ln = s.shapes.add_connector(1, coords[i][0], coords[i][1],
                                        coords[i + 1][0], coords[i + 1][1])
            _set_line(ln, col, 2.5)
        # markers
        for mx, my in coords:
            mk = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    mx - Inches(0.075),
                                    my - Inches(0.075),
                                    Inches(0.15), Inches(0.15))
            _set_fill(mk, col)
            _no_line(mk)
        end_points.append((col, name, coords[-1], pts[-1]))

    # in-chart series end-labels (slightly offset right of last marker)
    for col, name, (mx, my), (yr, v, lbl) in end_points:
        add_text(s, mx + Inches(0.1), my - Inches(0.18),
                 Inches(2.0), Inches(0.32),
                 name.split("·")[0].split("(")[0].strip(),
                 size=9, bold=True, color=col, font=FONT_MONO)

    # ============================================================
    # Right column · legend + key takeaway
    # ============================================================
    rx = Inches(9.4)
    ry = Inches(2.0)
    rw = Inches(3.35)
    rh = Inches(4.7)
    panel(s, rx, ry, rw, rh, BG_PANEL, stroke=None)
    add_text(s, rx + Inches(0.2), ry + Inches(0.18),
             rw - Inches(0.4), Inches(0.36),
             "曲线图例 · 公开数据", size=12, bold=True, color=ACCENT,
             font=FONT_MONO)

    legend_meta = [
        (ACCENT, "① 模型规模",
         "GPT-1 → LLaMA-4 / DSV3：~10⁴ 倍"),
        (ACCENT_RED, "② NVIDIA HW",
         "V100 → GB200 (BF16)：~20 倍"),
        (ACCENT_ALT, "③ 算子复杂度",
         "vanilla attn → FA-4：~10 倍 (定性)"),
    ]
    ly = ry + Inches(0.65)
    for col, t, d in legend_meta:
        sw = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                rx + Inches(0.25),
                                ly + Inches(0.16),
                                Inches(0.32), Inches(0.08))
        _set_fill(sw, col)
        _no_line(sw)
        add_text(s, rx + Inches(0.65), ly,
                 rw - Inches(0.85), Inches(0.32),
                 t, size=11.5, bold=True, color=col)
        add_text(s, rx + Inches(0.65), ly + Inches(0.3),
                 rw - Inches(0.85), Inches(0.5),
                 d, size=9.5, color=FG_MUTED)
        ly = ly + Inches(0.78)

    # KEY TAKEAWAY callout
    panel(s, rx + Inches(0.15), ry + Inches(3.1),
          rw - Inches(0.3), Inches(1.45),
          BG_PANEL_ALT, stroke=None)
    accent_bar(s, rx + Inches(0.15), ry + Inches(3.1),
               width=Inches(0.08), height=Inches(1.45),
               color=ACCENT_ALT)
    add_text(s, rx + Inches(0.32), ry + Inches(3.18),
             rw - Inches(0.5), Inches(0.3),
             "KEY TAKEAWAY", size=10, bold=True, color=ACCENT_ALT,
             font=FONT_MONO)
    add_text(s, rx + Inches(0.32), ry + Inches(3.5),
             rw - Inches(0.5), Inches(1.0),
             "8 年内 · 模型规模 ↑ ~10⁴×，算子复杂度 ↑ ~10×，"
             "但硬件 BF16 算力仅 ↑ ~20×。\n"
             "缺口必须靠编译 / 调度 / 算法填补——这正是本洞察的主题。",
             size=10.5, color=FG)

    source_note(s, "Sources: 各模型公开 paper/blog (GPT-1..4, LLaMA, "
                   "PaLM, DeepSeek-V3) · NVIDIA 官方 datasheet "
                   "(V100/A100/H100/H200/Blackwell) · PyTorch ATen "
                   "源码计数 · FlashAttention 系列 (Tri Dao et al.)。"
                   "* GPT-4 / LLaMA-4 参数为 SemiAnalysis 等公开估算。")
    glossary_hint(s, "→ Glossary P03 · ATen / MLA / GQA / FA-2/3/4")
    footer(s, page, total, "Part A")


def slide_a4_failure_modes(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART A", ACCENT)
    slide_header(s, "A4 · 传统路径为什么不够用",
                 "Failure modes of the status quo")

    cards = [
        ("人工手写 + autotune", ACCENT,
         ["不可扩展", "跨硬件不可迁移", "知识沉淀依赖个人"]),
        ("图编译器 (XLA / Inductor / TVM)", ACCENT_ALT,
         ["高层有效", "深层决策仍依赖启发式 / 专家",
          "tile / pipeline / 寄存器分配"]),
        ("LLM 直写 Triton / CUDA", ACCENT_RED,
         ["正确性差", "token 成本高",
          "可维护性差 + 知识无法跨 kernel 复用"]),
        ("库 / 编译器黑盒", ACCENT_BLUE,
         ["长尾覆盖不全", "不可解释",
          "跨硬件 / 定制扩展受限"]),
    ]
    w = Inches(5.95)
    h = Inches(1.7)
    for i, (t, c, bl) in enumerate(cards):
        x = Inches(0.6) + (i % 2) * (w + Inches(0.15))
        y = Inches(2.0) + (i // 2) * (h + Inches(0.18))
        panel(s, x, y, w, h, BG_PANEL, stroke=None)
        accent_bar(s, x, y, width=Inches(0.14), height=h, color=c)
        add_text(s, x + Inches(0.35), y + Inches(0.12),
                 w - Inches(0.5), Inches(0.4),
                 t, size=15, bold=True, color=c)
        add_bullets(s, x + Inches(0.35), y + Inches(0.6),
                    w - Inches(0.5), Inches(1.05),
                    bl, size=12, bullet_color=c, line_spacing=1.18)
    panel(s, Inches(0.6), Inches(5.85), Inches(12.15), Inches(1.15),
          BG_PANEL_ALT, stroke=None)
    accent_bar(s, Inches(0.6), Inches(5.85), width=Inches(0.14),
               height=Inches(1.15), color=ACCENT)
    add_text(s, Inches(0.95), Inches(5.96), Inches(11.7), Inches(0.45),
             "结论", size=12, bold=True, color=ACCENT, font=FONT_MONO)
    add_text(s, Inches(0.95), Inches(6.30), Inches(11.7), Inches(0.7),
             "单点工具不再够用——「能力底座 + 问题域复杂度」共同迫使我们重新思考算子工程栈。",
             size=14, bold=True, color=FG)
    footer(s, page, total, "Part A")


def slide_a5_output_form(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART A", ACCENT)
    slide_header(s, "A5 · 这份洞察的输出形态",
                 "六段式：案例 → 趋势 → 进展 → 难题 → 目标态 → 构建方案")

    steps = [
        ("案例", "7 个代表系统", ACCENT_BLUE),
        ("趋势", "T1–T6 收敛方向", ACCENT),
        ("进展", "G1–G6 已验证", ACCENT_ALT),
        ("难题", "H1–H9 待解", ACCENT_RED),
        ("目标态", "F1–F7 北极星", ACCENT_PURPLE),
        ("构建", "Arke 的回答", ACCENT),
    ]
    n = len(steps)
    sw = Inches(1.85)
    sh = Inches(1.6)
    sx = Inches(0.6)
    sy = Inches(2.3)
    for i, (t, d, c) in enumerate(steps):
        x = sx + i * (sw + Inches(0.18))
        panel(s, x, sy, sw, sh, BG_PANEL, stroke=None)
        accent_bar(s, x, sy, width=sw, height=Inches(0.08), color=c)
        add_text(s, x, sy + Inches(0.25), sw, Inches(0.45),
                 t, size=16, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), sy + Inches(0.8),
                 sw - Inches(0.2), Inches(0.7),
                 d, size=11, color=FG_MUTED, align=PP_ALIGN.CENTER)
        if i < n - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     x + sw + Inches(0.01),
                                     sy + Inches(0.65), Inches(0.16),
                                     Inches(0.3))
            _set_fill(arr, ACCENT)
            _no_line(arr)

    # value lines
    panel(s, Inches(0.6), Inches(4.5), Inches(12.15), Inches(2.4),
          BG_PANEL_ALT, stroke=None)
    accent_bar(s, Inches(0.6), Inches(4.5), width=Inches(0.14),
               height=Inches(2.4), color=ACCENT_ALT)
    add_text(s, Inches(0.95), Inches(4.65), Inches(11.7), Inches(0.4),
             "判断准则", size=12, bold=True, color=ACCENT_ALT,
             font=FONT_MONO)
    add_bullets(s, Inches(0.95), Inches(5.1), Inches(11.7), Inches(1.7),
                ["不止「看热闹」——每条趋势都要回答：工程上该不该跟、怎么跟",
                 "用六维坐标系评估：正确性 / 性能 / 成本 / 通用性 / 可移植性 / 可运营",
                 "目标是给出「可落地的范式判断」，而不是论文综述",
                 ],
                size=14, bullet_color=ACCENT_ALT, line_spacing=1.3)
    footer(s, page, total, "Part A")


# ============================================================
# PART B
# ============================================================
def slide_b1_panorama(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART B", ACCENT_BLUE)
    slide_header(s, "B1 · 全景图",
                 "LLM-driven kernel 工作链路：Spec ↔ Agent ↔ Toolchain ↔ HW")

    # 4 horizontal layers — all referenced cases are public industry systems
    layers = [
        ("① Spec / 抽象层", "kernel 表达 · 决策表达 · 知识表达",
         "代表抽象层：CuTe / Triton / 各家 DSL", ACCENT_PURPLE),
        ("② LLM Agent 层", "规划 · 决策 · 反思 · 回滚",
         "代表系统：KernelAgent / AVO / KernelFalcon", ACCENT),
        ("③ Toolchain 层", "compile · profile · verify · search",
         "代表系统：KernelEvolve / AutoKernel / K-Search", ACCENT_ALT),
        ("④ Hardware 层", "NVIDIA · Ascend · AMD · 多 NPU / DSA",
         "代表系统：AscendKernelGen / KernelEvolve (MTIA)", ACCENT_RED),
    ]
    lx = Inches(0.7)
    ly = Inches(2.0)
    lw = Inches(11.9)
    lh = Inches(1.05)
    for i, (t, d, c, col) in enumerate(layers):
        y = ly + i * (lh + Inches(0.12))
        panel(s, lx, y, lw, lh, BG_PANEL, stroke=None)
        accent_bar(s, lx, y, width=Inches(0.14), height=lh, color=col)
        add_text(s, lx + Inches(0.35), y + Inches(0.12),
                 Inches(3.4), Inches(0.4),
                 t, size=15, bold=True, color=col)
        add_text(s, lx + Inches(0.35), y + Inches(0.55),
                 Inches(4.5), Inches(0.45),
                 d, size=12, color=FG_MUTED)
        add_text(s, lx + Inches(5.5), y + Inches(0.34),
                 Inches(6.2), Inches(0.45),
                 d_case := c, size=11.5, color=FG, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        # vertical down arrow between layers
        if i < len(layers) - 1:
            ay = y + lh + Inches(0.0)
            ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    Inches(6.5), ay,
                                    Inches(0.3), Inches(0.12))
            _set_fill(ar, FG_MUTED)
            _no_line(ar)
    # caption
    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
             "案例只是证据，趋势才是主线——下一节用 7 个案例填到这张图的不同位置上。",
             size=12, color=FG_MUTED, italic=True)
    footer(s, page, total, "Part B")


def slide_b2_index(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART B", ACCENT_BLUE)
    slide_header(s, "B2 · 7 个代表案例索引",
                 "一页一个；含可归属性能信号（公开数据）")

    rows = [
        ("1", "KernelEvolve", "Meta · ISCA'26",
         "生产级异构 + RAG + 搜索",
         "KernelBench 250/250 · 最高 17× vs PyTorch",
         "T1/T3/T4", ACCENT),
        ("2", "KernelAgent / KernelFalcon", "PyTorch · Meta",
         "Deep Agents 分层 + 硬件信号",
         "L1/L2/L3 250/250 · 1.56× torch.compile · 89% H100 roofline",
         "T2/T3", ACCENT_BLUE),
        ("3", "AutoKernel", "RightNow AI",
         "autoresearch + Amdahl + 双后端",
         "H100：RMSNorm 5.29× · softmax 2.82× · cross-entropy 2.21×",
         "T1/T2", ACCENT_ALT),
        ("4", "K-Search", "UC Berkeley",
         "World-Model 规划 + 策略/实现解耦",
         "FlashInfer GQA/MLA/MoE 平均 2.10× · MoE 14.3× · TriMul SOTA 1030μs",
         "T1/T5", ACCENT_PURPLE),
        ("5", "AVO", "NVIDIA",
         "Agent-as-Variation-Operator",
         "Blackwell B200 · MHA 7-day 演化 · 超 cuDNN +3.5%, 超 FA-4 +10.5%",
         "T1/T2", ACCENT),
        ("6", "CuTeGen", "U. Toronto",
         "稳定抽象层 (CuTe) + 单 kernel 渐进精炼",
         "12 matmul + 14 activation kernels · token 成本可控",
         "T4", ACCENT_RED),
        ("7", "KernelGen-LM / AscendKernelGen", "PCL · 中山大学 · 华为",
         "领域数据 + 领域模型 (SFT + RLEF)",
         "Ascend NPU · L2 编译成功率 0%→95.5% (Pass@10)",
         "T6", ACCENT_ALT),
    ]
    # header row — reorganized to fit 6 columns
    hx = Inches(0.45)
    hw = Inches(12.45)
    rh = Inches(0.5)
    hy = Inches(1.95)
    panel(s, hx, hy, hw, rh, BG_PANEL_ALT, stroke=None)
    cols_x = [Inches(0.55), Inches(0.95), Inches(3.45),
              Inches(5.65), Inches(8.0), Inches(11.7)]
    cols_w = [Inches(0.4), Inches(2.5), Inches(2.2),
              Inches(2.35), Inches(3.7), Inches(1.2)]
    headers = ["#", "案例", "来源", "一句话定位",
               "公开性能 / 数据信号", "趋势"]
    for cx_, cw_, h in zip(cols_x, cols_w, headers):
        add_text(s, cx_, hy + Inches(0.11), cw_, Inches(0.3),
                 h, size=11, bold=True, color=ACCENT_ALT,
                 font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
    # rows
    for i, (n, name, src, desc, signal, t, col) in enumerate(rows):
        y = hy + Inches(0.55) + i * Inches(0.62)
        panel(s, hx, y, hw, Inches(0.58),
              BG_PANEL if i % 2 == 0 else BG_PANEL_DEEP, stroke=None)
        accent_bar(s, hx, y, width=Inches(0.06),
                   height=Inches(0.58), color=col)
        add_text(s, cols_x[0], y + Inches(0.14), cols_w[0], Inches(0.3),
                 n, size=11, bold=True, color=col, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cols_x[1], y + Inches(0.14), cols_w[1], Inches(0.3),
                 name, size=11, bold=True, color=FG,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cols_x[2], y + Inches(0.14), cols_w[2], Inches(0.3),
                 src, size=9.5, color=FG_MUTED,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cols_x[3], y + Inches(0.14), cols_w[3], Inches(0.3),
                 desc, size=10, color=FG, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cols_x[4], y + Inches(0.14), cols_w[4], Inches(0.3),
                 signal, size=9.5, color=ACCENT_BLUE,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
        add_text(s, cols_x[5], y + Inches(0.14), cols_w[5], Inches(0.3),
                 t, size=10, bold=True, color=col,
                 font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
    source_note(s, "Source: 各案例 arXiv / 官方 blog / GitHub repo (详见 References)；性能为各方公布数据，硬件 / shape 各异不可直接横比。")
    footer(s, page, total, "Part B")


def slide_b2_case(prs, page, total, idx, name, source, color,
                  what, key_tech, trend, gap, metrics, source_cite):
    """Single case page using fixed template — now with公开 metrics row."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART B", ACCENT_BLUE)
    slide_header(s, f"B2.{idx} · 案例 {idx}/7 · {name}",
                 f"{source}　|　代表趋势：{trend}")

    # Metrics strip (public, attributable data)
    panel(s, Inches(0.6), Inches(1.95), Inches(12.15), Inches(0.7),
          BG_PANEL_ALT, stroke=None)
    accent_bar(s, Inches(0.6), Inches(1.95), width=Inches(0.14),
               height=Inches(0.7), color=ACCENT_BLUE)
    add_text(s, Inches(0.95), Inches(2.02), Inches(2.4), Inches(0.32),
             "公开数据 / 代表性能", size=10, bold=True,
             color=ACCENT_BLUE, font=FONT_MONO,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.4), Inches(2.02), Inches(9.3), Inches(0.55),
             metrics, size=11.5, color=FG, anchor=MSO_ANCHOR.MIDDLE)

    # 2x2 grid (slightly compacted to make room for metrics strip)
    grid = [
        ("做了什么", what, ACCENT),
        ("关键技术点", key_tech, ACCENT_BLUE),
        ("代表的趋势", trend_long(trend), ACCENT_ALT),
        ("它没解决什么", gap, ACCENT_RED),
    ]
    cw = Inches(5.95)
    ch = Inches(2.05)
    cx = Inches(0.6)
    cy = Inches(2.8)
    for i, (t, items, col) in enumerate(grid):
        x = cx + (i % 2) * (cw + Inches(0.15))
        y = cy + (i // 2) * (ch + Inches(0.15))
        panel(s, x, y, cw, ch, BG_PANEL, stroke=None)
        accent_bar(s, x, y, width=Inches(0.14), height=ch, color=col)
        add_text(s, x + Inches(0.3), y + Inches(0.1),
                 cw - Inches(0.5), Inches(0.36),
                 t, size=14, bold=True, color=col)
        add_bullets(s, x + Inches(0.3), y + Inches(0.55),
                    cw - Inches(0.5), ch - Inches(0.65),
                    items, size=11, bullet_color=col, line_spacing=1.2)
    # bottom case-color stripe
    accent_bar(s, Inches(0.6), Inches(6.95), width=Inches(12.15),
               height=Inches(0.05), color=color)
    source_note(s, f"Source: {source_cite}")
    footer(s, page, total, f"Part B · 案例 {idx}/7")


def trend_long(code: str) -> list:
    mapping = {
        "T1": "搜索化：one-shot → 演化 / 规划 / MCTS",
        "T2": "工具化：tool-use + harness + 可复现评测",
        "T3": "硬件信号化：profile / roofline / NCU 入循环",
        "T4": "抽象层选择：稳定 + 可迭代的目标层",
        "T5": "知识资产化：rationale / playbook / 轨迹回放",
        "T6": "领域模型化：SFT + RLEF + 后训练",
    }
    parts = [p.strip() for p in code.replace("/", " / ").split(" / ")
             if p.strip()]
    return [mapping[p] for p in parts if p in mapping]


def slide_b3_trends(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART B", ACCENT_BLUE)
    slide_header(s, "B3 · 趋势归纳",
                 "业界正在收敛的 6 个方向（每条标注代表案例）")

    trends = [
        ("T1", "搜索化",
         "从 one-shot → 演化 / 规划 / MCTS 等结构化搜索",
         "代表：KernelEvolve · K-Search · AVO",
         ACCENT),
        ("T2", "工具化",
         "tool-use + harness + 可复现评测，把 LLM 关进编译器 / profiler 控制室",
         "代表：KernelAgent · AutoKernel",
         ACCENT_BLUE),
        ("T3", "硬件信号化",
         "profile / roofline / NCU 反馈直接进入决策循环",
         "代表：KernelAgent · AVO · KernelEvolve",
         ACCENT_ALT),
        ("T4", "抽象层选择",
         "放弃直写 PTX，押注稳定可迭代的抽象层 (Triton / CuTe / DSL)",
         "代表：CuTeGen · KernelEvolve",
         ACCENT_PURPLE),
        ("T5", "知识资产化",
         "RAG + playbook + decision rationale + 轨迹 → 经验从手感到资产",
         "代表：KernelEvolve (RAG) · K-Search (world model)",
         ACCENT_RED),
        ("T6", "领域模型化",
         "SFT + RLEF + 后训练，让小 / 中模型在 kernel 域具备专家级判断",
         "代表：KernelGen-LM / AscendKernelGen",
         ACCENT),
    ]
    cw = Inches(4.0)
    ch = Inches(2.4)
    for i, (code, t, d, case_ref, c) in enumerate(trends):
        x = Inches(0.6) + (i % 3) * (cw + Inches(0.1))
        y = Inches(2.0) + (i // 3) * (ch + Inches(0.18))
        panel(s, x, y, cw, ch, BG_PANEL, stroke=None)
        accent_bar(s, x, y, width=Inches(0.14), height=ch, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.13),
                 Inches(0.9), Inches(0.4),
                 code, size=18, bold=True, color=c, font=FONT_MONO)
        add_text(s, x + Inches(1.1), y + Inches(0.18),
                 cw - Inches(1.3), Inches(0.4),
                 t, size=15, bold=True, color=FG)
        add_text(s, x + Inches(0.3), y + Inches(0.65),
                 cw - Inches(0.4), Inches(1.0),
                 d, size=11.5, color=FG_MUTED)
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x + Inches(0.3),
                                  y + ch - Inches(0.55),
                                  cw - Inches(0.5), Inches(0.4))
        chip.adjustments[0] = 0.4
        _set_fill(chip, BG_PANEL_DEEP)
        _set_line(chip, c, 0.75)
        add_text(s, x + Inches(0.3), y + ch - Inches(0.55),
                 cw - Inches(0.5), Inches(0.4),
                 case_ref, size=10, color=c,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
                 font=FONT_MONO)
    glossary_hint(s, "→ Glossary P03 · 术语 RAG / harness / RLEF / NCU")
    source_note(s, "← 与 A2.5 范式判断对应：T1/T2/T3/T4/T5/T6 分别延伸自 ❶❷❸❹❺❻ 六条预判 (P07)")
    footer(s, page, total, "Part B")


def slide_b4_progress(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART B", ACCENT_BLUE)
    slide_header(s, "B4 · 已有的「好进展」",
                 "What's working — 客观陈述案例已经验证可行的部分（含真实数据）")

    progs = [
        ("G1", "正确性门禁",
         "数值等价 (numerical equivalence) 已是业界标配",
         "KernelAgent: KernelBench L1/L2/L3 250/250 全 PASS"),
        ("G2", "并行探索 + 早停 + 预算",
         "multi-worker + early-win 显著降低 token/wall-clock 成本",
         "AutoKernel: 长时无人值守演化 · KernelEvolve: tree search"),
        ("G3", "硬件 profile 反馈",
         "进入决策循环后稳定带来双位数性能提升",
         "KernelAgent: 89% H100 roofline · 1.56× torch.compile"),
        ("G4", "稳定抽象层",
         "Triton / CuTe 等上的 LLM 生成已能落地生产",
         "CuTeGen: 12 matmul + 14 act kernels · KernelEvolve: 多 DSL"),
        ("G5", "RAG + 知识库 + 后训练",
         "让「专家直觉」开始可累计、可迁移",
         "KernelEvolve: 硬件知识库 RAG · KernelGen-LM: SFT + RLEF"),
        ("G6", "Agent 编排",
         "分层 + 规划 + 回滚让长会话稳定性显著改善",
         "AVO: B200 上 MHA 7-day 演化超 FA-4 +10.5%"),
    ]
    rh = Inches(0.82)
    ry = Inches(2.0)
    rx = Inches(0.6)
    rw = Inches(12.15)
    for i, (code, t, d, evidence) in enumerate(progs):
        y = ry + i * (rh + Inches(0.05))
        panel(s, rx, y, rw, rh,
              BG_PANEL if i % 2 == 0 else BG_PANEL_DEEP, stroke=None)
        accent_bar(s, rx, y, width=Inches(0.14), height=rh, color=ACCENT)
        add_text(s, rx + Inches(0.4), y + Inches(0.2),
                 Inches(0.9), Inches(0.45),
                 code, size=16, bold=True, color=ACCENT, font=FONT_MONO)
        add_text(s, rx + Inches(1.45), y + Inches(0.13),
                 Inches(3.3), Inches(0.4),
                 t, size=14, bold=True, color=FG,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(1.45), y + Inches(0.43),
                 Inches(3.3), Inches(0.36),
                 d, size=10, color=FG_MUTED,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(4.95), y + Inches(0.18),
                 Inches(7.0), Inches(0.5),
                 "证据：" + evidence, size=11, color=ACCENT_BLUE,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
    source_note(s, "Source: 各案例公开 paper / blog / repo（详见 B2 与 References）")
    footer(s, page, total, "Part B")


def slide_b5_challenges(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART B", ACCENT_BLUE)
    slide_header(s, "B5 · 仍待解决的「关键技术难题」",
                 "What's hard — H1–H9（统一编号 · 每条带业界证据）")

    chs = [
        ("H1", "策略可迁移性",
         "策略仍埋在自由代码里，难审计 / 回滚 / 跨 kernel 迁移",
         "证据：KernelAgent / AVO 输出仍是 Triton/CUDA",
         ACCENT),
        ("H2", "多级验证",
         "事后对拍 → 静态 / 数值 / 性能的可剪枝多级门禁",
         "证据：各家自建 harness，无统一协议",
         ACCENT_BLUE),
        ("H3", "动态 / 符号 shape",
         "动态 shape 下泛化与不退化同时保证",
         "证据：shape guard / specialization 各家自做",
         ACCENT_ALT),
        ("H4", "跨硬件迁移",
         "NVIDIA → Ascend / AMD / NPU 策略与知识复用",
         "证据：KernelGen-LM 切硬件需重训",
         ACCENT_PURPLE),
        ("H5", "模型级自治",
         "从单 kernel 到模型级 (bottleneck → 优化 → 回归) 闭环",
         "证据：AutoKernel 仅做单 kernel；端到端少",
         ACCENT_RED),
        ("H6", "Token / 预算治理",
         "长会话中预算 / 上下文 / 稳定性可治理",
         "证据：AVO 7-day 演化的成本控制是开放问题",
         ACCENT),
        ("H7", "后端天花板",
         "Triton 抽象封顶 → 何时下沉 MLIR / LLVM 解锁深决策",
         "证据：CuTeGen 选 CuTe 才能逼近 cuBLAS",
         ACCENT_BLUE),
        ("H8", "数据稀缺 / 领域模型",
         "kernel × HW × profile 三元组数据稀缺",
         "证据：AscendKernelGen 自造 Ascend-CoT",
         ACCENT_ALT),
        ("H9", "评测可复现",
         "缺统一 benchmark / 形状层级 / 基线协议",
         "证据：各家性能数据硬件 / shape 各异，难横比",
         ACCENT_RED),
    ]
    cw = Inches(4.0)
    ch = Inches(1.55)
    for i, (code, t, d, ev, c) in enumerate(chs):
        x = Inches(0.6) + (i % 3) * (cw + Inches(0.1))
        y = Inches(2.0) + (i // 3) * (ch + Inches(0.13))
        panel(s, x, y, cw, ch, BG_PANEL, stroke=None)
        accent_bar(s, x, y, width=Inches(0.14), height=ch, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.1),
                 Inches(0.9), Inches(0.36),
                 code, size=15, bold=True, color=c, font=FONT_MONO)
        add_text(s, x + Inches(1.05), y + Inches(0.11),
                 cw - Inches(1.2), Inches(0.36),
                 t, size=13, bold=True, color=FG)
        add_text(s, x + Inches(0.3), y + Inches(0.5),
                 cw - Inches(0.4), Inches(0.5),
                 d, size=10.5, color=FG_MUTED)
        add_text(s, x + Inches(0.3), y + ch - Inches(0.42),
                 cw - Inches(0.4), Inches(0.36),
                 ev, size=9, color=ACCENT_BLUE,
                 italic=True, font=FONT_MONO)
    footer(s, page, total, "Part B")


# ============================================================
# TARGET-STATE
# ============================================================
def slide_s1_north_star(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "TARGET", ACCENT_ALT)
    slide_header(s, "S1 · 北极星",
                 "未来大模型算子生成 / 调优技术的目标态")

    # north-star statement panel
    panel(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.6),
          BG_PANEL, stroke=None)
    accent_bar(s, Inches(0.9), Inches(2.0), width=Inches(0.16),
               height=Inches(2.6), color=ACCENT_ALT)
    add_text(s, Inches(1.2), Inches(2.15), Inches(11), Inches(0.4),
             "目标态 · ONE-LINER", size=12, bold=True, color=ACCENT_ALT,
             font=FONT_MONO)
    add_runs(s, Inches(1.2), Inches(2.55), Inches(11), Inches(2.0),
             [
                 {"text": "让 ", "size": 22, "color": FG_MUTED},
                 {"text": "任意 LLM ", "size": 22, "bold": True,
                  "color": ACCENT},
                 {"text": "在", "size": 22, "color": FG_MUTED},
                 {"text": "最小 token 预算", "size": 22, "bold": True,
                  "color": ACCENT_ALT},
                 {"text": "下产出 ", "size": 22, "color": FG_MUTED},
                 {"text": "正确、最优、可迁移、可演进 ", "size": 22,
                  "bold": True, "color": FG},
                 {"text": "的算子。", "size": 22, "color": FG_MUTED},
             ], line_spacing=1.4)

    # 5 keywords strip
    kws = [
        ("语义/策略分层", ACCENT),
        ("有界动作空间", ACCENT_BLUE),
        ("多级编译器验证", ACCENT_ALT),
        ("结构化经验沉淀", ACCENT_PURPLE),
        ("跨硬件零迁移成本", ACCENT_RED),
    ]
    kw_y = Inches(5.0)
    kw_h = Inches(1.6)
    sw = Inches(2.35)
    sx = Inches(0.65)
    for i, (kw, c) in enumerate(kws):
        x = sx + i * (sw + Inches(0.07))
        panel(s, x, kw_y, sw, kw_h, BG_PANEL_ALT, stroke=None)
        accent_bar(s, x, kw_y, width=sw, height=Inches(0.08), color=c)
        add_text(s, x + Inches(0.15), kw_y + Inches(0.45),
                 sw - Inches(0.3), Inches(0.55),
                 f"#{i+1:02d}", size=12, color=c, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), kw_y + Inches(0.85),
                 sw - Inches(0.3), Inches(0.6),
                 kw, size=14, bold=True, color=FG,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, page, total, "Target State")


def slide_s2_features(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "TARGET", ACCENT_ALT)
    slide_header(s, "S2 · 7 大技术特征 F1–F7",
                 "由趋势 T1–T6 + 难题 H1–H9 收敛而来——目标态的 spec")

    feats = [
        ("F1", "语义与策略分离", "T4 · T5 / H1",
         "「做什么」是不可变数学；「怎么做」是可搜索决策。", ACCENT),
        ("F2", "有界动作空间", "T2 / H2 · H6",
         "LLM 不写自由代码，只在编译器枚举的合法动作上选择。",
         ACCENT_BLUE),
        ("F3", "多层 IR + 多级验证", "T2 · T4 / H2 · H7",
         "高层语义 → 策略 → 调度 → 指令；每层静态/数值/性能门禁。",
         ACCENT_ALT),
        ("F4", "硬件信号闭环", "T3 / H3 · H4",
         "profile / roofline / 占用率 / 寄存器压力进入决策循环。",
         ACCENT_PURPLE),
        ("F5", "结构化经验资产", "T5 / H1 · H8",
         "rationale / playbook / 轨迹 / 决策树成为一等公民。",
         ACCENT_RED),
        ("F6", "跨硬件统一表达", "T4 / H4",
         "同一份语义、同一套决策原语，target-aware 适配多硬件。",
         ACCENT),
        ("F7", "可复现评测协议", "T1 · T2 / H9",
         "统一形状层级 / baseline / 报表 schema，跨工作可比。",
         ACCENT_ALT),
    ]
    cw = Inches(4.0)
    ch = Inches(2.3)
    for i, (code, t, tag, d, c) in enumerate(feats):
        if i < 6:
            x = Inches(0.6) + (i % 3) * (cw + Inches(0.1))
            y = Inches(2.0) + (i // 3) * (ch + Inches(0.16))
        else:
            x = Inches(4.7)
            y = Inches(2.0) + 2 * (ch + Inches(0.16))
        panel(s, x, y, cw, ch, BG_PANEL, stroke=None)
        accent_bar(s, x, y, width=Inches(0.14), height=ch, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.12),
                 Inches(0.85), Inches(0.4),
                 code, size=16, bold=True, color=c, font=FONT_MONO)
        add_text(s, x + Inches(1.1), y + Inches(0.13),
                 cw - Inches(1.2), Inches(0.4),
                 t, size=14, bold=True, color=FG)
        add_text(s, x + Inches(0.3), y + Inches(0.6),
                 cw - Inches(0.4), Inches(0.4),
                 tag, size=11, color=ACCENT_ALT, font=FONT_MONO)
        add_text(s, x + Inches(0.3), y + Inches(1.0),
                 cw - Inches(0.4), ch - Inches(1.05),
                 d, size=11.5, color=FG_MUTED)
    footer(s, page, total, "Target State")


def slide_s3_implications(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "TARGET", ACCENT_ALT)
    slide_header(s, "S3 · 对通用编程的启示",
                 "把 kernel 域的洞察推广到更广的软件工程")

    items = [
        ("01", "编程语言会双层化",
         "What / How 显式分离——适用于性能敏感 + 正确性敏感的所有领域：" \
         "数据库执行计划 / 分布式调度 / 系统调优 / 形式化证明。",
         ACCENT),
        ("02", "模型不是代码生成器，而是决策代理",
         "通用编程的下一步：在结构化决策空间里做选择；" \
         "编译器 / 运行时 / 类型系统是它的「环境」。",
         ACCENT_BLUE),
        ("03", "编译器 / 工具链承担「可验证 RL 环境」的角色",
         "静态 / 数值 / 性能 三级验证不仅是 kernel 域的工程实践，" \
         "更是让 LLM 在通用程序合成中「安全地试错」的范式。",
         ACCENT_ALT),
        ("04", "经验资产化是 LLM 工程的真正护城河",
         "decision rationale / 轨迹 / 决策树可被检索、训练、跨团队复用——" \
         "比生成的代码本身更有长期价值。",
         ACCENT_PURPLE),
        ("05", "跨平台 / 跨架构「零迁移成本」正在变得可能",
         "抽象层稳定 + 策略可迁移 → 写一次跑多硬件，从图编译器口号变 LLM 时代现实。",
         ACCENT_RED),
        ("06", "评测协议是行业的公共物品",
         "与其每家堆 demo，不如共同推动可复现的评测协议——" \
         "这是行业从「PR 驱动」走向「工程化」的前提。",
         ACCENT),
    ]
    rh = Inches(0.78)
    ry = Inches(2.0)
    rx = Inches(0.6)
    rw = Inches(12.15)
    for i, (n, t, d, c) in enumerate(items):
        y = ry + i * (rh + Inches(0.04))
        panel(s, rx, y, rw, rh,
              BG_PANEL if i % 2 == 0 else BG_PANEL_DEEP, stroke=None)
        accent_bar(s, rx, y, width=Inches(0.14), height=rh, color=c)
        add_text(s, rx + Inches(0.4), y + Inches(0.18),
                 Inches(0.7), Inches(0.45),
                 n, size=16, bold=True, color=c, font=FONT_MONO)
        add_text(s, rx + Inches(1.25), y + Inches(0.13),
                 Inches(4.0), Inches(0.45),
                 t, size=13.5, bold=True, color=FG,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(5.4), y + Inches(0.13),
                 Inches(6.5), Inches(0.55),
                 d, size=11, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total, "Target State")


def slide_s4_bridge(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "TARGET", ACCENT_ALT)
    slide_header(s, "S4 · 从目标态走向构建",
                 "Agent 是手段，IR + 验证 + 编译栈才是底盘")

    # main statement
    panel(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.0),
          BG_PANEL, stroke=None)
    accent_bar(s, Inches(0.9), Inches(2.0), width=Inches(0.16),
               height=Inches(2.0), color=ACCENT_ALT)
    add_runs(s, Inches(1.2), Inches(2.2), Inches(11), Inches(1.7),
             [
                 {"text": "下一章给出一个 ", "size": 22, "color": FG_MUTED},
                 {"text": "可落地的总体设计", "size": 22, "bold": True,
                  "color": ACCENT},
                 {"text": "，", "size": 22, "color": FG_MUTED},
                 {"text": "而不是「再堆一个 agent」", "size": 22,
                  "bold": True, "color": ACCENT_ALT},
                 {"text": "。", "size": 22, "color": FG_MUTED},
             ], line_spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)

    # what next
    items = [
        ("设计逻辑", "5 条公理，每条对应 F1–F7 的一个或多个特征",
         ACCENT),
        ("整体架构", "一张系统总图——四件套 + Benchmark 的协作关系",
         ACCENT_BLUE),
        ("四件套 + Benchmark",
         "Language · IR · Compiler · Agent · Benchmark 各一页",
         ACCENT_ALT),
        ("最终对齐", "F1–F7 ↔ Arke 组件一一对应，闭环全 PPT 叙事",
         ACCENT_PURPLE),
    ]
    cw = Inches(2.95)
    cx = Inches(0.6)
    cy = Inches(4.5)
    for i, (t, d, c) in enumerate(items):
        x = cx + i * (cw + Inches(0.1))
        panel(s, x, cy, cw, Inches(2.2), BG_PANEL_ALT, stroke=None)
        accent_bar(s, x, cy, width=cw, height=Inches(0.08), color=c)
        add_text(s, x + Inches(0.2), cy + Inches(0.25),
                 cw - Inches(0.3), Inches(0.45),
                 t, size=14, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), cy + Inches(0.85),
                 cw - Inches(0.3), Inches(1.2),
                 d, size=12, color=FG_MUTED, align=PP_ALIGN.CENTER)
    footer(s, page, total, "Target State")


# ============================================================
# PART C
# ============================================================
def slide_c1_axioms(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART C", ACCENT_PURPLE)
    slide_header(s, "C1 · 设计逻辑",
                 "Arke 凭什么相信能逼近目标态——5 条设计公理")

    axioms = [
        ("公理 1", "算子是「可数学验证」的对象",
         "→ 必然要 语义/策略分离", "F1", ACCENT),
        ("公理 2", "LLM 是决策者，不是代码生成器",
         "→ 必然要 Bounded Action Space", "F2", ACCENT_BLUE),
        ("公理 3", "编译器是「可验证 RL 环境」",
         "→ 必然要 多层 IR + 多级验证", "F3", ACCENT_ALT),
        ("公理 4", "经验必须能沉淀",
         "→ 必然要 @rationale 作为一等公民", "F5", ACCENT_PURPLE),
        ("公理 5", "性能上限来自硬件，不是抽象层",
         "→ 必然要 渐进下沉到更深的 IR", "F3 + F6", ACCENT_RED),
    ]
    rh = Inches(0.85)
    ry = Inches(2.0)
    rx = Inches(0.6)
    rw = Inches(12.15)
    for i, (n, t, d, f, c) in enumerate(axioms):
        y = ry + i * (rh + Inches(0.06))
        panel(s, rx, y, rw, rh,
              BG_PANEL if i % 2 == 0 else BG_PANEL_DEEP, stroke=None)
        accent_bar(s, rx, y, width=Inches(0.14), height=rh, color=c)
        add_text(s, rx + Inches(0.4), y + Inches(0.18),
                 Inches(1.2), Inches(0.5),
                 n, size=14, bold=True, color=c, font=FONT_MONO)
        add_text(s, rx + Inches(1.7), y + Inches(0.13),
                 Inches(4.5), Inches(0.5),
                 t, size=14, bold=True, color=FG,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(6.3), y + Inches(0.13),
                 Inches(4.4), Inches(0.55),
                 d, size=12, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE)
        # F-tag chip
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  rx + Inches(10.85),
                                  y + Inches(0.22),
                                  Inches(1.15), Inches(0.42))
        chip.adjustments[0] = 0.4
        _set_fill(chip, c)
        _no_line(chip)
        add_text(s, rx + Inches(10.85), y + Inches(0.22),
                 Inches(1.15), Inches(0.42),
                 f, size=11, bold=True, color=BG_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_MONO)

    add_text(s, Inches(0.6), Inches(6.78), Inches(12), Inches(0.3),
             "设计公理 → 体系结构 → 四件套实现，逻辑自上而下闭合。",
             size=12, color=FG_MUTED, italic=True, align=PP_ALIGN.CENTER)
    source_note(s, "← 5 条公理直接回应 A2.5 范式判断：公理 1↔❺ · 公理 2↔❹ · 公理 3↔❷❺ · 公理 4↔❺❻ · 公理 5↔❻ (P07)")
    footer(s, page, total, "Part C")


def slide_c2_architecture(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART C", ACCENT_PURPLE)
    slide_header(s, "C2 · 整体架构",
                 "唯一的系统总图——一眼看懂 Arke 怎么承载目标态")

    # Layered architecture, 4 boxes + HW
    layers = [
        ("①", "Arke Language (.ak)",
         "kernel { 语义 }   strategy { 决策 + rationale }",
         "F1 · F5", ACCENT),
        ("②", "Arke IR — 4 层架构",
         "L4 SemanticIR · L3 StrategyIR · L2 ScheduleIR · L1 InstructionIR",
         "F2 · F3 · F6", ACCENT_BLUE),
        ("③", "Compiler Toolchain",
         "OpRegistry · PassPipeline · V0/V1/V2 · Backend (Triton/MLIR/LLVM)",
         "F3 · F4 · F6", ACCENT_ALT),
        ("④", "Agent Engineering",
         "Tool-use 协议 · Bounded Action · Budget · Checkpoint/Rollback",
         "F2 · F4 · F5", ACCENT_PURPLE),
    ]
    lx = Inches(0.7)
    ly = Inches(1.95)
    lw = Inches(11.9)
    lh = Inches(1.0)
    for i, (n, t, d, ftag, col) in enumerate(layers):
        y = ly + i * (lh + Inches(0.12))
        panel(s, lx, y, lw, lh, BG_PANEL, stroke=None)
        accent_bar(s, lx, y, width=Inches(0.14), height=lh, color=col)
        add_text(s, lx + Inches(0.4), y + Inches(0.15),
                 Inches(0.7), Inches(0.5),
                 n, size=22, bold=True, color=col, font=FONT_MONO)
        add_text(s, lx + Inches(1.2), y + Inches(0.13),
                 Inches(4.6), Inches(0.45),
                 t, size=14.5, bold=True, color=FG)
        add_text(s, lx + Inches(1.2), y + Inches(0.55),
                 Inches(8.5), Inches(0.45),
                 d, size=11.5, color=FG_MUTED)
        # F-tag chip on the right
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  lx + Inches(10.45),
                                  y + Inches(0.3),
                                  Inches(1.3), Inches(0.42))
        chip.adjustments[0] = 0.4
        _set_fill(chip, col)
        _no_line(chip)
        add_text(s, lx + Inches(10.45), y + Inches(0.3),
                 Inches(1.3), Inches(0.42),
                 ftag, size=10, bold=True, color=BG_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_MONO)
        if i < len(layers) - 1:
            ay = y + lh + Inches(0.0)
            ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    Inches(6.5), ay,
                                    Inches(0.3), Inches(0.12))
            _set_fill(ar, FG_MUTED)
            _no_line(ar)

    # HW row
    hw_y = ly + 4 * (lh + Inches(0.12))
    panel(s, lx, hw_y, lw, Inches(0.7), BG_PANEL_ALT, stroke=None)
    accent_bar(s, lx, hw_y, width=Inches(0.14), height=Inches(0.7),
               color=ACCENT_RED)
    add_text(s, lx + Inches(0.4), hw_y + Inches(0.18),
             Inches(11), Inches(0.4),
             "GPU / NPU 执行:　 NVIDIA · Ascend · AMD · 未来 NPU / DSA",
             size=14, bold=True, color=FG, anchor=MSO_ANCHOR.MIDDLE)

    # bottom caption: benchmark cross-cut
    bm_y = hw_y + Inches(0.85)
    panel(s, lx, bm_y, lw, Inches(0.4),
          BG_PANEL_ALT, stroke=None)
    add_text(s, lx + Inches(0.3), bm_y, lw, Inches(0.4),
             "⑤ Benchmark 体系（横切的「度量层」）：BL × OT × ST × L  ·  多 baseline  ·  统一 reporting schema  ·  ⇒ F7",
             size=11.5, bold=True, color=ACCENT_ALT,
             font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total, "Part C")


def _component_slide(prs, page, total, idx, name, oneliner,
                     color, points, trend_resp, hard_resp, f_resp):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART C", ACCENT_PURPLE)
    slide_header(s, f"C{2+idx} · 组件 {idx} · {name}",
                 oneliner)

    # left: design points
    panel(s, Inches(0.6), Inches(2.0), Inches(7.4), Inches(4.85),
          BG_PANEL, stroke=None)
    accent_bar(s, Inches(0.6), Inches(2.0), width=Inches(0.14),
               height=Inches(4.85), color=color)
    add_text(s, Inches(0.95), Inches(2.15), Inches(7), Inches(0.5),
             "设计要点", size=14, bold=True, color=color, font=FONT_MONO)
    add_bullets(s, Inches(0.95), Inches(2.7), Inches(7.0), Inches(4.0),
                points, size=13, bullet_color=color, line_spacing=1.32)

    # right: 3 mapping panels (trend / challenges / F)
    rx = Inches(8.2)
    rw = Inches(4.55)
    rows = [
        ("呼应趋势 / 进展", trend_resp, ACCENT),
        ("回应难题", hard_resp, ACCENT_RED),
        ("承载目标态特征", f_resp, ACCENT_ALT),
    ]
    rh = Inches(1.55)
    ry = Inches(2.0)
    for i, (t, items, c) in enumerate(rows):
        y = ry + i * (rh + Inches(0.12))
        panel(s, rx, y, rw, rh, BG_PANEL_ALT, stroke=None)
        accent_bar(s, rx, y, width=rw, height=Inches(0.06), color=c)
        add_text(s, rx + Inches(0.25), y + Inches(0.14),
                 rw - Inches(0.4), Inches(0.4),
                 t, size=12, bold=True, color=c, font=FONT_MONO)
        add_bullets(s, rx + Inches(0.25), y + Inches(0.55),
                    rw - Inches(0.4), rh - Inches(0.65),
                    items, size=11.5, bullet_color=c, line_spacing=1.2,
                    marker_l0="• ")
    footer(s, page, total, f"Part C · 组件 {idx}")


def slide_c3_language(prs, page, total):
    _component_slide(
        prs, page, total, idx=1,
        name="Arke Language (.ak)",
        oneliner="LLM 与人共同书写的「算子规约语言」",
        color=ACCENT,
        points=[
            "kernel { ... } 写纯数学语义（无 tile / 无 thread / 无 memory）",
            "strategy { ... } 写离散决策 + @rationale，每条决策可枚举/回滚",
            "where 子句承载符号维度 / 动态 shape 约束",
            ("token 效率：相对手写 Triton 节省一个数量级", 1),
            ("可读性：人和 LLM 同时可读、可写、可审计", 1),
            "通用算子抽象：算子通过 OpRegistry 注册，而非内建关键字",
            "single source of truth：.ak 是规范，JSON IR 仅作序列化",
        ],
        trend_resp=[
            "T4 · 稳定抽象层",
            "T5 · 知识资产化（@rationale）",
        ],
        hard_resp=[
            "H1 · 策略可迁移：策略从代码中抽出",
            "H3 · 动态 shape：where + symbolic dim",
        ],
        f_resp=[
            "F1 · 语义 / 策略分离",
            "F5 · 结构化经验资产",
        ],
    )


def slide_c4_ir(prs, page, total):
    _component_slide(
        prs, page, total, idx=2,
        name="Arke IR (多层架构)",
        oneliner="LLM-Native 中间表示——4 层各有「LLM 参与度」",
        color=ACCENT_BLUE,
        points=[
            "Layer 4 · SemanticIR：纯数学，硬件无关，作为正确性参考",
            "Layer 3 · StrategyIR：决策 + @rationale，LLM 主战场，可 checkpoint / rollback",
            "Layer 2 · ScheduleIR：thread / block / warp 映射，编译器主导",
            "Layer 1 · InstructionIR：近 LLVM IR，全自动",
            ("决策空间是有界的，可枚举、可剪枝", 1),
            ("每层之间下沉，配合 V0/V1/V2 三级验证", 1),
            "与 MLIR / LLVM 互操作而非依赖：可下沉、可绕过、可演进",
            "符号维度（symbolic dim）贯穿 4 层，动态 shape 一等公民",
        ],
        trend_resp=[
            "T2 · 工具化",
            "T4 · 抽象层选择",
        ],
        hard_resp=[
            "H2 · 多级验证",
            "H7 · 后端天花板",
            "H3 · 动态 / 符号 shape",
        ],
        f_resp=[
            "F2 · 有界动作空间",
            "F3 · 多层 IR + 多级验证",
            "F6 · 跨硬件统一表达",
        ],
    )


def slide_c5_compiler(prs, page, total):
    _component_slide(
        prs, page, total, idx=3,
        name="Arke Compiler Toolchain",
        oneliner='"Compiler-as-Verifier"——不优化，只把关',
        color=ACCENT_ALT,
        points=[
            "OpRegistry：算子单一事实源，新增算子的成本最小化",
            "Pass Pipeline：可组合的 IR 变换基础设施",
            "多级验证流水：",
            ("V0 静态 (<1ms)：结构 / 类型 / 约束", 1),
            ("V1 数值：与 NumPy / 参考实现对拍", 1),
            ("V2 性能：编译 + profile，进入 LLM 反馈循环", 1),
            "Backend 抽象：Triton / MLIR / LLVM 可插拔，渐进下沉",
            "硬件信号：roofline / occupancy / 寄存器压力可被 Agent 引用",
        ],
        trend_resp=[
            "G1 · 正确性门禁",
            "G3 · 硬件信号",
            "G4 · 稳定抽象层",
        ],
        hard_resp=[
            "H2 · 多级验证",
            "H4 · 跨硬件迁移",
            "H7 · 后端天花板",
        ],
        f_resp=[
            "F3 · 多级验证",
            "F4 · 硬件信号闭环",
            "F6 · 跨硬件统一表达",
        ],
    )


def slide_c6_agent(prs, page, total):
    _component_slide(
        prs, page, total, idx=4,
        name="Arke Agent Engineering",
        oneliner="让任意 LLM 都能稳定开 kernel 的「工程化运行时」",
        color=ACCENT_PURPLE,
        points=[
            "统一工具协议：analyze_compute / list_legal_actions / apply_decision",
            ("verify_correctness / compile_and_profile / observe", 1),
            ("checkpoint / rollback / nudge / compact", 1),
            "Bounded Action Space：动作来自编译器枚举，杜绝幻觉式自由代码",
            "预算 / 节流：budget · prompt cache · segmented compact",
            "双模集成：",
            ("Mode A 内置 agent（CLI / Python API）", 1),
            ("Mode B 外部 agent（Cursor / Claude Code 等可直接驱动）", 1),
            "轨迹 + rationale 沉淀：每次会话产出可检索 / 可训练资产",
        ],
        trend_resp=[
            "T1 · 搜索化",
            "T2 · 工具化",
            "G6 · Agent 编排",
        ],
        hard_resp=[
            "H5 · 模型级闭环",
            "H6 · token / 预算治理",
            "H8 · 数据 / 领域模型",
        ],
        f_resp=[
            "F2 · 有界动作空间",
            "F4 · 硬件信号闭环",
            "F5 · 经验资产",
        ],
    )


def slide_c7_benchmark(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART C", ACCENT_PURPLE)
    slide_header(s, "C7 · 组件 ⑤ · Benchmark 体系",
                 "「没有度量就没有优化」——Arke 的可复现评测协议")

    # 4-axis matrix
    axes = [
        ("BL", "Benchmark Level",
         "BL1 单算子正确性 → BL5 多算子融合 → BL6 模型 E2E",
         ACCENT),
        ("OT", "Op Tier",
         "OT0 elementwise → OT1 reduction → OT2 dense → OT3 fused → OT4 attention",
         ACCENT_BLUE),
        ("ST", "Shape Tier",
         "ST1 小 → ST2 中 → ST3 大 → ST4 极端（覆盖动态 shape）",
         ACCENT_ALT),
        ("L", "Layer",
         "L1 单算子 · L2 融合 · L3 模型级 E2E",
         ACCENT_PURPLE),
    ]
    cw = Inches(2.95)
    cy = Inches(2.0)
    cx = Inches(0.6)
    for i, (code, t, d, c) in enumerate(axes):
        x = cx + i * (cw + Inches(0.1))
        panel(s, x, cy, cw, Inches(2.05), BG_PANEL, stroke=None)
        accent_bar(s, x, cy, width=cw, height=Inches(0.08), color=c)
        add_text(s, x + Inches(0.2), cy + Inches(0.25),
                 cw - Inches(0.3), Inches(0.5),
                 code, size=22, bold=True, color=c,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)
        add_text(s, x + Inches(0.2), cy + Inches(0.85),
                 cw - Inches(0.3), Inches(0.4),
                 t, size=12, bold=True, color=FG,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), cy + Inches(1.25),
                 cw - Inches(0.3), Inches(0.75),
                 d, size=10.5, color=FG_MUTED,
                 align=PP_ALIGN.CENTER)

    # baselines + reporting
    panel(s, Inches(0.6), Inches(4.3), Inches(7.7), Inches(2.55),
          BG_PANEL, stroke=None)
    accent_bar(s, Inches(0.6), Inches(4.3), width=Inches(0.14),
               height=Inches(2.55), color=ACCENT_ALT)
    add_text(s, Inches(0.95), Inches(4.42), Inches(7), Inches(0.4),
             "多 baseline 矩阵", size=13, bold=True, color=ACCENT_ALT,
             font=FONT_MONO)
    add_bullets(s, Inches(0.95), Inches(4.85), Inches(7.2), Inches(2.0),
                ["P0 cuBLAS / cuDNN —— 厂商极致基线",
                 "P1 FlagGems / Liger —— 开源 Triton 算子库",
                 "P3 PyTorch eager —— 通用基线",
                 "P4 torch.compile (Inductor) —— 图编译器基线",
                 "P5 Arke 自身（KernelCache）—— 自我回归"],
                size=12, bullet_color=ACCENT_ALT, line_spacing=1.25)

    panel(s, Inches(8.45), Inches(4.3), Inches(4.3), Inches(2.55),
          BG_PANEL_ALT, stroke=None)
    accent_bar(s, Inches(8.45), Inches(4.3), width=Inches(0.14),
               height=Inches(2.55), color=ACCENT_RED)
    add_text(s, Inches(8.78), Inches(4.42), Inches(4), Inches(0.4),
             "回应 H9 · 承载 F7", size=13, bold=True, color=ACCENT_RED,
             font=FONT_MONO)
    add_bullets(s, Inches(8.78), Inches(4.85), Inches(3.85), Inches(2.0),
                ["统一 reporting schema (JSON + Markdown)",
                 "跨 baseline 横向可比",
                 "跨硬件 / 跨版本可复现",
                 "Arke 给行业的「公共物品提案」"],
                size=11.5, bullet_color=ACCENT_RED, line_spacing=1.25)
    footer(s, page, total, "Part C · 组件 5")


def slide_c8_alignment(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART C", ACCENT_PURPLE)
    slide_header(s, "C8 · 与目标态 F1–F7 最终对齐",
                 "Arke 不是「再做一个 agent」，而是把目标态变成可工程实现的栈")

    rows = [
        ("F1", "语义 / 策略分离",
         "Language + IR L4/L3", ".ak 双块 + IR 双层强约束", ACCENT),
        ("F2", "有界动作空间",
         "IR L3 + Agent",
         "StrategyIR 决策可枚举；Agent 工具协议禁止自由代码",
         ACCENT_BLUE),
        ("F3", "多层 IR + 多级验证",
         "IR L4–L1 + Compiler V0/V1/V2",
         "四层下沉 × 三级验证", ACCENT_ALT),
        ("F4", "硬件信号闭环",
         "Compiler + Agent",
         "compile_and_profile 反馈进入决策循环", ACCENT_PURPLE),
        ("F5", "结构化经验资产",
         "Language @rationale + Agent 轨迹",
         "一等公民 + 跨会话沉淀", ACCENT_RED),
        ("F6", "跨硬件统一表达",
         "Language + IR + Backend 抽象",
         "单语义多 target；Triton / MLIR / LLVM 渐进", ACCENT),
        ("F7", "可复现评测协议",
         "Benchmark (BL/OT/ST/L)",
         "多 baseline + 统一 schema", ACCENT_ALT),
    ]
    # header
    rx = Inches(0.6)
    rw = Inches(12.15)
    hy = Inches(2.0)
    panel(s, rx, hy, rw, Inches(0.5), BG_PANEL_ALT, stroke=None)
    cols_x = [Inches(0.85), Inches(1.4), Inches(4.6), Inches(8.5)]
    headers = ["F#", "目标态特征", "Arke 承载组件", "一句话说明"]
    for cx_, h in zip(cols_x, headers):
        add_text(s, cx_, hy + Inches(0.1), Inches(3), Inches(0.32),
                 h, size=11.5, bold=True, color=ACCENT_ALT,
                 font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
    for i, (code, t, comp, d, c) in enumerate(rows):
        y = hy + Inches(0.55) + i * Inches(0.62)
        panel(s, rx, y, rw, Inches(0.55),
              BG_PANEL if i % 2 == 0 else BG_PANEL_DEEP, stroke=None)
        accent_bar(s, rx, y, width=Inches(0.1),
                   height=Inches(0.55), color=c)
        add_text(s, cols_x[0], y + Inches(0.13), Inches(0.5), Inches(0.3),
                 code, size=12, bold=True, color=c, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cols_x[1], y + Inches(0.13), Inches(3.2), Inches(0.3),
                 t, size=12, bold=True, color=FG,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cols_x[2], y + Inches(0.13), Inches(3.7), Inches(0.3),
                 comp, size=11, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_MONO)
        add_text(s, cols_x[3], y + Inches(0.13), Inches(3.6), Inches(0.3),
                 d, size=11, color=FG, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total, "Part C")


def slide_c9_qa(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    chapter_chip(s, "PART C", ACCENT_PURPLE)
    slide_header(s, "C9 · 讨论与 Q&A",
                 "5 个研讨引导问题")

    questions = [
        ("01", "有界动作空间的「边界」由谁定义？",
         "编译器枚举 vs 模型外推——谁说了算？", ACCENT),
        ("02", "@rationale 是终极的知识沉淀介质吗？",
         "还是只是过渡形态？下一形态会是什么？", ACCENT_BLUE),
        ("03", "自研领域模型 vs 复用通用前沿模型",
         "分界点在哪里？什么时候必须自研？", ACCENT_ALT),
        ("04", "跨硬件策略迁移：是 IR 层的事，还是模型层的事？",
         "IR-as-portability vs Model-as-portability", ACCENT_PURPLE),
        ("05", "评测的「业界横向可比」由谁主导？",
         "是否应当推动一个开放协议作为公共物品？", ACCENT_RED),
    ]
    rh = Inches(0.92)
    ry = Inches(2.0)
    rx = Inches(0.6)
    rw = Inches(12.15)
    for i, (n, q, sub, c) in enumerate(questions):
        y = ry + i * (rh + Inches(0.04))
        panel(s, rx, y, rw, rh,
              BG_PANEL if i % 2 == 0 else BG_PANEL_DEEP, stroke=None)
        accent_bar(s, rx, y, width=Inches(0.14), height=rh, color=c)
        add_text(s, rx + Inches(0.4), y + Inches(0.22),
                 Inches(0.7), Inches(0.45),
                 n, size=18, bold=True, color=c, font=FONT_MONO)
        add_text(s, rx + Inches(1.3), y + Inches(0.16),
                 Inches(10.5), Inches(0.45),
                 q, size=14, bold=True, color=FG)
        add_text(s, rx + Inches(1.3), y + Inches(0.55),
                 Inches(10.5), Inches(0.4),
                 sub, size=11, color=FG_MUTED, italic=True)
    footer(s, page, total, "Part C")


# ============================================================
# References
# ============================================================
def slide_references(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_text(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35),
             "REFERENCES", size=12, bold=True, color=ACCENT,
             font=FONT_MONO)
    add_text(s, Inches(0.6), Inches(0.78), Inches(12), Inches(0.85),
             "参考资料（论文 · Repo · Blog）", size=26, bold=True, color=FG)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.62), Inches(0.6),
                              Emu(28000))
    _set_fill(line, ACCENT)
    _no_line(line)

    refs = [
        ("KernelEvolve (arXiv)", "https://arxiv.org/abs/2512.23236"),
        ("KernelEvolve · Meta blog",
         "https://engineering.fb.com/2026/04/02/developer-tools/"
         "kernelevolve-..."),
        ("KernelAgent · repo",
         "https://github.com/meta-pytorch/KernelAgent"),
        ("KernelFalcon · PyTorch blog",
         "https://pytorch.org/blog/kernelfalcon-autonomous-gpu-kernel-..."),
        ("KernelAgent · PyTorch blog",
         "https://pytorch.org/blog/kernelagent-hardware-guided-..."),
        ("AutoKernel (arXiv)", "https://arxiv.org/abs/2603.21331"),
        ("AutoKernel · repo",
         "https://github.com/RightNow-AI/autokernel"),
        ("K-Search (arXiv)", "https://arxiv.org/abs/2602.19128"),
        ("K-Search · repo", "https://github.com/caoshiyi/K-Search"),
        ("AVO (arXiv)", "https://arxiv.org/abs/2603.24517"),
        ("CuTeGen (arXiv)", "https://arxiv.org/abs/2604.01489"),
        ("AscendKernelGen (arXiv)", "https://arxiv.org/abs/2601.07160"),
        ("Awesome list",
         "https://github.com/flagos-ai/awesome-LLM-driven-kernel-generation"),
        ("Arke project (this work)",
         "docs/spec/arke-lang-spec-v2.md, docs/spec/arke-ir-spec-v2.md"),
    ]
    top = Inches(2.0)
    rh = Inches(0.36)
    for i, (name, url) in enumerate(refs):
        y = top + i * rh
        panel(s, Inches(0.9), y, Inches(11.5), Inches(0.32),
              BG_PANEL if i % 2 == 0 else BG_PANEL_DEEP, stroke=None)
        add_text(s, Inches(1.1), y + Inches(0.04),
                 Inches(3.6), Inches(0.28),
                 name, size=10.5, bold=True, color=ACCENT_ALT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.7), y + Inches(0.04),
                 Inches(7.6), Inches(0.28),
                 url, size=9, color=FG_MUTED,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
    footer(s, page, total, "References")


# ============================================================
# Build
# ============================================================
def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide plan (v0.5):
    # 01 cover · 02 TOC · 03 Glossary
    # 04 Part A divider · 05 A1 · 06 A2 · 07 A2.5 thesis (NEW) ·
    # 08 A3 · 09 A3 chart · 10 A4 · 11 A5
    # 12 Part B divider · 13 B1 · 14 B2 index · 15–21 B2.1–7 cases ·
    # 22 B3 trends · 23 B4 progress · 24 B5 challenges
    # 25 Target divider · 26 S1 · 27 S2 · 28 S3 · 29 S4
    # 30 Part C divider · 31 C1 · 32 C2 · 33 C3 · 34 C4 · 35 C5 · 36 C6 ·
    # 37 C7 · 38 C8 · 39 C9 · 40 References
    total = 40

    slide_cover(prs)                                                      # 1
    slide_toc(prs, 2, total)                                              # 2
    slide_glossary(prs, 3, total)                                         # 3

    slide_divider(prs, 4, total,
                  "PART A · BACKGROUND",
                  "为什么需要这次洞察",
                  "模型能力 · 范式判断 · 算子工程现状 · 失效 · 输出形态",
                  accent=ACCENT)                                          # 4
    slide_a1_questions(prs, 5, total)                                     # 5
    slide_a2_model_capability(prs, 6, total)                              # 6
    slide_a2p5_thesis(prs, 7, total)                                      # 7  ← NEW
    slide_a3_problem_domain(prs, 8, total)                                # 8
    slide_a3_chart(prs, 9, total)                                         # 9
    slide_a4_failure_modes(prs, 10, total)                                # 10
    slide_a5_output_form(prs, 11, total)                                  # 11

    slide_divider(prs, 12, total,
                  "PART B · INDUSTRY TRENDS",
                  "业界趋势：从 LLM-driven kernel 看技术走向",
                  "7 个代表案例 · 趋势 T1–T6 · 进展 G1–G6 · 难题 H1–H9",
                  accent=ACCENT_BLUE)                                     # 12
    slide_b1_panorama(prs, 13, total)                                     # 13
    slide_b2_index(prs, 14, total)                                        # 14

    # 7 cases (each with 公开 metrics + source citation)
    slide_b2_case(
        prs, 15, total, idx=1, name="KernelEvolve",
        source="Meta · ISCA 2026 · arXiv 2512.23236", color=ACCENT,
        what=[
            "在生产级 ranking infra 上做异构 kernel 生成与优化",
            "搜索式优化：tree search + 自改进 state",
            "RAG 注入硬件知识库（含专有 MTIA 架构）",
            "覆盖多 DSL：Triton / CuTe / FlyDSL / CUDA / HIP / MTIA C++",
        ],
        key_tech=[
            "演化 + 验证 + 反馈三段式 loop",
            "Job harness：规模化并行 compile + profile",
            "运行时 RAG 上下文 → 动态 prompt 合成",
            "生产部署：DLRM / Ads 推理吞吐提升",
        ],
        trend="T1 / T3 / T4",
        gap=[
            "策略仍偏自由代码，跨 kernel 不易抽象成可复用资产",
            "硬件知识库以文档为主，缺结构化决策资产",
        ],
        metrics=("KernelBench 250/250 全 PASS · ATen 160 ops × 3 HW · "
                 "最高 17× vs PyTorch · 已落地 Meta 生产 ranking 系统"),
        source_cite=("Meta engineering blog (2026-04) + arXiv 2512.23236"
                     " (KernelEvolve / Heterogeneous Agentic Kernel "
                     "Authoring at Scale)"),
    )                                                                     # 15
    slide_b2_case(
        prs, 16, total, idx=2,
        name="KernelAgent / KernelFalcon",
        source="PyTorch · Meta · 2025–2026", color=ACCENT_BLUE,
        what=[
            "Deep Agents 分层编排：Orchestrator / FuserAgent / "
            "ExtractorAgent / Workers / Profiler",
            "硬件信号驱动：NCU + roofline 进入决策",
            "严格门禁：编译 → 数值 → 性能 三段验证",
            "并行多 worker 探索 + early-win 收敛",
        ],
        key_tech=[
            "Deterministic orchestration（控制逻辑从 LLM 拿出来）",
            "Grounded tool use：每步都经编译器/硬件验证",
            "子图边界 + shape contract → JSON 合同",
            "用合成 kernel 重建 forward 端到端替换",
        ],
        trend="T2 / T3",
        gap=[
            "输出仍是 Triton 自由代码，缺独立策略层",
            "跨硬件迁移尚未系统化验证（仍以 H100 为主）",
        ],
        metrics=("KernelBench L1/L2/L3 250/250 · 1.56× vs torch.compile · "
                 "89% H100 roofline · 多算子并行合成"),
        source_cite=("PyTorch blog · KernelFalcon (autonomous GPU kernel "
                     "generation via deep agents) + KernelAgent "
                     "(hardware-guided GPU kernel optimization) · "
                     "github.com/meta-pytorch/KernelAgent"),
    )                                                                     # 16
    slide_b2_case(
        prs, 17, total, idx=3, name="AutoKernel",
        source="RightNow AI · arXiv 2603.21331", color=ACCENT_ALT,
        what=[
            "autoresearch 风格优化循环：profile → extract → "
            "edit → bench → keep/revert",
            "Amdahl 优先级排序：集中精力到收益最大的瓶颈",
            "双后端：Triton（快迭代）+ CUDA C++（深硬件）",
            "5 阶段正确性 harness + 6 层 playbook",
        ],
        key_tech=[
            "瓶颈定位 + 优先级调度（torch.profiler）",
            "Tool-use harness：固定基准 + 自动 keep/revert",
            "知识注入以 playbook（提示工程）为主",
            "可长时间无人值守探索",
        ],
        trend="T1 / T2",
        gap=[
            "优化粒度是「编辑文件」，缺中间层级可逆决策",
            "知识沉淀仍是代码片段 + prompt，不易跨 kernel 复用",
        ],
        metrics=("H100 实测：RMSNorm 5.29× · softmax 2.82× · "
                 "cross-entropy 2.21× vs PyTorch eager（公开 README 数据）"),
        source_cite=("arXiv 2603.21331 + github.com/RightNow-AI/autokernel"),
    )                                                                     # 17
    slide_b2_case(
        prs, 18, total, idx=4, name="K-Search",
        source="UC Berkeley · arXiv 2602.19128", color=ACCENT_PURPLE,
        what=[
            "把 kernel 生成视作规划问题：problem → search tree",
            "World Model：LLM 估计状态转移 + 价值，指导搜索",
            "Plan ↔ Codegen 显式解耦：策略与实现分离",
            "Co-Evolving：世界模型随搜索共同演化",
        ],
        key_tech=[
            "World-model + planning + 状态空间剪枝",
            "Stagnation-aware：无提升时自动切 action",
            "Exec feedback 更新世界模型",
            "对临时编译失败更鲁棒（不会「误杀」好策略）",
        ],
        trend="T1 / T5",
        gap=[
            "世界模型大多隐式在 LLM 内部，难跨 session 持久化",
            "硬件信号集成偏弱；输出仍是 Triton/CUDA 自由代码",
        ],
        metrics=("FlashInfer GQA/MLA/MoE 平均 2.10× · MoE 最高 14.3× · "
                 "GPUMode TriMul H100 SOTA 1030μs（公开 paper 数据）"),
        source_cite=("arXiv 2602.19128 + github.com/caoshiyi/K-Search"),
    )                                                                     # 18
    slide_b2_case(
        prs, 19, total, idx=5, name="AVO",
        source="NVIDIA · arXiv 2603.24517", color=ACCENT,
        what=[
            "Evolutionary Algorithm 骨架（population / lineage）",
            "把变异算子 Vary() 升级为自主 agent：propose → "
            "repair → critique → verify",
            "Domain KB 作为知识源（硬件规约 + FA-4 / CUTLASS 实现）",
            "长周期演化（多日尺度）发现微架构级优化",
        ],
        key_tech=[
            "Agent-as-Variation-Operator：从 pipeline → agent",
            "Self-directed loop · 多目标演化",
            "Long-horizon：在 attention 上发现专家未覆盖的优化",
            "Exec feedback 驱动下一轮变异",
        ],
        trend="T1 / T2",
        gap=[
            "缺独立的策略表征层；输出多是 CUDA / C++ 实现",
            "结果可解释性偏弱；长周期探索 token / 算力成本高",
        ],
        metrics=("Blackwell B200 · MHA 7-day 演化 · 超 cuDNN +3.5% · "
                 "超 FlashAttention-4 +10.5% · GQA 迁移仅需 30 min"),
        source_cite=("arXiv 2603.24517 (Agentic Variation Operators "
                     "for Evolutionary Kernel Optimization)"),
    )                                                                     # 19
    slide_b2_case(
        prs, 20, total, idx=6, name="CuTeGen",
        source="U. Toronto · arXiv 2604.01489", color=ACCENT_RED,
        what=[
            "选 CuTe（CUTLASS 4.x layout algebra）作为稳定抽象层",
            "单 kernel 渐进精炼（generate → test → refine）",
            "失败走结构化模板而非纯报错粘贴",
            "先正确性后性能，profile 反馈延迟注入",
        ],
        key_tech=[
            "稳定抽象层 + 渐进精炼",
            "对 tiling / data-movement 更显式",
            "Token 成本可控（不做大规模搜索）",
        ],
        trend="T4",
        gap=[
            "强绑定 NVIDIA / CuTe 生态，跨硬件能力受限",
            "仍是 C++ 模板体系，知识沉淀机制弱",
        ],
        metrics=("CUTLASS v4.3 (CuTe) 上 12 matmul + 14 activation kernels"
                 " · 渐进精炼 token 成本可控（vs 纯 LLM 直写 CUDA）"),
        source_cite=("arXiv 2604.01489 (LLM × CuTe abstraction for "
                     "iterative kernel synthesis)"),
    )                                                                     # 20
    slide_b2_case(
        prs, 21, total, idx=7,
        name="KernelGen-LM / AscendKernelGen",
        source="PCL · 中山大学 · 华为 · arXiv 2601.07160",
        color=ACCENT_ALT,
        what=[
            "Ascend-CoT 数据集：真实 kernel + 文档推理 + 代码推理链",
            "KernelGen-LM：Qwen3 backbone · NPU-aware SFT",
            "RLEF：执行反馈强化（compile / 正确性 / 性能信号）",
            "NPUKernelBench：编译/正确性/性能三维评测，难度分层",
        ],
        key_tech=[
            "走「训一个领域模型」路线：知识进入权重",
            "强耦合单硬件（AscendC / 昇腾）换取深度",
            "用 CoT + RLEF 填补 DSL 数据稀缺",
            "评测一等公民：编译 / 正确性 / 性能全量指标",
        ],
        trend="T6",
        gap=[
            "通用性受限：跨硬件需要再训 / 再适配",
            "工具链与决策协议尚未通用化",
        ],
        metrics=("Ascend NPU · AscendC DSL · NPUKernelBench L2 编译"
                 "成功率 0% → 95.5% (Pass@10)；权重 / 数据 / 评测全部开源"),
        source_cite=("arXiv 2601.07160 (KernelGen-LM / AscendKernelGen)"),
    )                                                                     # 21

    slide_b3_trends(prs, 22, total)                                       # 22
    slide_b4_progress(prs, 23, total)                                     # 23
    slide_b5_challenges(prs, 24, total)                                   # 24

    slide_divider(prs, 25, total,
                  "TARGET STATE",
                  "目标态总览：未来算子生成 / 调优技术",
                  "北极星 · 7 大技术特征 F1–F7 · 对通用编程的启示",
                  accent=ACCENT_ALT)                                      # 25
    slide_s1_north_star(prs, 26, total)                                   # 26
    slide_s2_features(prs, 27, total)                                     # 27
    slide_s3_implications(prs, 28, total)                                 # 28
    slide_s4_bridge(prs, 29, total)                                       # 29

    slide_divider(prs, 30, total,
                  "PART C · ARKE",
                  "Arke 的技术构建：设计逻辑 + 整体架构 + 四件套 + Benchmark",
                  "5 条公理 · 一张总图 · Language / IR / Compiler / Agent / Benchmark",
                  accent=ACCENT_PURPLE)                                   # 30
    slide_c1_axioms(prs, 31, total)                                       # 31
    slide_c2_architecture(prs, 32, total)                                 # 32
    slide_c3_language(prs, 33, total)                                     # 33
    slide_c4_ir(prs, 34, total)                                           # 34
    slide_c5_compiler(prs, 35, total)                                     # 35
    slide_c6_agent(prs, 36, total)                                        # 36
    slide_c7_benchmark(prs, 37, total)                                    # 37
    slide_c8_alignment(prs, 38, total)                                    # 38
    slide_c9_qa(prs, 39, total)                                           # 39

    slide_references(prs, 40, total)                                      # 40

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote: {out_path}")


if __name__ == "__main__":
    out = (Path(__file__).resolve().parents[2]
           / "docs" / "sharing" / "ai-native-insight-deck.pptx")
    build(out)
