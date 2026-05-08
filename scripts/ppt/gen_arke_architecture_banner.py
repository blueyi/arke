"""
Arke architecture banner PPT — single slide, 570×120 px.
LLM-Native edition: simplified to 4 boxes + feedback loop.

Layout (all left → right pipeline + feedback arrow above):
  feedback ←─ legal actions / profile / gate verdict ───────────
                                                                 │
  [① LLM Agent] ─→ [② Arke Lang + IR] ─→ [③ Compiler-Verifier] ──┴→ [④ Multi-HW]

Output: docs/sharing/ai-native-arke-architecture-banner.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ============================================================
# Slide size = 570 × 120 px @ 96 DPI; 1 px = 9525 EMU
# ============================================================
PX = 9525
SLIDE_W = Emu(570 * PX)   # 5.9375 in
SLIDE_H = Emu(120 * PX)   # 1.25  in


# ============================================================
# Theme — matches the main 41-slide insight deck
# ============================================================
BG_DARK = RGBColor(0x0B, 0x1B, 0x2E)
BG_PANEL = RGBColor(0x12, 0x2A, 0x43)
BG_PANEL_ALT = RGBColor(0x18, 0x35, 0x53)
FG = RGBColor(0xEA, 0xF2, 0xFA)
FG_MUTED = RGBColor(0x9F, 0xB5, 0xC9)
FG_DIM = RGBColor(0x6E, 0x86, 0x9F)

ACCENT = RGBColor(0x38, 0xD1, 0xB8)         # teal — Arke Lang+IR
ACCENT_BLUE = RGBColor(0x7A, 0xB8, 0xFF)    # LLM
ACCENT_ALT = RGBColor(0xF2, 0xC5, 0x5C)     # gold — Compiler-Verifier
ACCENT_RED = RGBColor(0xEF, 0x6E, 0x6E)     # HW
STROKE = RGBColor(0x24, 0x44, 0x66)

FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Consolas"


# ============================================================
# Helpers
# ============================================================
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _set_line(shape, color, width_pt=0.6):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def emu(px_value):
    return Emu(int(px_value * PX))


def add_text(slide, left, top, width, height, text, *,
             size=8, bold=False, italic=False,
             color=FG, font=FONT_BODY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


# ============================================================
# Build
# ============================================================
def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s = prs.slides.add_slide(prs.slide_layouts[6])

    # background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, BG_DARK)
    _no_line(bg)
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)

    # top + bottom hair-line accent bands (1 px each)
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, SLIDE_W, emu(1.5))
    _set_fill(top_bar, ACCENT)
    _no_line(top_bar)
    bot_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, SLIDE_H - emu(1.5),
                                 SLIDE_W, emu(1.5))
    _set_fill(bot_bar, ACCENT)
    _no_line(bot_bar)

    # ============================================================
    # Top band: title + LLM-NATIVE chip (the visual anchor)
    # ============================================================
    # Title
    add_text(s, emu(8), emu(4), emu(380), emu(13),
             "ARKE · AI-NATIVE OPERATOR COMPILER STACK",
             size=7, bold=True, color=ACCENT_BLUE, font=FONT_MONO,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Subtitle (right side)
    add_text(s, emu(360), emu(4), emu(202), emu(13),
             "LLM 做决策 · 编译器把关数学",
             size=6.5, color=FG_MUTED, font=FONT_MONO,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Feedback loop row — visualized as a single right-to-left arrow
    # spanning the top of the box row, labeled clearly
    # ============================================================
    fb_y = emu(20)
    fb_h = emu(8)
    # backing line
    backing = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 emu(120), fb_y + emu(3),
                                 emu(330), emu(2))
    _set_fill(backing, ACCENT_ALT)
    _no_line(backing)
    # left-pointing arrow head (feedback flows right→left)
    arr_head = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                                  emu(112), fb_y, emu(10), emu(8))
    _set_fill(arr_head, ACCENT_ALT)
    _no_line(arr_head)
    # label
    add_text(s, emu(140), fb_y - emu(2), emu(290), emu(11),
             "feedback loop  ·  legal actions  ·  hw signal (profile / roofline)  ·  gate verdict",
             size=5.5, italic=True, color=ACCENT_ALT, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Main pipeline — exactly 4 boxes
    # ============================================================
    PAD_X = 6
    GAP = 10  # arrow gap between boxes
    ARROW_W = 10
    BOX_TOP_PX = 32
    BOX_H_PX = 60

    # Width budget: 570 - 2*PAD_X - 3*GAP = 528 → 132 each
    W = (570 - 2 * PAD_X - 3 * GAP) // 4   # 132

    boxes = [
        # (id, title, line1, line2, color, fill_bg, glyph_color)
        ("①", "LLM as Decider",
         "决策 / 推理 / 反思",
         "NOT a code generator",
         ACCENT_BLUE, BG_PANEL_ALT),

        ("②", "Arke Lang + IR",
         "语义 / 策略分离  ·  rationale",
         "结构化 · 有界动作空间",
         ACCENT, BG_PANEL),

        ("③", "Compiler-Verifier",
         "静态 / 数值 / 性能 三级门禁",
         "checkpoint  ·  rollback",
         ACCENT_ALT, BG_PANEL),

        ("④", "Multi-Hardware",
         "NVIDIA · Ascend · AMD",
         "→ 多 NPU / DSA",
         ACCENT_RED, BG_PANEL_ALT),
    ]

    box_y = emu(BOX_TOP_PX)
    box_h = emu(BOX_H_PX)
    cur_x = PAD_X
    box_centers = []
    for i, (g, head, l1, l2, col, fill_bg) in enumerate(boxes):
        x = emu(cur_x)
        bw = emu(W)
        # rounded box
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, box_y, bw, box_h)
        box.adjustments[0] = 0.16
        _set_fill(box, fill_bg)
        _set_line(box, col, 0.9)
        # left accent stripe
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x, box_y, emu(3), box_h)
        _set_fill(bar, col)
        _no_line(bar)

        # glyph (top-left)
        add_text(s, x + emu(6), box_y + emu(3), emu(16), emu(14),
                 g, size=10, bold=True, color=col, font=FONT_MONO,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

        # title (bold)
        add_text(s, x + emu(6), box_y + emu(4),
                 bw - emu(12), emu(15),
                 head, size=8, bold=True, color=col,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)

        # body line 1 (key concept)
        add_text(s, x + emu(8), box_y + emu(22),
                 bw - emu(14), emu(14),
                 l1, size=6, bold=True, color=FG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # body line 2 (qualifier)
        add_text(s, x + emu(8), box_y + emu(36),
                 bw - emu(14), emu(20),
                 l2, size=5.5, color=FG_MUTED, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        box_centers.append((cur_x, cur_x + W))
        cur_x += W
        # arrow
        if i < len(boxes) - 1:
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                emu(cur_x + (GAP - ARROW_W) // 2),
                box_y + emu(BOX_H_PX // 2 - 5),
                emu(ARROW_W), emu(10))
            _set_fill(arr, FG_MUTED)
            _no_line(arr)
            cur_x += GAP

    # ============================================================
    # Bottom caption — LLM-Native scenarios
    # ============================================================
    cap_y = emu(BOX_TOP_PX + BOX_H_PX + 4)
    add_text(s, emu(8), cap_y, emu(560), emu(11),
             "解决场景  ·  LLM 驱动 kernel 生成 / 优化  ·  跨硬件统一表达 (SIMT ↔ SIMD ↔ NPU)  ·  最小 token 预算下逼近厂商库性能",
             size=6.2, bold=True, color=FG,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, emu(8), cap_y + emu(11), emu(560), emu(10),
             "LLM-Native  =  bounded action  +  rationale  +  multi-level gates  +  hw-signal closed loop",
             size=5.5, italic=True, color=ACCENT, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote: {out_path}")
    print(f"slide size: {prs.slide_width / PX:.0f} × "
          f"{prs.slide_height / PX:.0f} px "
          f"({prs.slide_width / 914400:.3f} × "
          f"{prs.slide_height / 914400:.3f} in)")
    print(f"box count: {len(boxes)}")


if __name__ == "__main__":
    out = (Path(__file__).resolve().parents[2]
           / "docs" / "sharing"
           / "ai-native-arke-architecture-banner.pptx")
    build(out)
