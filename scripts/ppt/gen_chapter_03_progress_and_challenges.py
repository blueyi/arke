"""
Chapter 03 deck generator: Progress & challenges distilled from 7 industry case studies.

Output: docs/sharing/chapter-03-progress-and-challenges.pptx
"""

from __future__ import annotations

from pathlib import Path

import importlib.util

from pptx import Presentation


def _load_master_module():
    master_path = Path(__file__).resolve().parents[1] / "gen_sharing_ppt.py"
    spec = importlib.util.spec_from_file_location("gen_sharing_ppt", master_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Failed to load module from {master_path}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)  # type: ignore[attr-defined]
    return mod


def build(out_path: Path) -> None:
    deck = _load_master_module()

    prs = Presentation()
    prs.slide_width = deck.SLIDE_W
    prs.slide_height = deck.SLIDE_H

    total = 9

    # 1 cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    deck.set_bg(s)
    deck.slide_header(s, "Chapter 03", "从 7 个方案抽象：好进展 & 关键难题清单")
    deck.panel(s, deck.Inches(0.9), deck.Inches(2.05), deck.Inches(11.5), deck.Inches(4.9), deck.BG_PANEL)
    deck.add_text(
        s,
        deck.Inches(1.2),
        deck.Inches(2.25),
        deck.Inches(11.0),
        deck.Inches(0.7),
        "目标：把“案例证据”抽象成可执行的技术问题列表，为下一章 Arke 对齐与 TODO 做输入。",
        size=16,
        color=deck.FG,
    )
    deck.add_bullets(
        s,
        deck.Inches(1.2),
        deck.Inches(3.05),
        deck.Inches(11.0),
        deck.Inches(3.6),
        [
            "输入：KernelEvolve / KernelAgent / AutoKernel / K-Search / AVO / CuTeGen / KernelGen-LM",
            "输出 1：已被证实有效的工程模式（可复用）",
            "输出 2：仍未体系化解决的难题（H1–H10）",
            "输出 3：面向目标态的能力清单（能力缺口 = TODO）",
        ],
        size=14,
        bullet_color=deck.ACCENT,
    )
    deck.footer(s, 1)

    # 2 progress (what works)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    deck.set_bg(s)
    deck.slide_header(s, "Chapter 03", "好进展（What’s working）：业界已验证有效的 8 个模式")
    deck.panel(s, deck.Inches(0.6), deck.Inches(1.9), deck.Inches(12.1), deck.Inches(5.2), deck.BG_PANEL)
    deck.accent_bar(s, deck.Inches(0.6), deck.Inches(1.9), width=deck.Inches(0.12), height=deck.Inches(5.2), color=deck.ACCENT)
    deck.add_bullets(
        s,
        deck.Inches(0.95),
        deck.Inches(2.05),
        deck.Inches(11.6),
        deck.Inches(4.9),
        [
            "P1  从 one-shot 走向 search / evolve / autoresearch（规模化探索是必须）",
            "P2  验证门禁成为共识：compile → correctness → performance 的 gated loop",
            "P3  硬件信号进入 loop：NCU / roofline / metrics 驱动优化方向",
            "P4  并行探索 + early-win：multi-worker 显著提高成功率与收敛速度",
            "P5  选择稳定抽象层（Triton / CuTe / DSL）降低迭代难度与幻觉概率",
            "P6  显式知识注入：RAG 硬件手册 / playbook / prompt template 成标配",
            "P7  模型/算子级“目标函数”被工程化：Amdahl 排序、端到端收益可解释",
            "P8  复现与可观测性开始成体系：artifact、trajectory、benchmark 固化流程",
        ],
        size=13,
        bullet_color=deck.ACCENT,
        line_spacing=1.18,
    )
    deck.footer(s, 2)

    # 3–5 challenges list (H1–H10)
    challenges = [
        ("H1", "策略与实现仍耦合在自由代码里", "难审计、难回滚、难迁移；好的策略会被临时编译失败“误杀”。"),
        ("H2", "知识沉淀介质不统一", "RAG / 权重 / 模板 / 日志各自为政；跨硬件迁移缺客观度量。"),
        ("H3", "动态 shape 的泛化与性能并存困难", "shape guard / specialization / cache 策略缺统一抽象。"),
        ("H4", "跨硬件（SIMT↔SIMD）复用成本高", "硬件约束差异导致策略重写；需要 target-aware 决策层。"),
        ("H5", "Triton 抽象天花板", "更深层（寄存器/屏障/指令调度）优化难表达。"),
        ("H6", "搜索成本与 token 成本失控", "长周期演化（如 7 天）昂贵；缺预算治理与状态压缩规范。"),
        ("H7", "从单 kernel 到模型级自治优化", "bottleneck 定位→回灌优化→回归验证 尚未产品化。"),
        ("H8", "评测维度不全导致“假提升”", "只看单 kernel latency 会牺牲端到端；需 BL6 类模型级指标。"),
        ("H9", "可组合性与融合边界难题", "子图合同/融合边界抽取不稳定；需要语义层合约。"),
        ("H10", "安全与可控性", "自由代码生成风险（UB/越界/资源滥用）；需要可验证的动作空间。"),
    ]

    def _challenges_slide(page: int, title: str, items):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        deck.set_bg(s)
        deck.slide_header(s, "Chapter 03", title)
        top = deck.Inches(1.95)
        row_h = deck.Inches(0.9)
        for i, (hid, htitle, hdesc) in enumerate(items):
            y = top + i * row_h
            deck.panel(s, deck.Inches(0.9), y, deck.Inches(11.5), deck.Inches(0.8), deck.BG_PANEL)
            badge = s.shapes.add_shape(deck.MSO_SHAPE.OVAL, deck.Inches(1.05), y + deck.Inches(0.15), deck.Inches(0.5), deck.Inches(0.5))
            deck._set_fill(badge, deck.ACCENT_RED)
            deck._no_line(badge)
            deck.add_text(s, deck.Inches(1.05), y + deck.Inches(0.15), deck.Inches(0.5), deck.Inches(0.5),
                          hid, size=11, bold=True, color=deck.BG_DARK, align=deck.PP_ALIGN.CENTER, anchor=deck.MSO_ANCHOR.MIDDLE,
                          font=deck.FONT_MONO)
            deck.add_text(s, deck.Inches(1.7), y + deck.Inches(0.12), deck.Inches(4.4), deck.Inches(0.32),
                          htitle, size=12, bold=True, color=deck.FG, anchor=deck.MSO_ANCHOR.MIDDLE)
            deck.add_text(s, deck.Inches(6.1), y + deck.Inches(0.12), deck.Inches(6.1), deck.Inches(0.32),
                          hdesc, size=10, color=deck.FG_MUTED, anchor=deck.MSO_ANCHOR.MIDDLE)
        deck.footer(s, page)

    _challenges_slide(3, "关键难题（H1–H5）：结构化表达 × 跨硬件 × 深层控制", challenges[:5])
    _challenges_slide(4, "关键难题（H6–H10）：成本治理 × 规模化落地 × 可控性", challenges[5:])

    # 5 ability map: required capabilities checklist (bridge to Arke)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    deck.set_bg(s)
    deck.slide_header(s, "Chapter 03", "从难题反推能力：目标态系统需要具备什么？")
    deck.panel(s, deck.Inches(0.6), deck.Inches(1.9), deck.Inches(12.1), deck.Inches(5.2), deck.BG_PANEL)
    deck.add_text(s, deck.Inches(0.9), deck.Inches(2.0), deck.Inches(11.5), deck.Inches(0.4),
                  "能力清单（对应 H1–H10）", size=14, bold=True, color=deck.ACCENT)
    deck.add_bullets(
        s,
        deck.Inches(0.9),
        deck.Inches(2.45),
        deck.Inches(11.5),
        deck.Inches(4.5),
        [
            "C1  语义合约层：可组合、可融合、可对拍（解决 H1/H9）",
            "C2  策略资产层：可审计、可回滚、可迁移（解决 H1/H2/H10）",
            "C3  多级验证器：V0 快剪枝 + V1 数值 + V2 性能（解决 H2/H8/H10）",
            "C4  符号维度与动态 shape：推断/约束/特化/缓存统一（解决 H3）",
            "C5  多后端路径：Triton 快迭代 + MLIR/LLVM 深控制（解决 H5）",
            "C6  成本治理：budget / compact / caching（解决 H6）",
            "C7  模型级自治循环：profile bottleneck → optimize → regress（解决 H7/H8）",
            "C8  跨硬件知识迁移机制：目标相关策略 + 可度量迁移收益（解决 H4/H2）",
        ],
        size=13,
        bullet_color=deck.ACCENT_ALT,
        line_spacing=1.18,
    )
    deck.footer(s, 5)

    # 6 bridge
    s = prs.slides.add_slide(prs.slide_layouts[6])
    deck.set_bg(s)
    deck.slide_header(s, "Chapter 03", "结论：下一章把 C1–C8 映射到 Arke 四件套，并列出 TODO")
    deck.panel(s, deck.Inches(0.9), deck.Inches(2.2), deck.Inches(11.5), deck.Inches(3.9), deck.BG_PANEL_ALT)
    deck.add_bullets(
        s,
        deck.Inches(1.2),
        deck.Inches(2.5),
        deck.Inches(11.0),
        deck.Inches(3.2),
        [
            "Arke 作为“AI 原生算子编译栈”，应覆盖 C1–C8 的大部分能力。",
            "凡是当前设计/实现未覆盖或缺验证指标的点，将以 TODO 标注（并给出建议验收方式）。",
            "按 Phase/Gate（BL/OT/ST/L）把能力落到可度量路径上，避免停留在趋势叙事。",
        ],
        size=16,
        bullet_color=deck.ACCENT,
    )
    deck.footer(s, 6)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote: {out_path}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "docs" / "sharing" / "chapter-03-progress-and-challenges.pptx"
    build(out)

