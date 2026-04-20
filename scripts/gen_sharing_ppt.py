"""Generate technical sharing deck for the AI-Native Compile Stack.

Output: docs/sharing/ai-native-compile-stack.pptx
Theme:  dark navy + teal accent, 16:9
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ------------------------------------------------------------------ palette

BG_DARK = RGBColor(0x0B, 0x1B, 0x2E)
BG_PANEL = RGBColor(0x12, 0x2A, 0x43)
BG_PANEL_ALT = RGBColor(0x18, 0x35, 0x53)
FG = RGBColor(0xEA, 0xF2, 0xFA)
FG_MUTED = RGBColor(0x9F, 0xB5, 0xC9)
ACCENT = RGBColor(0x38, 0xD1, 0xB8)
ACCENT_ALT = RGBColor(0xF2, 0xC5, 0x5C)
ACCENT_BLUE = RGBColor(0x7A, 0xB8, 0xFF)
ACCENT_RED = RGBColor(0xEF, 0x6E, 0x6E)
STROKE = RGBColor(0x24, 0x44, 0x66)

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TOTAL = 32

# ------------------------------------------------------------------ helpers


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _set_line(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_bg(slide, color: RGBColor = BG_DARK) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, color)
    _no_line(bg)
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = FG,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    *,
    size: int = 14,
    color: RGBColor = FG,
    bullet_color: RGBColor = ACCENT,
    line_spacing: float = 1.25,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        marker = "▸ " if level == 0 else "·  "
        r1 = p.add_run()
        r1.text = marker
        r1.font.name = FONT_BODY
        r1.font.size = Pt(size)
        r1.font.bold = level == 0
        r1.font.color.rgb = bullet_color if level == 0 else FG_MUTED
        r2 = p.add_run()
        r2.text = text
        r2.font.name = FONT_BODY
        r2.font.size = Pt(size - (0 if level == 0 else 1))
        r2.font.color.rgb = color if level == 0 else FG_MUTED
        p.level = level
    return tb


def panel(slide, left, top, width, height, color=BG_PANEL, stroke=STROKE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.06
    _set_fill(shp, color)
    _set_line(shp, stroke, 0.75)
    return shp


def accent_bar(slide, left, top, width=Inches(0.12), height=Inches(0.42), color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill(shp, color)
    _no_line(shp)
    return shp


def slide_header(slide, kicker: str, title: str):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35), kicker,
             size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8), title,
             size=28, bold=True, color=FG)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55),
                                  Inches(0.6), Emu(26000))
    _set_fill(line, ACCENT)
    _no_line(line)


def footer(slide, page_num: int):
    add_text(slide, Inches(0.6), Inches(7.15), Inches(6), Inches(0.28),
             "AI-Native Compile Stack · 四件套 · 技术分享",
             size=9, color=FG_MUTED)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.28),
             f"{page_num:02d} / {TOTAL:02d}",
             size=9, color=FG_MUTED, align=PP_ALIGN.RIGHT)


# ------------------------------------------------------------------ slides


def make_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    for y in (Inches(0.0), Inches(7.42)):
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, Inches(0.08))
        _set_fill(b, ACCENT)
        _no_line(b)

    add_text(slide, Inches(0.9), Inches(1.2), Inches(11), Inches(0.5),
             "AI-NATIVE COMPILE STACK · TECH SHARING",
             size=14, bold=True, color=ACCENT)
    add_text(slide, Inches(0.9), Inches(1.8), Inches(12), Inches(1.2),
             "AI 原生算子编译栈的「四件套」",
             size=44, bold=True, color=FG)
    add_text(slide, Inches(0.9), Inches(2.9), Inches(12), Inches(0.9),
             "Language × IR × Compiler Toolchain × Agent",
             size=30, bold=True, color=ACCENT)
    add_text(slide, Inches(0.9), Inches(3.8), Inches(12), Inches(0.6),
             "构建策略、必要性与工程权衡",
             size=22, color=FG_MUTED)

    panel(slide, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.4), BG_PANEL)
    add_text(slide, Inches(1.1), Inches(5.35), Inches(11), Inches(0.4),
             "从「让 LLM 写代码」到「让 LLM 做决策」—— 一套让 LLM 在有界动作空间中探索、编译器负责验证的 AI-Native 基础设施",
             size=14, color=FG)
    add_text(slide, Inches(1.1), Inches(5.8), Inches(11), Inches(0.4),
             "讨论焦点  ·  Language / IR / Compiler Toolchain / Agent 四件套的构建必要性、设计取舍与未来展开",
             size=12, color=FG_MUTED)
    add_text(slide, Inches(1.1), Inches(6.2), Inches(11), Inches(0.4),
             "案例参照  ·  Arke Project · 45 ops · OT0–OT4 · BL1–BL6 · Triton → MLIR → LLVM IR",
             size=12, color=FG_MUTED)


def make_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "AGENDA", "本次分享地图")

    items = [
        ("01", "背景：AI 编程编译的范式迁移", "为什么「四件套」不是可选项"),
        ("02–06", "四件套分述：Language / IR / Compiler / Agent",
         "必要性 · 构建策略 · 技术讨论点"),
        ("07", "业界技术洞察与方案全景",
         "LLM-for-Code · Autotune · MLIR · Agentic Infra 四条主线"),
        ("08", "AI 原生编译栈的 6 个必备特征 (T1–T6)",
         "由业界洞察抽取，作为方案评估基线"),
        ("09", "方案对比  ·  本方案 vs 业界",
         "T1–T6 六维度对比表 + 优势简述"),
        ("10", "整体技术架构  +  四件套各自架构图",
         "Arke 总图 · Language · IR · Compiler · Agent"),
        ("15", "贯通：职责矩阵 · 验证 · Benchmark · Token · 路线图 · Q&A",
         "把架构锁到可度量的坐标上"),
    ]
    top = Inches(1.85)
    row_h = Inches(0.72)
    for i, (no, title, sub) in enumerate(items):
        y = top + i * row_h
        panel(slide, Inches(0.9), y, Inches(11.5), Inches(0.62), BG_PANEL)
        add_text(slide, Inches(1.05), y + Inches(0.12), Inches(1.3), Inches(0.42),
                 no, size=15, bold=True, color=ACCENT, font=FONT_MONO)
        add_text(slide, Inches(2.35), y + Inches(0.07), Inches(7.0), Inches(0.3),
                 title, size=14, bold=True, color=FG)
        add_text(slide, Inches(2.35), y + Inches(0.34), Inches(9.1), Inches(0.28),
                 sub, size=11, color=FG_MUTED)
    footer(slide, 2)


def make_bg_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "01 · 背景", "AI 时代算子开发的结构性拐点")

    cols = [
        ("硬件膨胀", ACCENT, [
            "NVIDIA SIMT / Ascend · AMD · NPU / DSA",
            "架构维度持续扩张",
            "「一套代码一套调优」 成本不可承受",
        ]),
        ("算子复杂化", ACCENT_ALT, [
            "FlashAttention · GQA · MLA · RoPE / YaRN",
            "Fused Compound · 动态 shape",
            "传统 Triton / CUDA 人写人调成为瓶颈",
        ]),
        ("LLM 直写失灵", ACCENT_RED, [
            "正确率 83% vs 结构化 100%",
            "~3500 tokens / 迭代 · 幻觉 · 调试链路长",
            "性能方差大 · 决策不可解释",
        ]),
    ]
    col_w = Inches(4.0)
    col_h = Inches(3.6)
    gap = Inches(0.1)
    total_w = col_w * 3 + gap * 2
    start = (SLIDE_W - total_w) / 2

    for i, (title, col, lines) in enumerate(cols):
        x = start + i * (col_w + gap)
        panel(slide, x, Inches(1.95), col_w, col_h, BG_PANEL)
        accent_bar(slide, x, Inches(1.95), width=col_w, height=Inches(0.08), color=col)
        add_text(slide, x + Inches(0.3), Inches(2.2), col_w - Inches(0.6), Inches(0.5),
                 title, size=20, bold=True, color=col)
        add_bullets(slide, x + Inches(0.3), Inches(2.85), col_w - Inches(0.6), Inches(2.6),
                    lines, size=13, bullet_color=col)

    panel(slide, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.1), BG_PANEL_ALT)
    add_text(slide, Inches(1.1), Inches(5.95), Inches(11), Inches(0.4),
             "方向性共识", size=13, bold=True, color=ACCENT)
    add_text(slide, Inches(1.1), Inches(6.28), Inches(11.1), Inches(0.7),
             "AI 时代的编译基础设施不是「一个更强的编译器」，而是一套让 LLM 以「决策者」身份、让编译器以「验证者」身份协同工作的完整栈",
             size=14, color=FG)
    footer(slide, 3)


def make_bg_fourset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "01 · 背景", "四件套：把「AI 原生编译栈」拆成可独立演进的四块")

    items = [
        ("Language", "AI 友好 DSL", ACCENT, [
            "语义 / 策略分离",
            "算子级抽象 · 符号 shape",
            "Token 最小化入口",
        ]),
        ("IR", "多层 LLM-Native IR", ACCENT_ALT, [
            "L4 Semantic · L3 Strategy",
            "L2 Schedule · L1 Instruction",
            "按 LLM 参与度分层",
        ]),
        ("Compiler Toolchain", "OpRegistry + Pass + Verifier", ACCENT_BLUE, [
            "算子 / Pass / 后端 可插拔",
            "数据驱动替代 if/elif 地狱",
            "V0 / V1 / V2 三级验证",
        ]),
        ("Agent", "有界动作空间 + 协议", ACCENT_RED, [
            "Tool-use 结构化 API",
            "Budget · Trajectory · @rationale",
            "Mode A 内置 / Mode B 外部",
        ]),
    ]
    card_w = Inches(2.95)
    card_h = Inches(4.2)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start = (SLIDE_W - total_w) / 2
    for i, (tag, subtitle, col, bullets) in enumerate(items):
        x = start + i * (card_w + gap)
        panel(slide, x, Inches(1.95), card_w, card_h, BG_PANEL)
        accent_bar(slide, x, Inches(1.95), width=card_w, height=Inches(0.08), color=col)
        add_text(slide, x + Inches(0.25), Inches(2.15), Inches(0.6), Inches(0.4),
                 f"0{i + 1}", size=22, bold=True, color=col, font=FONT_MONO)
        add_text(slide, x + Inches(0.25), Inches(2.6), card_w - Inches(0.5), Inches(0.45),
                 tag, size=18, bold=True, color=FG)
        add_text(slide, x + Inches(0.25), Inches(3.05), card_w - Inches(0.5), Inches(0.35),
                 subtitle, size=11, color=FG_MUTED)
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.25),
                                     Inches(3.45), Inches(0.8), Emu(13000))
        _set_fill(sep, col)
        _no_line(sep)
        add_bullets(slide, x + Inches(0.25), Inches(3.6), card_w - Inches(0.5), Inches(2.5),
                    bullets, size=12, bullet_color=col)

    panel(slide, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.7), BG_PANEL_ALT)
    add_text(slide, Inches(1.1), Inches(6.42), Inches(11), Inches(0.5),
             "耦合而解耦  ·  Language ↔ IR 贴近  /  Compiler ↔ IR 贴近  /  Agent 只通过工具协议访问前三者  —— 任何一件变动都可以在界面处被隔离",
             size=13, color=FG, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, 4)


def make_piece_intro(prs, piece_no, piece_name, english, accent, tagline, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.9), Inches(1.1), Inches(5), Inches(1.0),
             f"件套 {piece_no}", size=18, bold=True, color=accent)
    add_text(slide, Inches(0.9), Inches(1.45), Inches(12), Inches(1.6),
             piece_name, size=52, bold=True, color=FG)
    add_text(slide, Inches(0.9), Inches(3.0), Inches(12), Inches(0.8),
             english, size=22, color=accent, bold=True)

    bignum = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(1.4),
                                    Inches(2.3), Inches(2.3))
    _set_fill(bignum, BG_PANEL)
    _set_line(bignum, accent, 2.0)
    add_text(slide, Inches(10.2), Inches(1.6), Inches(2.3), Inches(1.9),
             f"0{piece_no}", size=72, bold=True, color=accent,
             font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    panel(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(1.7), BG_PANEL)
    accent_bar(slide, Inches(0.9), Inches(4.3), width=Inches(0.14), height=Inches(1.7),
               color=accent)
    add_text(slide, Inches(1.2), Inches(4.5), Inches(11), Inches(0.4),
             "议题核心", size=12, bold=True, color=accent)
    add_text(slide, Inches(1.2), Inches(4.85), Inches(11), Inches(1.0),
             tagline, size=16, color=FG)

    tags = ["必要性 · 为什么必须存在", "构建策略 · 如何做对", "技术讨论点 · 开放问题"]
    for i, t in enumerate(tags):
        x = Inches(0.9) + i * Inches(3.9)
        tbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, Inches(6.2), Inches(3.7), Inches(0.6))
        tbox.adjustments[0] = 0.3
        _set_fill(tbox, BG_PANEL_ALT)
        _set_line(tbox, accent, 0.75)
        add_text(slide, x, Inches(6.25), Inches(3.7), Inches(0.5),
                 t, size=12, bold=True, color=FG,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def make_two_col_content(prs, kicker, title,
                         left_title, left_bullets,
                         right_title, right_bullets,
                         page, accent=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, kicker, title)

    panel(slide, Inches(0.6), Inches(1.85), Inches(6.1), Inches(5.1), BG_PANEL)
    accent_bar(slide, Inches(0.6), Inches(1.85),
               width=Inches(0.12), height=Inches(5.1), color=accent)
    add_text(slide, Inches(0.9), Inches(2.0), Inches(5.8), Inches(0.45),
             left_title, size=16, bold=True, color=accent)
    add_bullets(slide, Inches(0.9), Inches(2.55), Inches(5.8), Inches(4.3),
                left_bullets, size=13, bullet_color=accent)

    panel(slide, Inches(6.9), Inches(1.85), Inches(6.0), Inches(5.1), BG_PANEL_ALT)
    accent_bar(slide, Inches(6.9), Inches(1.85),
               width=Inches(0.12), height=Inches(5.1), color=ACCENT_ALT)
    add_text(slide, Inches(7.2), Inches(2.0), Inches(5.6), Inches(0.45),
             right_title, size=16, bold=True, color=ACCENT_ALT)
    add_bullets(slide, Inches(7.2), Inches(2.55), Inches(5.6), Inches(4.3),
                right_bullets, size=13, bullet_color=ACCENT_ALT)
    footer(slide, page)


def make_discussion(prs, kicker, title, questions, page, accent=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, kicker, title)

    top = Inches(1.95)
    row_h = Inches(0.88)
    for i, q in enumerate(questions):
        y = top + i * row_h
        panel(slide, Inches(0.9), y, Inches(11.5), Inches(0.78), BG_PANEL)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.05),
                                       y + Inches(0.14), Inches(0.5), Inches(0.5))
        _set_fill(badge, accent)
        _no_line(badge)
        add_text(slide, Inches(1.05), y + Inches(0.14), Inches(0.5), Inches(0.5),
                 f"Q{i + 1}", size=11, bold=True, color=BG_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.75), y + Inches(0.15), Inches(10.5), Inches(0.6),
                 q, size=13, color=FG, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def make_integration(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "15 · 贯通", "四件套如何协同 —— 一张职责矩阵")

    rows = [
        ("Language", ".ak v2 · where · tuple · @rationale",
         "句法层 token 最小化与语义显式", ACCENT),
        ("IR", "L4 / L3 / L2 / L1 · 符号维度贯穿",
         "按 LLM 参与度解耦责任", ACCENT_ALT),
        ("Compiler", "OpRegistry · Pass · Backend · Verifier",
         "从硬编码到数据驱动", ACCENT_BLUE),
        ("Agent", "Tools · Budget · Trajectory",
         "有界决策 · 可审计 · 可回放", ACCENT_RED),
    ]
    interfaces = [
        "parse → SemanticIR L4",
        "Strategy 决策空间喂给 Agent",
        "Verifier 回注给 Agent / Pass 操作 IR",
        "仅通过工具协议访问 IR / Compiler",
    ]

    headers = ["件套", "承载物", "对 AI 原生的贡献", "与其它件套的接口"]
    col_xs = [Inches(0.9), Inches(2.5), Inches(5.3), Inches(9.1)]
    col_ws = [Inches(1.6), Inches(2.8), Inches(3.8), Inches(3.5)]
    panel(slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(0.5), BG_PANEL_ALT)
    for x, w, h in zip(col_xs, col_ws, headers):
        add_text(slide, x + Inches(0.1), Inches(1.98), w, Inches(0.35),
                 h, size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

    top = Inches(2.45)
    row_h = Inches(1.05)
    for i, (name, carry, contrib, col) in enumerate(rows):
        y = top + i * row_h
        panel(slide, Inches(0.9), y, Inches(11.5), row_h - Inches(0.08), BG_PANEL)
        accent_bar(slide, Inches(0.9), y, width=Inches(0.12),
                   height=row_h - Inches(0.08), color=col)
        add_text(slide, col_xs[0] + Inches(0.1), y + Inches(0.25),
                 col_ws[0], Inches(0.5), name, size=15, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, col_xs[1] + Inches(0.1), y + Inches(0.25),
                 col_ws[1], Inches(0.5), carry, size=11, color=FG,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, col_xs[2] + Inches(0.1), y + Inches(0.25),
                 col_ws[2], Inches(0.5), contrib, size=11, color=FG_MUTED,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, col_xs[3] + Inches(0.1), y + Inches(0.25),
                 col_ws[3], Inches(0.5), interfaces[i], size=11, color=FG_MUTED,
                 anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def make_verification(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "15 · 贯通", "三级验证 V0 / V1 / V2 在四件套中的归属")

    layers = [
        ("V0  静态验证  <1ms",
         "类型 · Shape · SSA · HW Constraint 即时剪枝",
         "归属 · IR + Pass", ACCENT),
        ("V1  数值验证",
         "SemanticInterpreter + Op Registry 逐 op 对参考实现",
         "归属 · Compiler Toolchain", ACCENT_ALT),
        ("V2  硬件性能验证",
         "真实 compile + profile + 对比 P0~P5 多级 baseline",
         "归属 · Backend + Benchmark", ACCENT_BLUE),
    ]
    card_w = Inches(3.9)
    card_h = Inches(4.2)
    gap = Inches(0.2)
    start = (SLIDE_W - card_w * 3 - gap * 2) / 2
    for i, (title, body, tag, col) in enumerate(layers):
        x = start + i * (card_w + gap)
        panel(slide, x, Inches(2.0), card_w, card_h, BG_PANEL)
        accent_bar(slide, x, Inches(2.0), width=card_w, height=Inches(0.1), color=col)
        add_text(slide, x + Inches(0.3), Inches(2.25), card_w - Inches(0.6), Inches(0.6),
                 title, size=17, bold=True, color=col)
        add_text(slide, x + Inches(0.3), Inches(3.0), card_w - Inches(0.6), Inches(2.3),
                 body, size=13, color=FG)
        tagbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x + Inches(0.3), Inches(5.4),
                                        card_w - Inches(0.6), Inches(0.5))
        tagbox.adjustments[0] = 0.3
        _set_fill(tagbox, BG_PANEL_ALT)
        _set_line(tagbox, col, 0.75)
        add_text(slide, x + Inches(0.3), Inches(5.4), card_w - Inches(0.6), Inches(0.5),
                 tag, size=12, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    panel(slide, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.55), BG_PANEL_ALT)
    add_text(slide, Inches(1.1), Inches(6.58), Inches(11.1), Inches(0.5),
             "门禁顺序 · Function ▸ Accuracy ▸ Performance  —— 让 LLM「大胆试、安全退」",
             size=13, color=FG, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def make_benchmark(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "15 · 贯通", "Benchmark 驱动 · Gate 把架构锁到可度量的坐标上")

    add_text(slide, Inches(0.9), Inches(1.9), Inches(12), Inches(0.5),
             "OT × ST → BL × L  坐标", size=14, bold=True, color=ACCENT)

    gates = [
        ("G2", "BL1 × L1", "matmul ≥ 70% cuBLAS", "语言 + IR MVP"),
        ("G3", "BL1 × L1", "LLM matmul ≥ 100% cuBLAS", "Agent 闭环"),
        ("G5", "BL3 × L1 + BL6/GPT-2 × L3", "GPT-2 top-1 正确", "端到端模型跑通"),
        ("G6", "BL4 × L1", "45 ops 正确 + ≥ 1.00× P3", "编译栈基础设施"),
        ("G7", "BL5 × L1+L2", "OT0–4 × ST1–4 + 4 融合", "Lang / IR v2 + 融合"),
        ("G8", "BL5 继承 + BL6 × L3", "LLaMA-2 ≥ 0.90× · DS-V2 ≥ 0.85×", "Agent 自主化"),
        ("G9", "BL6 × L3 (4 模型)", "Qwen ≥ 0.90× · Arke ≥ 1.05× P5", "Phase 1 v1.0"),
    ]
    col_xs = [Inches(1.0), Inches(1.95), Inches(4.85), Inches(8.85)]
    col_ws = [Inches(0.95), Inches(2.85), Inches(4.0), Inches(3.5)]
    headers = ["Gate", "BL × L", "关键指标", "对应四件套演进"]

    panel(slide, Inches(0.9), Inches(2.45), Inches(11.5), Inches(0.45), BG_PANEL_ALT)
    for x, w, h in zip(col_xs, col_ws, headers):
        add_text(slide, x, Inches(2.47), w, Inches(0.42),
                 h, size=12, bold=True, color=ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)

    row_h = Inches(0.55)
    for i, (g, bl, metric, stack) in enumerate(gates):
        y = Inches(2.95) + i * row_h
        bg = BG_PANEL if i % 2 == 0 else BG_PANEL_ALT
        panel(slide, Inches(0.9), y, Inches(11.5), row_h - Inches(0.05), bg)
        add_text(slide, col_xs[0], y + Inches(0.08), col_ws[0], Inches(0.4),
                 g, size=13, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, col_xs[1], y + Inches(0.08), col_ws[1], Inches(0.4),
                 bl, size=11, color=FG, anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
        add_text(slide, col_xs[2], y + Inches(0.08), col_ws[2], Inches(0.4),
                 metric, size=11, color=FG, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, col_xs[3], y + Inches(0.08), col_ws[3], Inches(0.4),
                 stack, size=11, color=FG_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def make_token_economics(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "15 · 贯通",
                 "Token 经济学  ·  为什么「让 LLM 做决策」比「让 LLM 写代码」便宜一个量级")

    data = [
        (".ak  (kernel only)", 72, "1×", ACCENT),
        (".ak  (kernel + strategy)", 160, "2×", ACCENT),
        ("LLM direct-write Triton", 563, "8×", ACCENT_ALT),
        ("Hand-tuned Triton (autotuned)", 1102, "15×", ACCENT_RED),
    ]
    top = Inches(2.1)
    bar_left = Inches(4.5)
    bar_max = Inches(7.2)
    max_tokens = max(d[1] for d in data)
    row_h = Inches(0.75)

    for i, (label, tokens, ratio, col) in enumerate(data):
        y = top + i * row_h
        panel(slide, Inches(0.9), y, Inches(11.5), row_h - Inches(0.12), BG_PANEL)
        add_text(slide, Inches(1.05), y + Inches(0.12), Inches(3.3), Inches(0.4),
                 label, size=13, bold=True, color=FG, anchor=MSO_ANCHOR.MIDDLE)
        bar_w = int(bar_max * (tokens / max_tokens))
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left,
                                     y + Inches(0.17), bar_w, Inches(0.3))
        _set_fill(bar, col)
        _no_line(bar)
        add_text(slide, bar_left + bar_w + Inches(0.15), y + Inches(0.15),
                 Inches(1.4), Inches(0.4), f"{tokens}  ·  {ratio}",
                 size=12, bold=True, color=FG,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)

    panel(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.4), BG_PANEL_ALT)
    add_text(slide, Inches(1.1), Inches(5.65), Inches(11), Inches(0.4),
             "结论", size=13, bold=True, color=ACCENT)
    add_bullets(slide, Inches(1.1), Inches(6.0), Inches(11), Inches(0.9), [
        "每次优化决策 ~10 tokens（结构化 action），不是 ~500 tokens（代码重写）",
        "编译器确定性验证替代多轮「修 bug」对话 → 每迭代从 ~3500 tokens 压到 ~500 tokens",
        "50 轮会话 · 分段 prompt cache 实测节省 ~176K tokens",
    ], size=12)
    footer(slide, page)


def make_roadmap(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "15 · 贯通", "路线图  ·  四件套与多后端同步演进")

    phases = [
        ("Phase 1", "Arke → Triton → NVIDIA GPU",
         "SIMT MVP  ✅  S6–S9 v1.0 推进中", ACCENT,
         "语言 v2 · IR 四层 · 编译栈 · Agent 自主化"),
        ("Phase 2", "Arke → Ascend Triton / NPU IR",
         "SIMD 泛化 · Cat B/D/E 完备", ACCENT_ALT,
         "@rationale 跨架构迁移验证  ≥10% lift"),
        ("Phase 3", "Arke → MLIR Dialect",
         "Triton 天花板拆除 · L2 决策", ACCENT_BLUE,
         "Agent 进入 loop-nest / memory layout 层"),
        ("Phase 4", "Arke → LLVM IR",
         "100% 硬件完整度 · 3+ 后端", ACCENT_RED,
         "Agent L3 决策：寄存器 / barrier / 指令调度"),
    ]
    top = Inches(2.0)
    row_h = Inches(1.18)
    for i, (ph, path, milestone, col, impact) in enumerate(phases):
        y = top + i * row_h
        panel(slide, Inches(0.9), y, Inches(11.5), row_h - Inches(0.15), BG_PANEL)
        accent_bar(slide, Inches(0.9), y, width=Inches(0.14),
                   height=row_h - Inches(0.15), color=col)
        add_text(slide, Inches(1.2), y + Inches(0.12), Inches(1.6), Inches(0.4),
                 ph, size=17, bold=True, color=col)
        add_text(slide, Inches(2.8), y + Inches(0.12), Inches(5.2), Inches(0.4),
                 path, size=13, bold=True, color=FG)
        add_text(slide, Inches(8.0), y + Inches(0.12), Inches(4.3), Inches(0.4),
                 milestone, size=11, color=FG_MUTED, anchor=MSO_ANCHOR.TOP)
        add_text(slide, Inches(2.8), y + Inches(0.55), Inches(9.5), Inches(0.4),
                 f"→  {impact}", size=12, color=ACCENT)
    footer(slide, page)


# ------------------------------------------------------------------ new blocks
# Industry insights · Required traits · Comparison · Architecture diagrams


def _arrow(slide, x1, y1, x2, y2, color=ACCENT, width_pt=1.5):
    """Draw a simple arrow between two points."""
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    # arrow head
    ln = line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


def _box(slide, left, top, width, height, text, *,
         fill=BG_PANEL, stroke=ACCENT, text_color=FG,
         size=12, bold=True, align=PP_ALIGN.CENTER,
         shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, stroke_pt=1.0):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        shp.adjustments[0] = 0.15
    _set_fill(shp, fill)
    _set_line(shp, stroke, stroke_pt)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _small_label(slide, left, top, width, text, *, color=FG_MUTED, size=10, bold=False):
    add_text(slide, left, top, width, Inches(0.25), text,
             size=size, bold=bold, color=color, align=PP_ALIGN.CENTER)


def make_industry_insight_1(prs, page):
    """Slide 1/2 of 业界技术洞察：四条主线趋势."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "07 · 业界洞察", "AI 编程编译前沿 · 四条主线正在交汇")

    items = [
        ("LLM-for-Code 的再定位", ACCENT, [
            "GPT-4 / Claude / Gemini 直写 CUDA / Triton 已可用但不稳定",
            "Meta CodeCompose · DeepMind AlphaCode · KernelBench 正基准化能力",
            "共识：LLM 作为「代码生成器」→「决策者 + 验证器」组合",
        ]),
        ("编译器自动调优的极限", ACCENT_ALT, [
            "Triton autotune · TVM MetaSchedule · AKG · Ansor 靠搜索空间",
            "搜索空间爆炸 · 启发式难泛化 · 无专家知识表达能力",
            "共识：需要把「人/LLM 的先验」接回调优循环",
        ]),
        ("MLIR 生态的繁荣与门槛", ACCENT_BLUE, [
            "NVIDIA Triton · IREE · torch-mlir · XLA · Modular Mojo 持续扩张",
            "dialect 爆炸 · C++/TableGen 门槛 · LLM 不可读",
            "共识：MLIR 是底座，但需要一层「LLM-facing」中间层",
        ]),
        ("Agentic AI Infra 范式成型", ACCENT_RED, [
            "Cursor · Claude Code · Devin · OpenClaw 验证 Tool-Use 范式",
            "MCP / OpenAI Tools / Anthropic Tool Use 协议趋同",
            "共识：结构化工具协议 + 预算 + Trajectory 成为 AI 系统标配",
        ]),
    ]
    card_w = Inches(5.9)
    card_h = Inches(2.35)
    positions = [
        (Inches(0.6), Inches(1.9)),
        (Inches(6.7), Inches(1.9)),
        (Inches(0.6), Inches(4.4)),
        (Inches(6.7), Inches(4.4)),
    ]
    for (x, y), (title, col, lines) in zip(positions, items):
        panel(slide, x, y, card_w, card_h, BG_PANEL)
        accent_bar(slide, x, y, width=Inches(0.12), height=card_h, color=col)
        add_text(slide, x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.4),
                 Inches(0.4), title, size=15, bold=True, color=col)
        add_bullets(slide, x + Inches(0.3), y + Inches(0.55),
                    card_w - Inches(0.4), card_h - Inches(0.7),
                    lines, size=11, bullet_color=col, line_spacing=1.18)
    footer(slide, page)


def make_industry_insight_2(prs, page):
    """Slide 2/2: 业界方案矩阵."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "07 · 业界洞察", "业界方案全景 · 各自解决一部分问题，但缺「AI-Native」整体性")

    rows = [
        ("Triton (OpenAI/NVIDIA)", "GPU Python DSL + autotune",
         "门槛低 · 生态广", "抽象天花板 · LLM 不友好 · 人写人调"),
        ("TVM / Ansor / MetaSchedule", "搜索 + cost model + schedule",
         "搜索泛化 · 多后端", "搜索空间大 · 专家知识难注入"),
        ("MLIR + Linalg / Transform", "多 dialect · 可组合 lowering",
         "工业级编译基础设施", "C++ / TableGen 门槛 · 非 LLM-Native"),
        ("FlagGems / Liger-Kernel", "Triton 算子库（200+）",
         "现成高性能实现 · 模型可直用", "手写维护 · 无决策/验证协议"),
        ("torch.compile / Inductor", "PyTorch 内置图编译 + Triton",
         "易集成 · 自动化", "黑盒 · 难定制 · 硬件扩展受限"),
        ("Mojo / Hidet", "新语言 / 新图编译",
         "可控性强 · 性能好", "生态新 · 缺 LLM-Native 协议"),
        ("Agent 框架 (Cursor · CC · Devin)", "Tool-Use + Context 工程",
         "闭环成熟 · token 工程化", "通用 agent · 非编译专用"),
    ]
    col_xs = [Inches(0.9), Inches(4.4), Inches(7.7), Inches(10.4)]
    col_ws = [Inches(3.4), Inches(3.2), Inches(2.6), Inches(2.4)]
    headers = ["方案 / 代表", "定位", "优势", "局限（从 AI-Native 视角）"]

    panel(slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(0.45), BG_PANEL_ALT)
    for x, w, h in zip(col_xs, col_ws, headers):
        add_text(slide, x, Inches(1.94), w, Inches(0.36),
                 h, size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

    row_h = Inches(0.61)
    for i, row in enumerate(rows):
        y = Inches(2.4) + i * row_h
        bg = BG_PANEL if i % 2 == 0 else BG_PANEL_ALT
        panel(slide, Inches(0.9), y, Inches(11.5), row_h - Inches(0.05), bg)
        for j, cell in enumerate(row):
            color = FG if j == 0 else (FG_MUTED if j == 3 else FG)
            bold = j == 0
            add_text(slide, col_xs[j], y + Inches(0.12), col_ws[j], Inches(0.4),
                     cell, size=10, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)

    panel(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5), BG_PANEL_ALT)
    add_text(slide, Inches(1.1), Inches(6.75), Inches(11), Inches(0.4),
             "结论 · 没有任何现有方案同时具备：LLM-Native 入口 + 多层 IR + 决策验证闭环 + 跨硬件泛化",
             size=12, bold=True, color=ACCENT_ALT, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def make_required_traits(prs, page):
    """应有特征：六条."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "08 · 应有特征", "AI 原生算子编译栈的 6 个必备特征")

    traits = [
        ("T1", "LLM-Native 入口", ACCENT,
         "人 / LLM 共同可读可写的算子级语言，算法无关、token 最小化"),
        ("T2", "语义 / 策略解耦", ACCENT_ALT,
         "What 与 How 独立演化，正确性与性能分别验证"),
        ("T3", "多层 IR · 按参与度分层", ACCENT_BLUE,
         "从 Semantic 到 Instruction，每层 LLM 参与粒度清晰"),
        ("T4", "有界动作空间 + 结构化协议", ACCENT_RED,
         "LLM 做选择不写代码 · 编译器枚举合法动作 · Tool-Use API"),
        ("T5", "Compiler-as-Verifier 三级闭环", ACCENT,
         "V0 静态 / V1 数值 / V2 性能 · checkpoint / rollback 安全探索"),
        ("T6", "多硬件可泛化 + 知识可迁移", ACCENT_ALT,
         "单一 Semantic 对多后端 · @rationale 跨架构可复用"),
    ]
    card_w = Inches(5.95)
    card_h = Inches(1.6)
    positions = [(i % 2, i // 2) for i in range(6)]
    for i, (tag, title, col, desc) in enumerate(traits):
        cx, cy = positions[i]
        x = Inches(0.6) + cx * Inches(6.2)
        y = Inches(1.95) + cy * Inches(1.72)
        panel(slide, x, y, card_w, card_h, BG_PANEL)
        accent_bar(slide, x, y, width=Inches(0.12), height=card_h, color=col)
        add_text(slide, x + Inches(0.3), y + Inches(0.1), Inches(1.0), Inches(0.4),
                 tag, size=16, bold=True, color=col, font=FONT_MONO)
        add_text(slide, x + Inches(1.25), y + Inches(0.1), card_w - Inches(1.4),
                 Inches(0.45), title, size=15, bold=True, color=FG)
        add_text(slide, x + Inches(0.3), y + Inches(0.7), card_w - Inches(0.45),
                 Inches(0.85), desc, size=11, color=FG_MUTED)

    panel(slide, Inches(0.6), Inches(7.02), Inches(12.15), Inches(0.0), BG_PANEL_ALT)
    footer(slide, page)


def make_comparison(prs, page):
    """方案对比表：在 T1–T6 六个维度上对比各方案."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "09 · 方案对比", "在 6 个维度上 · 本方案 vs 业界代表方案")

    traits_short = ["T1 入口", "T2 语义/策略", "T3 多层 IR", "T4 有界动作", "T5 三级验证", "T6 跨硬件"]
    #  '●' 完全满足   '◐' 部分满足   '○' 不满足
    FULL, PART, NONE = "●", "◐", "○"
    rows = [
        ("本方案 · Arke",               [FULL, FULL, FULL, FULL, FULL, FULL], ACCENT),
        ("Triton + autotune",          [PART, NONE, NONE, NONE, PART, PART], FG_MUTED),
        ("TVM / Ansor",                [NONE, PART, PART, NONE, PART, FULL], FG_MUTED),
        ("MLIR + Linalg / Transform",  [NONE, PART, FULL, NONE, PART, FULL], FG_MUTED),
        ("FlagGems / Liger",           [NONE, NONE, NONE, NONE, NONE, PART], FG_MUTED),
        ("torch.compile / Inductor",   [NONE, NONE, PART, NONE, PART, PART], FG_MUTED),
        ("LLM 直写 Triton",            [PART, NONE, NONE, NONE, NONE, NONE], FG_MUTED),
        ("Cursor / CC 通用 Agent",     [NONE, NONE, NONE, PART, NONE, NONE], FG_MUTED),
    ]

    # header row
    col_name_w = Inches(3.2)
    cell_w = Inches(1.35)
    start_x = Inches(0.8)
    header_y = Inches(1.95)
    panel(slide, start_x, header_y, col_name_w + cell_w * 6 + Inches(0.2),
          Inches(0.45), BG_PANEL_ALT)
    add_text(slide, start_x + Inches(0.1), header_y + Inches(0.05), col_name_w,
             Inches(0.35), "方案 / 维度", size=11, bold=True, color=ACCENT,
             anchor=MSO_ANCHOR.MIDDLE)
    for j, th in enumerate(traits_short):
        x = start_x + col_name_w + j * cell_w
        add_text(slide, x, header_y + Inches(0.05), cell_w, Inches(0.35),
                 th, size=10, bold=True, color=ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # rows
    row_h = Inches(0.48)
    for i, (name, marks, _) in enumerate(rows):
        y = header_y + Inches(0.5) + i * row_h
        highlight = i == 0
        bg = BG_PANEL if not highlight else BG_PANEL_ALT
        panel(slide, start_x, y, col_name_w + cell_w * 6 + Inches(0.2),
              row_h - Inches(0.05), bg)
        if highlight:
            accent_bar(slide, start_x, y, width=Inches(0.1),
                       height=row_h - Inches(0.05), color=ACCENT)
        add_text(slide, start_x + Inches(0.2), y + Inches(0.06), col_name_w,
                 Inches(0.35), name, size=11,
                 bold=highlight, color=ACCENT if highlight else FG,
                 anchor=MSO_ANCHOR.MIDDLE)
        for j, m in enumerate(marks):
            x = start_x + col_name_w + j * cell_w
            color = {FULL: ACCENT, PART: ACCENT_ALT, NONE: ACCENT_RED}[m]
            add_text(slide, x, y + Inches(0.06), cell_w, Inches(0.35),
                     m, size=18, bold=True, color=color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font=FONT_MONO)

    # legend
    lg_y = Inches(6.55)
    legend = [("●", ACCENT, "完全满足"), ("◐", ACCENT_ALT, "部分满足"),
              ("○", ACCENT_RED, "不满足")]
    lx = Inches(0.8)
    for sym, col, txt in legend:
        add_text(slide, lx, lg_y, Inches(0.3), Inches(0.3),
                 sym, size=16, bold=True, color=col, font=FONT_MONO)
        add_text(slide, lx + Inches(0.35), lg_y + Inches(0.04), Inches(1.2),
                 Inches(0.3), txt, size=10, color=FG_MUTED)
        lx += Inches(1.6)
    footer(slide, page)


def make_advantage_detail(prs, page):
    """六条优势详解（配套对比表使用）."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "09 · 方案对比", "六个维度上的优势简述 —— 为什么能在所有维度上「完全满足」")

    items = [
        ("T1 · LLM-Native 入口", ACCENT,
         ".ak 语言算子级抽象 · 算法无关 · 72–160 tokens · 人 / LLM 共同可写",
         "业界多数方案以 Python 或 C++ 为入口，LLM 需兼顾语义 + 硬件细节"),
        ("T2 · 语义 / 策略解耦", ACCENT_ALT,
         "kernel (semantics) ↔ strategy (@rationale decisions) 句法级分离",
         "Triton / CUDA 策略与实现混写；TVM 通过 schedule 部分分离但仍耦合"),
        ("T3 · 多层 IR 按参与度分层", ACCENT_BLUE,
         "Semantic → Strategy → Schedule → Instruction · 每层 LLM 参与度明确",
         "MLIR 多 dialect 但非按 LLM 参与度切；LLVM / Triton IR 单层密集"),
        ("T4 · 有界动作空间 + 结构化协议", ACCENT_RED,
         "compiler 枚举 legal_actions · Agent 选择 · tool-use 闭环",
         "业界 agent 方案多为自由文本生成，缺编译器先验约束"),
        ("T5 · Compiler-as-Verifier 三级闭环", ACCENT,
         "V0 <1ms · V1 对参考 · V2 HW profile · checkpoint / rollback",
         "autotune 只验证性能；直写方案事后对拍；均不具备 V0 即时剪枝"),
        ("T6 · 多硬件可泛化 + 知识可迁移", ACCENT_ALT,
         "单 Semantic 对多后端 · @rationale 以自然语言跨架构迁移",
         "FlagGems 每硬件一份实现；TVM 跨架构但知识不显式"),
    ]
    card_w = Inches(5.95)
    card_h = Inches(1.58)
    for i, (title, col, pro, vs) in enumerate(items):
        cx, cy = i % 2, i // 2
        x = Inches(0.6) + cx * Inches(6.2)
        y = Inches(1.95) + cy * Inches(1.68)
        panel(slide, x, y, card_w, card_h, BG_PANEL)
        accent_bar(slide, x, y, width=Inches(0.12), height=card_h, color=col)
        add_text(slide, x + Inches(0.3), y + Inches(0.08), card_w - Inches(0.4),
                 Inches(0.35), title, size=13, bold=True, color=col)
        add_text(slide, x + Inches(0.3), y + Inches(0.5), card_w - Inches(0.4),
                 Inches(0.5), "▸ " + pro, size=10, color=FG)
        add_text(slide, x + Inches(0.3), y + Inches(1.02), card_w - Inches(0.4),
                 Inches(0.5), "⋯ vs " + vs, size=10, color=FG_MUTED)
    footer(slide, page)


def make_overall_arch(prs, page):
    """Arke 整体技术架构图."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "10 · 整体架构", "Arke 整体技术架构 · 从用户入口到硬件执行")

    # layer backdrops
    layers = [
        ("用户入口层", ACCENT, Inches(0.7), Inches(1.9), Inches(0.6)),
        ("Language 层", ACCENT, Inches(0.7), Inches(2.6), Inches(0.6)),
        ("IR 层（四层）", ACCENT_ALT, Inches(0.7), Inches(3.3), Inches(1.35)),
        ("Agent & Compiler Toolchain 层", ACCENT_BLUE,
         Inches(0.7), Inches(4.75), Inches(1.3)),
        ("Backend 层", ACCENT_RED, Inches(0.7), Inches(6.15), Inches(0.55)),
        ("硬件执行层", ACCENT_RED, Inches(0.7), Inches(6.78), Inches(0.45)),
    ]
    for name, col, x, y, h in layers:
        p = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.8), h)
        _set_fill(p, BG_PANEL_ALT)
        _set_line(p, col, 1.0)
        add_text(slide, x, y, Inches(1.8), h, name,
                 size=10, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    content_x = Inches(2.7)
    content_w = Inches(10.0)

    # 用户入口层
    entries = [".ak 文件", "自然语言", "CLI 参数", "Python API"]
    ew = (content_w - Inches(0.3)) / 4
    for i, e in enumerate(entries):
        x = content_x + i * (ew + Inches(0.1))
        _box(slide, x, Inches(1.95), ew, Inches(0.5), e,
             fill=BG_PANEL, stroke=ACCENT, size=11, text_color=FG)

    # Language 层
    _box(slide, content_x, Inches(2.65), content_w, Inches(0.5),
         "Arke Language v2  ·  kernel { semantics }  +  strategy { @rationale decisions }  ·  where 符号维度",
         fill=BG_PANEL, stroke=ACCENT, size=11)

    # IR 层 — 四个横向框
    ir_layers = [
        ("L4  SemanticIR", "what to compute", ACCENT_ALT),
        ("L3  StrategyIR", "how to optimize", ACCENT_ALT),
        ("L2  ScheduleIR", "schedule mapping", ACCENT_ALT),
        ("L1  InstructionIR", "near-LLVM IR", ACCENT_ALT),
    ]
    iw = (content_w - Inches(0.3)) / 4
    for i, (t, sub, col) in enumerate(ir_layers):
        x = content_x + i * (iw + Inches(0.1))
        _box(slide, x, Inches(3.4), iw, Inches(0.55), t,
             fill=BG_PANEL, stroke=col, size=11)
        add_text(slide, x, Inches(3.98), iw, Inches(0.3), sub,
                 size=9, color=FG_MUTED, align=PP_ALIGN.CENTER)
    # LLM 参与度 bar
    add_text(slide, content_x, Inches(4.28), content_w, Inches(0.28),
             "← LLM 参与度：主作者  /  主决策者  /  review only  /  不参与 →",
             size=9, color=FG_MUTED, align=PP_ALIGN.CENTER)

    # Agent & Compiler 层
    agent_w = content_w * 0.45
    comp_w = content_w - agent_w - Inches(0.2)
    # Agent block
    _box(slide, content_x, Inches(4.85), agent_w, Inches(1.1),
         "Agent  ·  Bounded Action Space\nTool-Use API · Budget · Trajectory · @rationale",
         fill=BG_PANEL, stroke=ACCENT_BLUE, size=10)
    # Compiler block
    cx = content_x + agent_w + Inches(0.2)
    _box(slide, cx, Inches(4.85), comp_w, Inches(1.1),
         "Compiler Toolchain  ·  OpRegistry · Pass Pipeline · BackendRegistry\nV0 静态 · V1 数值 · V2 性能  三级验证",
         fill=BG_PANEL, stroke=ACCENT_BLUE, size=10)

    # Backend 层
    backends = ["Triton  (Phase 1–2)", "MLIR Dialect  (Phase 3)", "LLVM IR  (Phase 4)"]
    bw = (content_w - Inches(0.2)) / 3
    for i, b in enumerate(backends):
        x = content_x + i * (bw + Inches(0.1))
        _box(slide, x, Inches(6.2), bw, Inches(0.45), b,
             fill=BG_PANEL, stroke=ACCENT_RED, size=10)

    # Hardware 层
    hws = ["NVIDIA GPU", "Huawei Ascend NPU", "AMD GPU", "未来 NPU / DSA"]
    hw_w = (content_w - Inches(0.3)) / 4
    for i, h in enumerate(hws):
        x = content_x + i * (hw_w + Inches(0.1))
        _box(slide, x, Inches(6.8), hw_w, Inches(0.4), h,
             fill=BG_PANEL_ALT, stroke=ACCENT_RED, size=10, bold=False)

    footer(slide, page)


def make_lang_arch(prs, page):
    """Language 件套流程图."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "11 · 件套架构 · Language", ".ak 源文件  →  Parser  →  AST  →  SemanticIR / StrategyIR")

    y0 = Inches(2.0)
    stages = [
        (".ak 源文件", "kernel { ... }\nstrategy { @rationale ... }", ACCENT),
        ("Lexer + Parser", "EBNF 文法  ·  正则词法\n歧义 0 · token 最小化", ACCENT),
        ("AST", "kernel_def · strategy_def\nwhere_clause · @rationale", ACCENT),
        ("Resolver", "Op Registry 查表\n符号维度求解 · 类型推断", ACCENT),
        ("Layer 4\nSemanticIR", "纯数学 DAG\nSSA by construction", ACCENT_ALT),
        ("Layer 3\nStrategyIR", "决策列表 + @rationale\n可枚举 · 可回滚", ACCENT_ALT),
    ]
    w = Inches(1.9)
    h = Inches(1.6)
    gap = Inches(0.12)
    start_x = (SLIDE_W - (w * 6 + gap * 5)) / 2
    prev_right = None
    for i, (title, body, col) in enumerate(stages):
        x = start_x + i * (w + gap)
        panel(slide, x, y0, w, h, BG_PANEL)
        accent_bar(slide, x, y0, width=w, height=Inches(0.08), color=col)
        add_text(slide, x + Inches(0.1), y0 + Inches(0.15), w - Inches(0.2),
                 Inches(0.5), title, size=12, bold=True, color=col,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), y0 + Inches(0.65), w - Inches(0.2),
                 h - Inches(0.7), body, size=9, color=FG_MUTED,
                 align=PP_ALIGN.CENTER)
        if prev_right is not None:
            mid_y = y0 + h / 2
            _arrow(slide, prev_right, mid_y, x, mid_y, color=ACCENT, width_pt=1.25)
        prev_right = x + w

    # bottom: core design
    panel(slide, Inches(0.7), Inches(4.1), Inches(11.9), Inches(2.8), BG_PANEL)
    add_text(slide, Inches(0.95), Inches(4.25), Inches(11.5), Inches(0.4),
             "核心设计要素", size=13, bold=True, color=ACCENT)
    design = [
        ("算子级抽象", "不暴露 loop / thread / memory"),
        ("语义 / 策略双段式", "kernel vs strategy 句法级切开"),
        ("Universal Op", "语言不枚举算子 · Op Registry 注入"),
        ("符号维度", "where B: dynamic(max=64) 一等公民"),
        ("@rationale", "每决策强制带自然语言理由"),
        ("Token 最小化", "`.ak` 72 tokens · Triton 1102 tokens"),
    ]
    card_w = Inches(3.75)
    card_h = Inches(0.95)
    for i, (k, v) in enumerate(design):
        cx, cy = i % 3, i // 3
        x = Inches(0.95) + cx * (card_w + Inches(0.15))
        y = Inches(4.75) + cy * (card_h + Inches(0.15))
        panel(slide, x, y, card_w, card_h, BG_PANEL_ALT)
        accent_bar(slide, x, y, width=Inches(0.08), height=card_h, color=ACCENT)
        add_text(slide, x + Inches(0.2), y + Inches(0.08), card_w - Inches(0.3),
                 Inches(0.35), k, size=11, bold=True, color=ACCENT)
        add_text(slide, x + Inches(0.2), y + Inches(0.42), card_w - Inches(0.3),
                 Inches(0.5), v, size=10, color=FG_MUTED)

    footer(slide, page)


def make_ir_arch(prs, page):
    """IR 件套架构图."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "12 · 件套架构 · IR", "多层 IR · 职责 × LLM 参与度 × 验证点")

    layers = [
        ("Layer 4 · SemanticIR",   "什么",   "纯数学 DAG · SSA 不可变",          "LLM 主作者", "V1 数值参考", ACCENT),
        ("Layer 3 · StrategyIR",   "怎么优化", "决策列表 · @rationale · 可回滚",   "LLM 主决策者", "V0 静态", ACCENT_ALT),
        ("Layer 2 · ScheduleIR",   "怎么映射", "thread / block / warp / vector",  "LLM review", "V0+V2", ACCENT_BLUE),
        ("Layer 1 · InstructionIR", "怎么成指令", "近 LLVM IR",                 "LLM 不参与", "V2 性能", ACCENT_RED),
    ]

    col_xs = [Inches(0.7), Inches(4.1), Inches(5.4), Inches(8.6), Inches(10.4), Inches(12.1)]
    col_ws = [Inches(3.35), Inches(1.25), Inches(3.15), Inches(1.75), Inches(1.65)]
    headers = ["IR 层", "职责", "承载内容", "LLM 参与度", "主验证级"]

    panel(slide, Inches(0.7), Inches(1.9), Inches(11.95), Inches(0.45), BG_PANEL_ALT)
    for x, w, h in zip(col_xs, col_ws, headers):
        add_text(slide, x, Inches(1.95), w, Inches(0.36),
                 h, size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

    row_h = Inches(0.78)
    for i, (name, duty, body, llm, v, col) in enumerate(layers):
        y = Inches(2.4) + i * row_h
        panel(slide, Inches(0.7), y, Inches(11.95), row_h - Inches(0.08), BG_PANEL)
        accent_bar(slide, Inches(0.7), y, width=Inches(0.12),
                   height=row_h - Inches(0.08), color=col)
        vals = [name, duty, body, llm, v]
        for j, (x, w, val) in enumerate(zip(col_xs, col_ws, vals)):
            align = PP_ALIGN.LEFT
            size = 12 if j == 0 else (11 if j == 2 else 11)
            bold = j == 0
            color = col if j == 0 else FG
            add_text(slide, x + Inches(0.08), y + Inches(0.18), w, Inches(0.45),
                     val, size=size, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)

    # side arrow: LLM 参与度渐退
    add_text(slide, Inches(0.7), Inches(5.7), Inches(11.95), Inches(0.3),
             "↓  lowering 方向  ↓     LLM 参与度递减      V0→V1→V2 验证递进",
             size=11, bold=True, color=FG_MUTED, align=PP_ALIGN.CENTER)

    # bottom: 下接 MLIR / LLVM IR
    _box(slide, Inches(0.7), Inches(6.15), Inches(5.8), Inches(0.55),
         "下接  ·  MLIR Standard Dialects (linalg / transform / scf / gpu)",
         fill=BG_PANEL_ALT, stroke=ACCENT_BLUE, size=11)
    _box(slide, Inches(6.85), Inches(6.15), Inches(5.8), Inches(0.55),
         "下接  ·  LLVM IR  →  PTX / ISA",
         fill=BG_PANEL_ALT, stroke=ACCENT_RED, size=11)
    footer(slide, page)


def make_compiler_arch(prs, page):
    """Compiler Toolchain 架构图."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "13 · 件套架构 · Compiler Toolchain",
                 "OpRegistry（SSOT）× Pass Pipeline × BackendRegistry × 三级验证器")

    # Central: Pass Pipeline horizontal chain
    chain = [
        ("Parse",            ACCENT),
        ("ShapeInference",   ACCENT),
        ("SemanticVerify",   ACCENT),
        ("StrategyLower",    ACCENT_ALT),
        ("Schedule",         ACCENT_ALT),
        ("Codegen",          ACCENT_BLUE),
        ("Profile",          ACCENT_RED),
    ]
    chain_y = Inches(3.3)
    w = Inches(1.55)
    h = Inches(0.75)
    gap = Inches(0.1)
    start_x = (SLIDE_W - (w * len(chain) + gap * (len(chain) - 1))) / 2
    prev_right = None
    for i, (label, col) in enumerate(chain):
        x = start_x + i * (w + gap)
        _box(slide, x, chain_y, w, h, label, fill=BG_PANEL, stroke=col, size=11)
        if prev_right is not None:
            _arrow(slide, prev_right, chain_y + h / 2, x, chain_y + h / 2,
                   color=FG_MUTED, width_pt=1.2)
        prev_right = x + w

    add_text(slide, start_x, chain_y - Inches(0.4),
             w * len(chain) + gap * (len(chain) - 1), Inches(0.3),
             "Pass Pipeline  ·  Analysis  /  Transform  /  Verification 三族 · 可组合 · 可替换",
             size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Above: OpRegistry SSOT
    reg_y = Inches(1.9)
    reg_w = Inches(11.9)
    reg_h = Inches(0.95)
    panel(slide, Inches(0.7), reg_y, reg_w, reg_h, BG_PANEL)
    accent_bar(slide, Inches(0.7), reg_y, width=reg_w, height=Inches(0.08),
               color=ACCENT)
    add_text(slide, Inches(0.9), reg_y + Inches(0.13), reg_w - Inches(0.3),
             Inches(0.35), "OpRegistry  ·  Single Source of Truth",
             size=13, bold=True, color=ACCENT)
    add_text(slide, Inches(0.9), reg_y + Inches(0.48), reg_w - Inches(0.3),
             Inches(0.45),
             "signature · shape_rule · template_hint · reference_impl · input_gen · attrs · strategy_space · hw_variants",
             size=10, color=FG_MUTED)

    # arrows from SSOT down to each pipeline stage that consumes it
    for i in (0, 1, 2, 4, 5):  # Parse/ShapeInf/SemV/Schedule/Codegen 都依赖 Registry
        x = start_x + i * (w + gap) + w / 2
        _arrow(slide, x, reg_y + reg_h, x, chain_y, color=ACCENT, width_pt=0.9)

    # Below: BackendRegistry
    bx_y = Inches(4.4)
    bx_h = Inches(0.8)
    panel(slide, Inches(0.7), bx_y, reg_w, bx_h, BG_PANEL)
    accent_bar(slide, Inches(0.7), bx_y, width=Inches(0.12),
               height=bx_h, color=ACCENT_BLUE)
    add_text(slide, Inches(0.9), bx_y + Inches(0.1), Inches(4),
             Inches(0.35), "BackendRegistry",
             size=13, bold=True, color=ACCENT_BLUE)
    # three backends inside
    bk_x0 = Inches(4.8)
    bk_w = Inches(2.5)
    for i, b in enumerate(["Triton Backend", "MLIR Backend", "LLVM IR Backend"]):
        x = bk_x0 + i * (bk_w + Inches(0.1))
        _box(slide, x, bx_y + Inches(0.2), bk_w, Inches(0.42),
             b, fill=BG_PANEL_ALT, stroke=ACCENT_BLUE, size=10)

    # Verifier row (three cards)
    v_y = Inches(5.4)
    v_w = Inches(3.85)
    verifiers = [
        ("V0  静态 · <1ms",
         "类型 / Shape / SSA / HW constraint", ACCENT),
        ("V1  数值",
         "SemanticInterpreter + 参考对拍", ACCENT_ALT),
        ("V2  性能",
         "compile + profile + P0~P5 baseline", ACCENT_RED),
    ]
    for i, (t, b, col) in enumerate(verifiers):
        x = Inches(0.7) + i * (v_w + Inches(0.15))
        panel(slide, x, v_y, v_w, Inches(1.0), BG_PANEL)
        accent_bar(slide, x, v_y, width=Inches(0.1), height=Inches(1.0), color=col)
        add_text(slide, x + Inches(0.25), v_y + Inches(0.1), v_w - Inches(0.3),
                 Inches(0.35), t, size=12, bold=True, color=col)
        add_text(slide, x + Inches(0.25), v_y + Inches(0.45), v_w - Inches(0.3),
                 Inches(0.5), b, size=10, color=FG_MUTED)

    # bottom band
    panel(slide, Inches(0.7), Inches(6.55), Inches(11.95), Inches(0.45), BG_PANEL_ALT)
    add_text(slide, Inches(0.9), Inches(6.58), Inches(11.6), Inches(0.4),
             "门禁顺序  ·  Function ▸ Accuracy ▸ Performance    /    Checkpoint · Rollback · 让 LLM 「大胆试、安全退」",
             size=11, bold=True, color=FG, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def make_agent_arch(prs, page):
    """Agent 件套架构图."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "14 · 件套架构 · Agent Engineering",
                 "LLM ↔ ArkeEnv 结构化闭环  ·  Tool-Use × Budget × Trajectory")

    # Left: LLM Agent box
    lx = Inches(0.7)
    ly = Inches(2.0)
    lw = Inches(3.2)
    lh = Inches(4.5)
    panel(slide, lx, ly, lw, lh, BG_PANEL)
    accent_bar(slide, lx, ly, width=Inches(0.12), height=lh, color=ACCENT_RED)
    add_text(slide, lx + Inches(0.25), ly + Inches(0.15), lw - Inches(0.3),
             Inches(0.4), "LLM Agent", size=15, bold=True, color=ACCENT_RED)
    add_text(slide, lx + Inches(0.25), ly + Inches(0.55), lw - Inches(0.3),
             Inches(0.35), "Decides (not generates code)",
             size=10, color=FG_MUTED)
    # inside items
    steps = [
        ("1", "analyze kernel  ·  read HW profile"),
        ("2", "list_legal_actions()"),
        ("3", "apply_decision(... @rationale)"),
        ("4", "verify → compile → profile"),
        ("5", "iterate  /  rollback  /  done"),
    ]
    for i, (n, t) in enumerate(steps):
        y = ly + Inches(1.1) + i * Inches(0.6)
        panel(slide, lx + Inches(0.25), y, lw - Inches(0.5), Inches(0.5), BG_PANEL_ALT)
        add_text(slide, lx + Inches(0.35), y + Inches(0.08), Inches(0.3),
                 Inches(0.35), n, size=11, bold=True, color=ACCENT_RED,
                 font=FONT_MONO)
        add_text(slide, lx + Inches(0.7), y + Inches(0.08),
                 lw - Inches(1.0), Inches(0.35), t,
                 size=10, color=FG)

    # Right: ArkeEnv box
    rx = Inches(9.3)
    rw = Inches(3.4)
    ry = ly
    rh = lh
    panel(slide, rx, ry, rw, rh, BG_PANEL)
    accent_bar(slide, rx, ry, width=Inches(0.12), height=rh, color=ACCENT_BLUE)
    add_text(slide, rx + Inches(0.25), ry + Inches(0.15), rw - Inches(0.3),
             Inches(0.4), "ArkeEnv", size=15, bold=True, color=ACCENT_BLUE)
    add_text(slide, rx + Inches(0.25), ry + Inches(0.55), rw - Inches(0.3),
             Inches(0.35), "Verifies + Executes", size=10, color=FG_MUTED)
    tools = [
        "analyze_compute / get_hw_profile",
        "list_legal_actions",
        "apply_decision",
        "verify_correctness  (V1)",
        "compile_and_profile (V2)",
        "checkpoint / rollback",
    ]
    for i, t in enumerate(tools):
        y = ry + Inches(1.1) + i * Inches(0.5)
        panel(slide, rx + Inches(0.25), y, rw - Inches(0.5), Inches(0.42), BG_PANEL_ALT)
        add_text(slide, rx + Inches(0.35), y + Inches(0.05),
                 rw - Inches(0.7), Inches(0.35), t,
                 size=10, color=FG, font=FONT_MONO)

    # Middle arrows (two directions, annotated)
    mid_y1 = Inches(3.2)
    mid_y2 = Inches(4.6)
    _arrow(slide, lx + lw, mid_y1, rx, mid_y1, color=ACCENT_RED, width_pt=1.5)
    _arrow(slide, rx, mid_y2, lx + lw, mid_y2, color=ACCENT_BLUE, width_pt=1.5)
    add_text(slide, lx + lw, mid_y1 - Inches(0.32),
             rx - (lx + lw), Inches(0.3),
             "tool_use  (structured API)",
             size=11, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
    add_text(slide, lx + lw, mid_y2 + Inches(0.02),
             rx - (lx + lw), Inches(0.3),
             "tool_result (legal actions · verification · profile)",
             size=11, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # Middle bottom: Trajectory + Budget + Prompt Cache
    mx = lx + lw + Inches(0.2)
    mw = rx - mx - Inches(0.2)
    my = Inches(5.3)
    mh = Inches(1.1)
    panel(slide, mx, my, mw, mh, BG_PANEL_ALT)
    add_text(slide, mx + Inches(0.2), my + Inches(0.08), mw - Inches(0.4),
             Inches(0.35),
             "Optimization State  ·  Ground Truth 不被压缩",
             size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    sub = [
        ("Strategy IR", ACCENT),
        ("Compile Results", ACCENT_ALT),
        ("Decision Log + @rationale", ACCENT_BLUE),
        ("Budget  (decisions / compiles)", ACCENT_RED),
    ]
    ssw = (mw - Inches(0.5)) / 4
    for i, (t, col) in enumerate(sub):
        sx = mx + Inches(0.15) + i * (ssw + Inches(0.1))
        _box(slide, sx, my + Inches(0.5), ssw, Inches(0.45),
             t, fill=BG_PANEL, stroke=col, size=9)

    # Bottom: dual mode integration + trajectory
    panel(slide, Inches(0.7), Inches(6.55), Inches(11.95), Inches(0.45), BG_PANEL_ALT)
    add_text(slide, Inches(0.9), Inches(6.58), Inches(11.6), Inches(0.4),
             "双模集成  ·  Mode A 内置 (CLI / Python API)    /    Mode B 外部 (Cursor · Claude Code · OpenClaw)    ·    Trajectory JSONL 全量落盘",
             size=11, bold=True, color=FG, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def make_closing(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    for y in (0, Inches(7.42)):
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, Inches(0.08))
        _set_fill(b, ACCENT)
        _no_line(b)

    add_text(slide, Inches(0.9), Inches(1.3), Inches(12), Inches(0.6),
             "THANK YOU", size=18, bold=True, color=ACCENT)
    add_text(slide, Inches(0.9), Inches(1.9), Inches(12), Inches(1.4),
             "Questions & Discussion", size=52, bold=True, color=FG)

    prompts = [
        ("从零开始", "你会先建四件套里的哪一件？为什么？"),
        ("动作粒度", "Bounded Action Space 在 Triton→MLIR→LLVM 各阶段如何动态扩张？"),
        ("验证投入", "V0 / V1 / V2 哪一级最值得重仓？哪一级最容易被过度工程？"),
        ("知识沉淀", "@rationale 应该是检索式 / RAG / 微调？如何真正形成跨硬件知识库？"),
    ]
    top = Inches(3.6)
    for i, (tag, q) in enumerate(prompts):
        y = top + (i // 2) * Inches(1.3)
        x = Inches(0.9) + (i % 2) * Inches(5.9)
        panel(slide, x, y, Inches(5.7), Inches(1.15), BG_PANEL)
        accent_bar(slide, x, y, width=Inches(0.12), height=Inches(1.15), color=ACCENT)
        add_text(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.3), Inches(0.35),
                 tag, size=12, bold=True, color=ACCENT)
        add_text(slide, x + Inches(0.3), y + Inches(0.5), Inches(5.3), Inches(0.6),
                 q, size=13, color=FG)

    add_text(slide, Inches(0.9), Inches(6.7), Inches(12), Inches(0.4),
             "AI-Native Compile Stack  ·  Language × IR × Compiler × Agent  ·  案例参照 Arke Project",
             size=11, color=FG_MUTED)
    footer(slide, page)


# ------------------------------------------------------------------ content bullets

LANG_NECESSITY = [
    "Python / Triton 暴露循环 / 线程 / 地址 → LLM 需同时懂 Python 语义 + CUDA 概念",
    "一个简单 matmul Triton 手写 ~1102 tokens，.ak ~72 tokens  (15×)",
    ("正确率：LLM 直写 Triton ≈ 83%，结构化协议 ≈ 100%", 1),
    ("Token 成本：~3500 tokens / 迭代 难以收敛", 1),
    ("调试链路：错误反馈是编译报错而非结构化 delta", 1),
    "结论：AI 时代必须有一个 LLM 友好的「入口语言」",
]

LANG_STRATEGY = [
    "语义 / 策略双段式：kernel { semantics } / strategy { @rationale decisions }",
    "算子级抽象：不做 loop-nest 语言，循环 / 线程 / 显存是编译器的事",
    "Universal Op + 外部 Op Registry：语言不枚举算子，支持任意扩展",
    "符号 shape 一等公民：where B: dynamic(max=64), S: dynamic(max=8192), D: static",
    "@rationale 作为语法一等公民，每个决策强制携带自然语言理由",
    "Token 最小化：新特性（symbolic shape / type inference）必须进一步降 token",
]

LANG_DISCUSSION = [
    "DSL vs 嵌入式 Python DSL vs 注解式：给 LLM 用的语言，是否必须正则 + 无歧义文法？",
    "算子级抽象天花板：什么场景必须暴露 loop / schedule？是否需要 escape hatch？",
    "条件策略（when / otherwise）：应放在 Language 表达还是 Strategy 表达？",
    "多入口（.ak / 自然语言 / Python API）的 LLM 友好度与 token 成本如何量化对比？",
    "Language 演进（v2 where / tuple / type inference）如何保证不破坏已训练 LLM 的先验？",
]

IR_NECESSITY = [
    "MLIR / LLVM IR 为编译器工程师而设计：C++ / TableGen / dialect 爆炸 / SSA 心智负担",
    "单层 IR 混合「算什么 / 怎么调度 / 指令怎么排」，LLM 决策粒度被混淆",
    "动态 shape 外挂（Python 侧 shape 信息 / 运行时 guard）与编译期分析割裂",
    "现状：算子知识散落在 6 个文件、~3000 行冗余；加 1 op 要改 6 处；shape 推断 401 行 if/elif",
    "结论：需要为 LLM 设计一套多层、算法无关、可序列化的 IR —— LLM-Native IR",
]

IR_STRATEGY = [
    "四层 IR · 按 LLM 参与度分层",
    ("L4 SemanticIR · LLM 主作者 · 不可变 · SSA by construction", 1),
    ("L3 StrategyIR · LLM 主决策者 · 可搜索 / 可回滚 · L1/L2/L3 决策级别", 1),
    ("L2 ScheduleIR · LLM 仅 review · thread / block / warp / vector 映射", 1),
    ("L1 InstructionIR · LLM 不介入 · 近 LLVM IR", 1),
    "JSON 是序列化格式，不是 IR —— IR 以类型化对象存在",
    "算法无关 · 节点 / 边 / 决策原语通用；具体算子由 Op Registry 注入",
    "MLIR / LLVM 作为后端而非依赖：可经 linalg / transform / scf / gpu 下降，也可直接 emit",
]

IR_DISCUSSION = [
    "为什么是 4 层？「LLM 参与度」能否作为分层的一等依据？",
    "StrategyIR L1/L2/L3 决策在 Triton 天花板下能否全部表达？何时必须切 MLIR？",
    "符号维度（M / seq_len）如何在四层之间一致传播？约束求解放在哪一层？",
    "ConditionalNode vs 任意 CFG：AI 生成的 kernel 是否永远不需要 goto 式控制流？",
    "多后端共存时，IR 的「可 lower 性」如何不被单一后端绑架？",
]

COMP_NECESSITY = [
    "反面教材：加一个算子改 6 处 · shape 推断 401 行 if/elif · 45 个手写 NumPy 参考",
    "编译器如果又优化又下降又 codegen，会侵蚀 LLM 的决策空间并稀释正确性保证",
    "Pass 逻辑串行耦合 · 无 pre/post condition · 无 analysis cache",
    "后端切换（Triton → MLIR → LLVM）没有统一协议，target routing 靠 if/elif",
    "结论：工具链必须从「硬编码」走向「数据驱动 · 可插拔 · 可度量」",
]

COMP_STRATEGY = [
    "OpRegistry as Single Source of Truth · shape rule / template / reference / input gen / strategy 派生",
    "Pass Infrastructure · Pass 协议 + Pipeline + PassContext · Analysis / Transform / Verification 三族",
    "SemanticInterpreter · 用 PyTorch eager 执行 IR graph，替代 45 个手写 NumPy 参考",
    "ShapeInferenceEngine · 声明式规则替代 401 行 if/elif 链",
    "Backend Abstraction + BackendRegistry · Triton / MLIR / LLVM 同一协议 · target routing 注册中心化",
    "Compiler-as-Verifier · V0 <1ms 静态 / V1 数值对参考 / V2 HW 真实 profile",
    "Checkpoint / Rollback · 门禁顺序 Function ▸ Accuracy ▸ Performance",
]

COMP_DISCUSSION = [
    "「编译器不优化只验证」这条线能否贯穿到 LLVM IR 层？哪里必须编译器自己做（如 regalloc）？",
    "OpRegistry 边界：硬件特化 variant 写进 Registry 还是下沉到 Backend？",
    "V0 <1ms 如何工程化保证？是否需要再加 L0 类型 / 语法层的即时反馈？",
    "Pass 依赖与可重入性：Analysis 结果缓存失效？pre/post-condition 是否应编译期检查？",
    "Triton 抽象天花板具体卡在哪些决策上？何时切 MLIR？何时直出 LLVM？",
]

AGENT_NECESSITY = [
    "LLM 三病 · 幻觉代码 / 上下文爆炸 / 决策不可追溯",
    "自由文本生成 token 经济学不可持续：~3500 tokens/迭代 × 数十轮",
    "即便 LLM 能力提升，「无界生成 + 事后 diff」模型无法提供审计 / 回放 / 迁移能力",
    "外部 Agent（Cursor / Claude Code）已具 LLM 能力，缺的是结构化工具协议 + 身份注入",
    "结论：Agent 必须被工程化为「有界动作空间 + 结构化协议 + 预算治理」的系统",
]

AGENT_STRATEGY = [
    "Bounded Action Space · LLM 只从编译器枚举的合法动作中选，不生成自由代码",
    "Tool-Use 协议 · analyze_compute / list_legal_actions / apply_decision / verify / profile / checkpoint / rollback",
    "双模集成 · Mode A 内置 (CLI / Python API) / Mode B 外部 (Cursor · Claude Code · OpenClaw)",
    "Budget 管理 · decisions / compiles 分开计数 · 实时回注到 tool_result",
    "Context 工程 · 分段 Prompt Cache (Role / HW / Kernel / Strategy) · 50 轮节省 ~176K tokens",
    "AsyncGenerator 优化循环 · typed OptimizationEvent · 流式 + 可取消",
    "@rationale 沉淀为 trajectory JSONL · 跨硬件迁移与回放的燃料",
]

AGENT_DISCUSSION = [
    "Bounded Action Space 的宽窄：太窄（LLM 被困）vs 太宽（幻觉回潮）—— 能否量化调优？",
    "预算设计：按 decision / compile / token 哪个是第一性的？三者如何组合？",
    "长会话 Compact 如何不破坏优化连续性？短期记忆 vs 长期知识库如何分层？",
    "@rationale 能否跨硬件真正迁移？Phase 2 目标 ≥10% cross-arch lift —— 如何设计证伪实验？",
    "Agent 自主化（kernel-only → 自动 strategy + 自动 3 轮 compile→profile→adjust）可靠性边界？",
    "当 LLM 能力年年提升，Bounded Action Space 是否会变成枷锁？什么情况下放宽？",
]


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    make_cover(prs)                                           # 1
    make_agenda(prs)                                          # 2
    make_bg_problem(prs)                                      # 3
    make_bg_fourset(prs)                                      # 4

    make_piece_intro(prs, 1, "Language", "AI-Native DSL", ACCENT,
                     "给 LLM 一个「算子级抽象 × 语义策略分离 × token 最小化」的编程表面 —— 让 LLM 从「写代码」走向「做决策」",
                     page=5)                                  # 5
    make_two_col_content(prs, "02 · Language", "必要性  ·  构建策略",
                         "必要性  ·  为什么必须存在", LANG_NECESSITY,
                         "构建策略  ·  如何做对", LANG_STRATEGY,
                         page=6, accent=ACCENT)               # 6
    make_discussion(prs, "02 · Language", "技术讨论点",
                    LANG_DISCUSSION, page=7, accent=ACCENT)   # 7

    make_piece_intro(prs, 2, "IR", "Multi-Layer LLM-Native IR", ACCENT_ALT,
                     "把「算什么 / 怎么优化 / 怎么调度 / 怎么成指令」按 LLM 参与度分为 4 层 —— 每一层只做一件事，每一层都有清晰的验证责任",
                     page=8)                                  # 8
    make_two_col_content(prs, "03 · IR", "必要性  ·  构建策略",
                         "必要性  ·  为什么必须存在", IR_NECESSITY,
                         "构建策略  ·  如何做对", IR_STRATEGY,
                         page=9, accent=ACCENT_ALT)           # 9
    make_discussion(prs, "03 · IR", "技术讨论点",
                    IR_DISCUSSION, page=10, accent=ACCENT_ALT)  # 10

    make_piece_intro(prs, 3, "Compiler Toolchain",
                     "OpRegistry × Pass × Verifier", ACCENT_BLUE,
                     "把「算子 / Pass / 后端」从硬编码变成数据驱动 —— 让「加一个算子 / 换一个后端」不再是重写，而是注册",
                     page=11)                                 # 11
    make_two_col_content(prs, "04 · Compiler Toolchain", "必要性  ·  构建策略",
                         "必要性  ·  为什么必须存在", COMP_NECESSITY,
                         "构建策略  ·  如何做对", COMP_STRATEGY,
                         page=12, accent=ACCENT_BLUE)         # 12
    make_discussion(prs, "04 · Compiler Toolchain", "技术讨论点",
                    COMP_DISCUSSION, page=13, accent=ACCENT_BLUE)  # 13

    make_piece_intro(prs, 4, "Agent Engineering",
                     "Bounded Action × Tool Protocol", ACCENT_RED,
                     "把 LLM 从「自由文本生成器」约束成「可审计 / 可回放 / 可学习」的优化决策者 —— 有界动作空间 + 结构化协议 + 预算治理",
                     page=14)                                 # 14
    make_two_col_content(prs, "05 · Agent", "必要性  ·  构建策略",
                         "必要性  ·  为什么必须存在", AGENT_NECESSITY,
                         "构建策略  ·  如何做对", AGENT_STRATEGY,
                         page=15, accent=ACCENT_RED)          # 15
    make_discussion(prs, "05 · Agent", "技术讨论点",
                    AGENT_DISCUSSION, page=16, accent=ACCENT_RED)  # 16

    # 业界洞察 / 应有特征 / 方案对比
    make_industry_insight_1(prs, page=17)                     # 17
    make_industry_insight_2(prs, page=18)                     # 18
    make_required_traits(prs, page=19)                        # 19
    make_comparison(prs, page=20)                             # 20
    make_advantage_detail(prs, page=21)                       # 21

    # 架构图
    make_overall_arch(prs, page=22)                           # 22
    make_lang_arch(prs, page=23)                              # 23
    make_ir_arch(prs, page=24)                                # 24
    make_compiler_arch(prs, page=25)                          # 25
    make_agent_arch(prs, page=26)                             # 26

    # 贯通 + 收尾
    make_integration(prs, page=27)                            # 27
    make_verification(prs, page=28)                           # 28
    make_benchmark(prs, page=29)                              # 29
    make_token_economics(prs, page=30)                        # 30
    make_roadmap(prs, page=31)                                # 31
    make_closing(prs, page=32)                                # 32

    prs.save(out_path)
    print(f"wrote: {out_path}")


if __name__ == "__main__":
    out_path = (Path(__file__).resolve().parent.parent
                / "docs" / "sharing" / "ai-native-compile-stack.pptx")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    build(out_path)
